import os
from anthropic import Anthropic
import streamlit as st
import io
import google.generativeai as genai
from PIL import Image
import datetime
import os
from pymongo import MongoClient
from bson import ObjectId
import json
import hashlib
from google.genai import types
import PyPDF2
from pptx import Presentation
import docx
import openai
from typing import List, Dict, Tuple
import hashlib
import pandas as pd
import re
from pypdf import PdfReader, PdfWriter
from pypdf.annotations import Text
import requests

# Configuração inicial
st.set_page_config(
    layout="wide",
    page_title="Agente Social",
    page_icon="🤖"
)

import os
import PyPDF2
import pdfplumber
from pathlib import Path

# --- CONFIGURAÇÃO DOS MODELOS ---
# Configuração da API do Anthropic (Claude)
anthropic_api_key = os.getenv("ANTHROPIC_API_KEY")
if anthropic_api_key:
    anthropic_client = Anthropic(api_key=anthropic_api_key)
else:
    st.error("ANTHROPIC_API_KEY não encontrada nas variáveis de ambiente")
    anthropic_client = None

# Configuração da API do Gemini
gemini_api_key = os.getenv("GEM_API_KEY")
if gemini_api_key:
    genai.configure(api_key=gemini_api_key)
    modelo_vision = genai.GenerativeModel("gemini-2.0-flash", generation_config={"temperature": 0.0})
    modelo_texto = genai.GenerativeModel("gemini-2.0-flash")
else:
    st.error("GEM_API_KEY não encontrada nas variáveis de ambiente")
    modelo_vision = None
    modelo_texto = None

import os
import PyPDF2
import pdfplumber
from pathlib import Path

# --- FUNÇÕES AUXILIARES MELHORADAS ---

def criar_prompt_validacao_preciso(texto, nome_arquivo, contexto_agente):
    """Cria um prompt de validação muito mais preciso para evitar falsos positivos"""
    
    prompt = f"""
{contexto_agente}


###BEGIN TEXTO PARA VALIDAÇÃO###
**Arquivo:** {nome_arquivo}
**Conteúdo:**
{texto[:12000]}
###END TEXTO PARA VALIDAÇÃO###

## FORMATO DE RESPOSTA OBRIGATÓRIO:



### ✅ CONFORMIDADE COM DIRETRIZES
- [Itens que estão alinhados com as diretrizes de branding]



**INCONSISTÊNCIAS COM BRANDING:**
- [Só liste desvios REAIS das diretrizes de branding]

### 💡 TEXTO REVISADO
- [Sugestões para aprimorar]

### 📊 STATUS FINAL
**Documento:** [Aprovado/Necessita ajustes/Reprovado]
**Principais ações necessárias:** [Lista resumida]

"""
    return prompt


# --- FUNÇÃO PARA ESCOLHER ENTRE GEMINI E CLAUDE ---
def gerar_resposta_modelo(prompt: str, modelo_escolhido: str = "Gemini", contexto_agente: str = None) -> str:
    """
    Gera resposta usando Gemini ou Claude baseado na escolha do usuário
    """
    try:
        if modelo_escolhido == "Gemini" and modelo_texto:
            if contexto_agente:
                prompt_completo = f"{contexto_agente}\n\n{prompt}"
            else:
                prompt_completo = prompt
            
            resposta = modelo_texto.generate_content(prompt_completo)
            return resposta.text
            
        elif modelo_escolhido == "Claude" and anthropic_client:
            if contexto_agente:
                system_prompt = contexto_agente
            else:
                system_prompt = "Você é um assistente útil."
            
            message = anthropic_client.messages.create(
                max_tokens=4000,
                messages=[{"role": "user", "content": prompt}],
                model="claude-haiku-4-5-20251001",
                system=system_prompt
            )
            return message.content[0].text
            
        else:
            return f"❌ Modelo {modelo_escolhido} não disponível. Verifique as configurações da API."
            
    except Exception as e:
        return f"❌ Erro ao gerar resposta com {modelo_escolhido}: {str(e)}"

def analisar_documento_por_slides(doc, contexto_agente):
    """Analisa documento slide por slide com alta precisão"""
    
    resultados = []
    
    for i, slide in enumerate(doc['slides']):
        with st.spinner(f"Analisando slide {i+1}..."):
            try:
                prompt_slide = f"""
{contexto_agente}

## ANÁLISE POR SLIDE - PRECISÃO ABSOLUTA

###BEGIN TEXTO PARA VALIDAÇÃO###
**SLIDE {i+1}:**
{slide['conteudo'][:2000]}
###END TEXTO PARA VALIDAÇÃO###


**ANÁLISE DO SLIDE {i+1}:**

### ✅ Pontos Fortes:
[O que está bom neste slide]

### ⚠️ Problemas REAIS:
- [Lista CURTA de problemas]

### 💡 Sugestões Específicas:
[Melhorias para ESTE slide específico]

Considere que slides que são introdutórios ou apenas de títulos não precisam de tanto rigor de branding

**STATUS:** [✔️ Aprovado / ⚠️ Ajustes Menores / ❌ Problemas Sérios]
"""
                
                resposta = modelo_texto.generate_content(prompt_slide)
                resultados.append({
                    'slide_num': i+1,
                    'analise': resposta.text,
                    'tem_alteracoes': '❌' in resposta.text or '⚠️' in resposta.text
                })
                
            except Exception as e:
                resultados.append({
                    'slide_num': i+1,
                    'analise': f"❌ Erro na análise do slide: {str(e)}",
                    'tem_alteracoes': False
                })
    
    # Construir relatório consolidado
    relatorio = f"# 📊 RELATÓRIO DE VALIDAÇÃO - {doc['nome']}\n\n"
    relatorio += f"**Total de Slides:** {len(doc['slides'])}\n"
    relatorio += f"**Slides com Alterações:** {sum(1 for r in resultados if r['tem_alteracoes'])}\n\n"
    
    # Slides que precisam de atenção
    slides_com_problemas = [r for r in resultados if r['tem_alteracoes']]
    if slides_com_problemas:
        relatorio += "## 🚨 SLIDES QUE PRECISAM DE ATENÇÃO:\n\n"
        for resultado in slides_com_problemas:
            relatorio += f"### 📋 Slide {resultado['slide_num']}\n"
            relatorio += f"{resultado['analise']}\n\n"
    
    # Resumo executivo
    relatorio += "## 📈 RESUMO EXECUTIVO\n\n"
    if slides_com_problemas:
        relatorio += f"**⚠️ {len(slides_com_problemas)} slide(s) necessitam de ajustes**\n"
        relatorio += f"**✅ {len(doc['slides']) - len(slides_com_problemas)} slide(s) estão adequados**\n"
    else:
        relatorio += "**🎉 Todos os slides estão em conformidade com as diretrizes!**\n"
    
    return relatorio

def extract_text_from_pdf_com_slides(arquivo_pdf):
    """Extrai texto de PDF com informação de páginas"""
    try:
        import PyPDF2
        pdf_reader = PyPDF2.PdfReader(arquivo_pdf)
        slides_info = []
        
        for pagina_num, pagina in enumerate(pdf_reader.pages):
            texto = pagina.extract_text()
            slides_info.append({
                'numero': pagina_num + 1,
                'conteudo': texto,
                'tipo': 'página'
            })
        
        texto_completo = "\n\n".join([f"--- PÁGINA {s['numero']} ---\n{s['conteudo']}" for s in slides_info])
        return texto_completo, slides_info
        
    except Exception as e:
        return f"Erro na extração PDF: {str(e)}", []

def extract_text_from_pptx_com_slides(arquivo_pptx):
    """Extrai texto de PPTX com informação de slides"""
    try:
        from pptx import Presentation
        import io
        
        prs = Presentation(io.BytesIO(arquivo_pptx.read()))
        slides_info = []
        
        for slide_num, slide in enumerate(prs.slides):
            texto_slide = f"--- SLIDE {slide_num + 1} ---\n"
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texto_slide += shape.text + "\n"
            
            slides_info.append({
                'numero': slide_num + 1,
                'conteudo': texto_slide,
                'tipo': 'slide'
            })
        
        texto_completo = "\n\n".join([s['conteudo'] for s in slides_info])
        return texto_completo, slides_info
        
    except Exception as e:
        return f"Erro na extração PPTX: {str(e)}", []

def extrair_texto_arquivo(arquivo):
    """Extrai texto de arquivos TXT e DOCX"""
    try:
        if arquivo.type == "text/plain":
            return str(arquivo.read(), "utf-8")
        elif arquivo.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            import docx
            import io
            doc = docx.Document(io.BytesIO(arquivo.read()))
            texto = ""
            for para in doc.paragraphs:
                texto += para.text + "\n"
            return texto
        else:
            return f"Tipo não suportado: {arquivo.type}"
    except Exception as e:
        return f"Erro na extração: {str(e)}"

def extract_text_from_pdf(pdf_path):
    """
    Extract text from a PDF file using multiple methods for better coverage
    """
    text = ""

    # Method 1: Try with pdfplumber (better for some PDFs)
    try:
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text
    except Exception as e:
        print(f"pdfplumber failed for {pdf_path}: {e}")

    # Method 2: Fallback to PyPDF2 if pdfplumber didn't extract much text
    if len(text.strip()) < 100:  # If very little text was extracted
        try:
            with open(pdf_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text 
        except Exception as e:
            print(f"PyPDF2 also failed for {pdf_path}: {e}")

    return text
    

# --- Sistema de Autenticação MELHORADO ---
def make_hashes(password):
    return hashlib.sha256(str.encode(password)).hexdigest()

def check_hashes(password, hashed_text):
    return make_hashes(password) == hashed_text

# Dados de usuário (em produção, isso deve vir de um banco de dados seguro)
users_db = {
    "admin": {
        "password": make_hashes("senha1234"),
        "squad": "admin",
        "nome": "Administrador"
    }
}

# Conexão MongoDB
client = MongoClient("mongodb+srv://gustavoromao3345:RqWFPNOJQfInAW1N@cluster0.5iilj.mongodb.net/auto_doc?retryWrites=true&w=majority&ssl=true&ssl_cert_reqs=CERT_NONE&tlsAllowInvalidCertificates=true")
db = client['agentes_personalizados']
collection_agentes = db['agentes']
collection_conversas = db['conversas']
collection_usuarios = db['usuarios']  # Nova coleção para usuários

# --- FUNÇÕES DE CADASTRO E LOGIN ---
def criar_usuario(email, senha, nome, squad):
    """Cria um novo usuário no banco de dados"""
    try:
        # Verificar se usuário já existe
        if collection_usuarios.find_one({"email": email}):
            return False, "Usuário já existe"
        
        # Criar hash da senha
        senha_hash = make_hashes(senha)
        
        novo_usuario = {
            "email": email,
            "senha": senha_hash,
            "nome": nome,
            "squad": squad,
            "data_criacao": datetime.datetime.now(),
            "ultimo_login": None,
            "ativo": True
        }
        
        result = collection_usuarios.insert_one(novo_usuario)
        return True, "Usuário criado com sucesso"
        
    except Exception as e:
        return False, f"Erro ao criar usuário: {str(e)}"

def verificar_login(email, senha):
    """Verifica as credenciais do usuário"""
    try:
        # Primeiro verificar no banco de dados
        usuario = collection_usuarios.find_one({"email": email, "ativo": True})
        
        if usuario:
            if check_hashes(senha, usuario["senha"]):
                # Atualizar último login
                collection_usuarios.update_one(
                    {"_id": usuario["_id"]},
                    {"$set": {"ultimo_login": datetime.datetime.now()}}
                )
                return True, usuario, "Login bem-sucedido"
            else:
                return False, None, "Senha incorreta"
        
        # Fallback para usuários hardcoded (apenas para admin)
        if email in users_db:
            user_data = users_db[email]
            if check_hashes(senha, user_data["password"]):
                usuario_fallback = {
                    "email": email,
                    "nome": user_data["nome"],
                    "squad": user_data["squad"],
                    "_id": "admin"
                }
                return True, usuario_fallback, "Login bem-sucedido"
            else:
                return False, None, "Senha incorreta"
        
        return False, None, "Usuário não encontrado"
        
    except Exception as e:
        return False, None, f"Erro no login: {str(e)}"

def get_current_user():
    """Retorna o usuário atual da sessão"""
    return st.session_state.get('user', {})

def get_current_squad():
    """Retorna o squad do usuário atual"""
    user = get_current_user()
    return user.get('squad', 'unknown')

def login():
    """Formulário de login e cadastro"""
    st.title("🔒 Agente Social - Login")
    
    tab_login, tab_cadastro = st.tabs(["Login", "Cadastro"])
    
    with tab_login:
        with st.form("login_form"):
            email = st.text_input("Email")
            password = st.text_input("Senha", type="password")
            submit_button = st.form_submit_button("Login")
            
            if submit_button:
                if email and password:
                    sucesso, usuario, mensagem = verificar_login(email, password)
                    if sucesso:
                        st.session_state.logged_in = True
                        st.session_state.user = usuario
                        st.success("Login realizado com sucesso!")
                        st.rerun()
                    else:
                        st.error(mensagem)
                else:
                    st.error("Por favor, preencha todos os campos")
    
    with tab_cadastro:
        with st.form("cadastro_form"):
            st.subheader("Criar Nova Conta")
            
            nome = st.text_input("Nome Completo")
            email = st.text_input("Email")
            squad = st.selectbox(
                "Selecione seu Squad:",
                ["Syngenta", "SME", "Enterprise"],
                help="Escolha o squad ao qual você pertence"
            )
            senha = st.text_input("Senha", type="password")
            confirmar_senha = st.text_input("Confirmar Senha", type="password")
            
            submit_cadastro = st.form_submit_button("Criar Conta")
            
            if submit_cadastro:
                if not all([nome, email, squad, senha, confirmar_senha]):
                    st.error("Por favor, preencha todos os campos")
                elif senha != confirmar_senha:
                    st.error("As senhas não coincidem")
                elif len(senha) < 6:
                    st.error("A senha deve ter pelo menos 6 caracteres")
                else:
                    sucesso, mensagem = criar_usuario(email, senha, nome, squad)
                    if sucesso:
                        st.success("Conta criada com sucesso! Faça login para continuar.")
                    else:
                        st.error(mensagem)

# Verificar se o usuário está logado
if "logged_in" not in st.session_state:
    st.session_state.logged_in = False

if not st.session_state.logged_in:
    login()
    st.stop()

# --- CONFIGURAÇÕES APÓS LOGIN ---
gemini_api_key = os.getenv("GEM_API_KEY")
if not gemini_api_key:
    st.error("GEMINI_API_KEY não encontrada nas variáveis de ambiente")
    st.stop()

genai.configure(api_key=gemini_api_key)
modelo_vision = genai.GenerativeModel("gemini-2.5-flash", generation_config={"temperature": 0.1})
modelo_texto = genai.GenerativeModel("gemini-2.5-flash")

# Configuração da API do Perplexity
perp_api_key = os.getenv("PERP_API_KEY")
if not perp_api_key:
    st.error("PERP_API_KEY não encontrada nas variáveis de ambiente")

# --- Configuração de Autenticação de Administrador ---
def check_admin_password():
    """Retorna True para usuários admin sem verificação de senha."""
    return st.session_state.user.get('squad') == "admin"

# --- FUNÇÕES CRUD PARA AGENTES (MODIFICADAS PARA SQUADS) ---
def criar_agente(nome, system_prompt, base_conhecimento, comments, planejamento, categoria, squad_permitido, agente_mae_id=None, herdar_elementos=None):
    """Cria um novo agente no MongoDB com squad permitido"""
    agente = {
        "nome": nome,
        "system_prompt": system_prompt,
        "base_conhecimento": base_conhecimento,
        "comments": comments,
        "planejamento": planejamento,
        "categoria": categoria,
        "squad_permitido": squad_permitido,  # Novo campo
        "agente_mae_id": agente_mae_id,
        "herdar_elementos": herdar_elementos or [],
        "data_criacao": datetime.datetime.now(),
        "ativo": True,
        "criado_por": get_current_user().get('email', 'unknown'),
        "criado_por_squad": get_current_squad()  # Novo campo
    }
    result = collection_agentes.insert_one(agente)
    return result.inserted_id

def listar_agentes():
    """Retorna todos os agentes ativos que o usuário atual pode ver"""
    current_squad = get_current_squad()
    
    # Admin vê todos os agentes
    if current_squad == "admin":
        return list(collection_agentes.find({"ativo": True}).sort("data_criacao", -1))
    
    # Usuários normais veem apenas agentes do seu squad ou squad "Todos"
    return list(collection_agentes.find({
        "ativo": True,
        "$or": [
            {"squad_permitido": current_squad},
            {"squad_permitido": "Todos"},
            {"criado_por_squad": current_squad}  # Usuário pode ver seus próprios agentes
        ]
    }).sort("data_criacao", -1))

def listar_agentes_para_heranca(agente_atual_id=None):
    """Retorna todos os agentes ativos que podem ser usados como mãe (com filtro de squad)"""
    current_squad = get_current_squad()
    
    query = {"ativo": True}
    
    # Filtro por squad
    if current_squad != "admin":
        query["$or"] = [
            {"squad_permitido": current_squad},
            {"squad_permitido": "Todos"},
            {"criado_por_squad": current_squad}
        ]
    
    if agente_atual_id:
        # Excluir o próprio agente da lista de opções para evitar auto-herança
        if isinstance(agente_atual_id, str):
            agente_atual_id = ObjectId(agente_atual_id)
        query["_id"] = {"$ne": agente_atual_id}
    
    return list(collection_agentes.find(query).sort("data_criacao", -1))

def obter_agente(agente_id):
    """Obtém um agente específico pelo ID com verificação de permissão por squad"""
    if isinstance(agente_id, str):
        agente_id = ObjectId(agente_id)
    
    agente = collection_agentes.find_one({"_id": agente_id})
    
    # Verificar permissão baseada no squad
    if agente and agente.get('ativo', True):
        current_squad = get_current_squad()
        
        # Admin pode ver tudo
        if current_squad == "admin":
            return agente
        
        # Usuários normais só podem ver agentes do seu squad ou "Todos"
        squad_permitido = agente.get('squad_permitido')
        criado_por_squad = agente.get('criado_por_squad')
        
        if squad_permitido == current_squad or squad_permitido == "Todos" or criado_por_squad == current_squad:
            return agente
    
    return None

def atualizar_agente(agente_id, nome, system_prompt, base_conhecimento, comments, planejamento, categoria, squad_permitido, agente_mae_id=None, herdar_elementos=None):
    """Atualiza um agente existente com verificação de permissão"""
    if isinstance(agente_id, str):
        agente_id = ObjectId(agente_id)
    
    # Verificar se o usuário tem permissão para editar este agente
    agente_existente = obter_agente(agente_id)
    if not agente_existente:
        raise PermissionError("Agente não encontrado ou sem permissão de edição")
    
    return collection_agentes.update_one(
        {"_id": agente_id},
        {
            "$set": {
                "nome": nome,
                "system_prompt": system_prompt,
                "base_conhecimento": base_conhecimento,
                "comments": comments,
                "planejamento": planejamento,
                "categoria": categoria,
                "squad_permitido": squad_permitido,  # Novo campo
                "agente_mae_id": agente_mae_id,
                "herdar_elementos": herdar_elementos or [],
                "data_atualizacao": datetime.datetime.now()
            }
        }
    )

def desativar_agente(agente_id):
    """Desativa um agente (soft delete) com verificação de permissão"""
    if isinstance(agente_id, str):
        agente_id = ObjectId(agente_id)
    
    # Verificar se o usuário tem permissão para desativar este agente
    agente_existente = obter_agente(agente_id)
    if not agente_existente:
        raise PermissionError("Agente não encontrado ou sem permissão para desativar")
    
    return collection_agentes.update_one(
        {"_id": agente_id},
        {"$set": {"ativo": False, "data_desativacao": datetime.datetime.now()}}
    )

def obter_agente_com_heranca(agente_id):
    """Obtém um agente com os elementos herdados aplicados"""
    agente = obter_agente(agente_id)
    if not agente or not agente.get('agente_mae_id'):
        return agente
    
    agente_mae = obter_agente(agente['agente_mae_id'])
    if not agente_mae:
        return agente
    
    elementos_herdar = agente.get('herdar_elementos', [])
    agente_completo = agente.copy()
    
    for elemento in elementos_herdar:
        if elemento == 'system_prompt' and not agente_completo.get('system_prompt'):
            agente_completo['system_prompt'] = agente_mae.get('system_prompt', '')
        elif elemento == 'base_conhecimento' and not agente_completo.get('base_conhecimento'):
            agente_completo['base_conhecimento'] = agente_mae.get('base_conhecimento', '')
        elif elemento == 'comments' and not agente_completo.get('comments'):
            agente_completo['comments'] = agente_mae.get('comments', '')
        elif elemento == 'planejamento' and not agente_completo.get('planejamento'):
            agente_completo['planejamento'] = agente_mae.get('planejamento', '')
    
    return agente_completo

def salvar_conversa(agente_id, mensagens, segmentos_utilizados=None):
    """Salva uma conversa no histórico"""
    if isinstance(agente_id, str):
        agente_id = ObjectId(agente_id)
    conversa = {
        "agente_id": agente_id,
        "mensagens": mensagens,
        "segmentos_utilizados": segmentos_utilizados,
        "data_criacao": datetime.datetime.now()
    }
    return collection_conversas.insert_one(conversa)

def obter_conversas(agente_id, limite=10):
    """Obtém o histórico de conversas de um agente"""
    if isinstance(agente_id, str):
        agente_id = ObjectId(agente_id)
    return list(collection_conversas.find(
        {"agente_id": agente_id}
    ).sort("data_criacao", -1).limit(limite))

# --- Função para construir contexto com segmentos selecionados ---
def construir_contexto(agente, segmentos_selecionados, historico_mensagens=None):
    """Constrói o contexto com base nos segmentos selecionados"""
    contexto = ""
    
    if "system_prompt" in segmentos_selecionados and agente.get('system_prompt'):
        contexto += f"### INSTRUÇÕES DO SISTEMA ###\n{agente['system_prompt']}\n\n"
    
    if "base_conhecimento" in segmentos_selecionados and agente.get('base_conhecimento'):
        contexto += f"### BASE DE CONHECIMENTO ###\n{agente['base_conhecimento']}\n\n"
    
    if "comments" in segmentos_selecionados and agente.get('comments'):
        contexto += f"### COMENTÁRIOS DO CLIENTE ###\n{agente['comments']}\n\n"
    
    if "planejamento" in segmentos_selecionados and agente.get('planejamento'):
        contexto += f"### PLANEJAMENTO ###\n{agente['planejamento']}\n\n"
    
    # Adicionar histórico se fornecido
    if historico_mensagens:
        contexto += "### HISTÓRICO DA CONVERSA ###\n"
        for msg in historico_mensagens:
            contexto += f"{msg['role']}: {msg['content']}\n"
        contexto += "\n"
    
    contexto += "### RESPOSTA ATUAL ###\nassistant:"
    
    return contexto

# --- MODIFICAÇÃO: SELECTBOX PARA SELEÇÃO DE AGENTE ---
def selecionar_agente_interface():
    """Interface para seleção de agente usando selectbox"""
    st.title("🤖 Agente Social")
    
    # Carregar agentes disponíveis
    agentes = listar_agentes()
    
    if not agentes:
        st.error("❌ Nenhum agente disponível. Crie um agente primeiro na aba de Gerenciamento.")
        return None
    
    # Preparar opções para o selectbox
    opcoes_agentes = []
    for agente in agentes:
        agente_completo = obter_agente_com_heranca(agente['_id'])
        if agente_completo:  # Só adiciona se tiver permissão
            descricao = f"{agente['nome']} - {agente.get('categoria', 'Social')}"
            if agente.get('agente_mae_id'):
                descricao += " 🔗"
            # Adicionar indicador de squad
            squad_permitido = agente.get('squad_permitido', 'Todos')
            descricao += f" 👥{squad_permitido}"
            opcoes_agentes.append((descricao, agente_completo))
    
    if opcoes_agentes:
        # Selectbox para seleção de agente
        agente_selecionado_desc = st.selectbox(
            "Selecione uma base de conhecimento para usar o sistema:",
            options=[op[0] for op in opcoes_agentes],
            index=0,
            key="selectbox_agente_principal"
        )
        
        # Encontrar o agente completo correspondente
        agente_completo = None
        for desc, agente in opcoes_agentes:
            if desc == agente_selecionado_desc:
                agente_completo = agente
                break
        
        if agente_completo and st.button("✅ Confirmar Seleção", key="confirmar_agente"):
            st.session_state.agente_selecionado = agente_completo
            st.session_state.messages = []
            st.session_state.segmentos_selecionados = ["system_prompt", "base_conhecimento", "comments", "planejamento"]
            st.success(f"✅ Agente '{agente_completo['nome']}' selecionado!")
            st.rerun()
        
        return agente_completo
    else:
        st.info("Nenhum agente disponível com as permissões atuais.")
        return None

# --- Verificar se o agente já foi selecionado ---
if "agente_selecionado" not in st.session_state:
    st.session_state.agente_selecionado = None

# Se não há agente selecionado, mostrar interface de seleção
if not st.session_state.agente_selecionado:
    selecionar_agente_interface()
    st.stop()

# --- INTERFACE PRINCIPAL (apenas se agente estiver selecionado) ---
agente_selecionado = st.session_state.agente_selecionado

def is_syn_agent(agent_name):
    """Verifica se o agente é da baseado no nome"""
    return agent_name and any(keyword in agent_name.upper() for keyword in ['SYN'])

PRODUCT_DESCRIPTIONS = {
    "FORTENZA": "Tratamento de sementes inseticida, focado no Cerrado e posicionado para controle do complexo de lagartas e outras pragas iniciais. Comunicação focada no mercado 'on farm' (tratamento feito na fazenda).",
    "ALADE": "Fungicida para controle de doenças em soja, frequentemente posicionado em programa com Mitrion para controle de podridões de vagens e grãos.",
    "VERDAVIS": "Inseticida e acaricida composto por PLINAZOLIN® technology (nova molécula, novo grupo químico, modo de ação inédito) + lambda-cialotrina. KBFs: + mais choque, + mais espectro e + mais dias de controle.",
    "ENGEO PLENO S": "Inseticida de tradição, referência no controle de percevejos. Mote: 'Nunca foi sorte. Sempre foi Engeo Pleno S'.",
    "MEGAFOL": "Bioativador da Syn Biologicals. Origem 100% natural (extratos vegetais e de algas Ascophyllum nodosum). Desenvolvido para garantir que a planta alcance todo seu potencial produtivo.",
    "MIRAVIS DUO": "Fungicida da família Miravis. Traz ADEPIDYN technology (novo ingrediente ativo, novo grupo químico). Focado no controle de manchas foliares.",
    "AVICTA COMPLETO": "Oferta comercial de tratamento industrial de sementes (TSI). Composto por inseticida, fungicida e nematicida.",
    "MITRION": "Fungicida para controle de doenças em soja, frequentemente posicionado em programa com Alade.",
    "AXIAL": "Herbicida para trigo. Composto por um novo ingrediente ativo. Foco no controle do azevém.",
    "CERTANO": "Bionematicida e biofungicida. Composto pela bactéria Bacillus velezensis. Controla nematoides e fungos de solo.",
    "MANEJO LIMPO": "Programa da Syn para manejo integrado de plantas daninhas.",
    "ELESTAL NEO": "Fungicida para controle de doenças em soja e algodão.",
    "FRONDEO": "Inseticida para cana-de-açúcar com foco no controle da broca da cana.",
    "FORTENZA ELITE": "Oferta comercial de TSI. Solução robusta contre pragas, doenças e nematoides do Cerrado.",
    "REVERB": "Produto para manejo de doenças em soja e milho com ação prolongada ou de espectro amplo.",
    "YIELDON": "Produto focado em maximizar a produtividade das lavouras.",
    "ORONDIS FLEXI": "Fungicida com flexibilidade de uso para controle de requeima, míldios e manchas.",
    "RIZOLIQ LLI": "Inoculante ou produto para tratamento de sementes que atua na rizosfera.",
    "ARVATICO": "Fungicida ou inseticida com ação específica para controle de doenças foliares ou pragas.",
    "VERDADERO": "Produto relacionado à saúde do solo ou nutrição vegetal.",
    "MIRAVIS": "Fungicida da família Miravis para controle de doenças.",
    "MIRAVIS PRO": "Fungicida premium da família Miravis para controle avançado de doenças.",
    "INSTIVO": "Lagarticida posicionado como especialista no controle de lagartas do gênero Spodoptera.",
    "CYPRESS": "Fungicida posicionado para últimas aplicações na soja, consolidando o manejo de doenças.",
    "CALARIS": "Herbicida composto por atrazina + mesotriona para controle de plantas daninhas no milho.",
    "SPONTA": "Inseticida para algodão com PLINAZOLIN® technology para controle de bicudo e outras pragas.",
    "INFLUX": "Inseticida lagarticida premium para controle de todas as lagartas, especialmente helicoverpa.",
    "JOINER": "Inseticida acaricida com tecnologia PLINAZOLIN para culturas hortifrúti.",
    "DUAL GOLD": "Herbicida para manejo de plantas daninhas.",
}

def extract_product_info(text: str) -> Tuple[str, str, str]:
    """Extrai informações do produto do texto da célula"""
    if not text or not text.strip():
        return None, None, None
    
    text = str(text).strip()
    
    # Remover emojis e marcadores
    clean_text = re.sub(r'[🔵🟠🟢🔴🟣🔃📲]', '', text).strip()
    
    # Padrões para extração
    patterns = {
        'product': r'\b([A-Z][A-Za-z\s]+(?:PRO|S|NEO|LLI|ELITE|COMPLETO|DUO|FLEXI|PLENO|XTRA)?)\b',
        'culture': r'\b(soja|milho|algodão|cana|trigo|HF|café|citrus|batata|melão|uva|tomate|multi)\b',
        'action': r'\b(depoimento|resultados|série|reforço|controle|lançamento|importância|jornada|conceito|vídeo|ação|diferenciais|awareness|problemática|glossário|manejo|aplicação|posicionamento)\b'
    }
    
    product_match = re.search(patterns['product'], clean_text, re.IGNORECASE)
    culture_match = re.search(patterns['culture'], clean_text, re.IGNORECASE)
    action_match = re.search(patterns['action'], clean_text, re.IGNORECASE)
    
    product = product_match.group(1).strip().upper() if product_match else None
    culture = culture_match.group(0).lower() if culture_match else "multi"
    action = action_match.group(0).lower() if action_match else "conscientização"
    
    return product, culture, action

def generate_context(content, product_name, culture, action, data_input, formato_principal):
    """Gera o texto de contexto discursivo usando LLM"""
    if not gemini_api_key:
        return "API key do Gemini não configurada. Contexto não disponível."
    
    # Determinar mês em português
    meses = {
        1: "janeiro", 2: "fevereiro", 3: "março", 4: "abril",
        5: "maio", 6: "junho", 7: "julho", 8: "agosto",
        9: "setembro", 10: "outubro", 11: "novembro", 12: "dezembro"
    }
    mes = meses[data_input.month]
    
    prompt = f"""
    Como redator especializado em agronegócio da Syn, elabore um texto contextual discursivo de 3-4 parágrafos para uma pauta de conteúdo.

    Informações da pauta:
    - Produto: {product_name}
    - Cultura: {culture}
    - Ação/tema: {action}
    - Mês de publicação: {mes}
    - Formato principal: {formato_principal}
    - Conteúdo original: {content}

    Descrição do produto: {PRODUCT_DESCRIPTIONS.get(product_name, 'Produto agrícola')}

    Instruções:
    - Escreva em formato discursivo e fluido, com 3-4 parágrafos bem estruturados
    - Mantenha tom técnico mas acessível, adequado para produtores rurais
    - Contextualize a importância do tema para a cultura e época do ano
    - Explique por que este conteúdo é relevante neste momento
    - Inclua considerações sobre o público-alvo e objetivos da comunicação
    - Não repita literalmente a descrição do produto, mas a incorpore naturalmente no texto
    - Use linguagem persuasiva mas factual, baseada em dados técnicos

    Formato: Texto corrido em português brasileiro
    """
    
    try:
        response = modelo_texto.generate_content(prompt)
        return response.text
    except Exception as e:
        return f"Erro ao gerar contexto: {str(e)}"

def generate_platform_strategy(product_name, culture, action, content):
    """Gera estratégia por plataforma usando Gemini"""
    if not gemini_api_key:
        return "API key do Gemini não configurada. Estratégias por plataforma não disponíveis."
    
    prompt = f"""
    Como especialista em mídias sociais para o agronegócio, crie uma estratégia de conteúdo detalhada:

    PRODUTO: {product_name}
    CULTURA: {culture}
    AÇÃO: {action}
    CONTEÚDO ORIGINAL: {content}
    DESCRIÇÃO DO PRODUTO: {PRODUCT_DESCRIPTIONS.get(product_name, 'Produto agrícola')}

    FORNECER ESTRATÉGIA PARA:
    - Instagram (Feed, Reels, Stories)
    - Facebook 
    - LinkedIn
    - WhatsApp Business
    - YouTube
    - Portal Mais Agro (blog)

    INCLUIR PARA CADA PLATAFORMA:
    1. Tipo de conteúdo recomendado
    2. Formato ideal (vídeo, carrossel, estático, etc.)
    3. Tom de voz apropriado
    4. CTA específico
    5. Melhores práticas

    Formato: Texto claro com seções bem definidas
    """
    
    try:
        response = modelo_texto.generate_content(prompt)
        return response.text
    except Exception as e:
        return f"Erro ao gerar estratégia: {str(e)}"

def generate_briefing(content, product_name, culture, action, data_input, formato_principal):
    """Gera um briefing completo em formato de texto puro"""
    description = PRODUCT_DESCRIPTIONS.get(product_name, "Descrição do produto não disponível.")
    context = generate_context(content, product_name, culture, action, data_input, formato_principal)
    platform_strategy = generate_platform_strategy(product_name, culture, action, content)
    
    briefing = f"""
BRIEFING DE CONTEÚDO - {product_name} - {culture.upper()} - {action.upper()}

CONTEXTO E OBJETIVO
{context}

DESCRIÇÃO DO PRODUTO
{description}

ESTRATÉGIA POR PLATAFORMA
{platform_strategy}

FORMATOS SUGERIDOS
- Instagram: Reels + Stories + Feed post
- Facebook: Carrossel + Link post
- LinkedIn: Artigo + Post informativo
- WhatsApp: Card informativo + Link
- YouTube: Shorts + Vídeo explicativo
- Portal Mais Agro: Blog post + Webstories

CONTATOS E OBSERVAÇões
- Validar com especialista técnico
- Checar disponibilidade de imagens/vídeos
- Incluir CTA para portal Mais Agro
- Seguir guidelines de marca
- Revisar compliance regulatório

DATA PREVISTA: {data_input.strftime('%d/%m/%Y')}
FORMATO PRINCIPAL: {formato_principal}
"""
    return briefing

# --- Interface Principal ---
st.sidebar.title(f"🤖 Bem-vindo, {get_current_user().get('nome', 'Usuário')}!")
st.sidebar.info(f"**Squad:** {get_current_squad()}")
st.sidebar.info(f"**Agente selecionado:** {agente_selecionado['nome']}")

# Botão de logout na sidebar
if st.sidebar.button("🚪 Sair", key="logout_btn"):
    for key in ["logged_in", "user", "admin_password_correct", "admin_user", "agente_selecionado"]:
        if key in st.session_state:
            del st.session_state[key]
    st.rerun()

# Botão para trocar agente
if st.sidebar.button("🔄 Trocar Agente", key="trocar_agente_global"):
    st.session_state.agente_selecionado = None
    st.session_state.messages = []
    st.rerun()

# --- SELECTBOX PARA TROCAR AGENTE ACIMA DAS ABAS ---
st.title("🤖 Agente Social")

# Carregar agentes disponíveis
agentes = listar_agentes()

if agentes:
    # Preparar opções para o selectbox
    opcoes_agentes = []
    for agente in agentes:
        agente_completo = obter_agente_com_heranca(agente['_id'])
        if agente_completo:  # Só adiciona se tiver permissão
            descricao = f"{agente['nome']} - {agente.get('categoria', 'Social')}"
            if agente.get('agente_mae_id'):
                descricao += " 🔗"
            # Adicionar indicador de squad
            squad_permitido = agente.get('squad_permitido', 'Todos')
            descricao += f" 👥{squad_permitido}"
            opcoes_agentes.append((descricao, agente_completo))
    
    if opcoes_agentes:
        # Encontrar o índice atual
        indice_atual = 0
        for i, (desc, agente) in enumerate(opcoes_agentes):
            if agente['_id'] == st.session_state.agente_selecionado['_id']:
                indice_atual = i
                break
        
        # Selectbox para trocar agente
        col1, col2 = st.columns([3, 1])
        with col1:
            novo_agente_desc = st.selectbox(
                "Selecionar Agente:",
                options=[op[0] for op in opcoes_agentes],
                index=indice_atual,
                key="selectbox_trocar_agente"
            )
        with col2:
            if st.button("🔄 Trocar", key="botao_trocar_agente"):
                # Encontrar o agente completo correspondente
                for desc, agente in opcoes_agentes:
                    if desc == novo_agente_desc:
                        st.session_state.agente_selecionado = agente
                        st.session_state.messages = []
                        st.success(f"✅ Agente alterado para '{agente['nome']}'!")
                        st.rerun()
                        break
    else:
        st.info("Nenhum agente disponível com as permissões atuais.")

# Menu de abas - DETERMINAR QUAIS ABAS MOSTRAR
abas_base = [
    "💬 Chat", 
    "⚙️ Gerenciar Agentes", 
    "✅ Validação Unificada",
    "✨ Geração de Conteúdo",
    "📝 Revisão Ortográfica",
    "Monitoramento de Redes"
]

if is_syn_agent(agente_selecionado['nome']):
    abas_base.append("📋 Briefing")

# Criar abas dinamicamente
tabs = st.tabs(abas_base)

# Mapear abas para suas respectivas funcionalidades
tab_mapping = {}
for i, aba in enumerate(abas_base):
    tab_mapping[aba] = tabs[i]

# --- ABA: CHAT ---
with tab_mapping["💬 Chat"]:
    st.header("💬 Chat com Agente")
    
    # Inicializar session_state se não existir
    if 'messages' not in st.session_state:
        st.session_state.messages = []
    if 'segmentos_selecionados' not in st.session_state:
        st.session_state.segmentos_selecionados = []
    if 'show_historico' not in st.session_state:
        st.session_state.show_historico = False
    if 'modelo_chat' not in st.session_state:
        st.session_state.modelo_chat = "Gemini"
    
    agente = st.session_state.agente_selecionado
    st.subheader(f"Conversando com: {agente['nome']}")
    
    # Seletor de modelo na sidebar do chat
    st.sidebar.subheader("🤖 Configurações do Modelo")
    modelo_chat = st.sidebar.selectbox(
        "Escolha o modelo:",
        ["Gemini", "Claude"],
        key="modelo_chat_selector",
        index=0 if st.session_state.modelo_chat == "Gemini" else 1
    )
    st.session_state.modelo_chat = modelo_chat
    
    # Status dos modelos
    if modelo_chat == "Gemini" and not gemini_api_key:
        st.sidebar.error("❌ Gemini não disponível")
    elif modelo_chat == "Claude" and not anthropic_api_key:
        st.sidebar.error("❌ Claude não disponível")
    else:
        st.sidebar.success(f"✅ {modelo_chat} ativo")
    
    # Controles de navegação no topo
    col1, col2, col3 = st.columns([1, 1, 1])
    
    with col1:
        if st.button("📚 Carregar Histórico", key="carregar_historico"):
            st.session_state.show_historico = not st.session_state.show_historico
            st.rerun()
    
    with col2:
        if st.button("🔄 Limpar Chat", key="limpar_chat"):
            st.session_state.messages = []
            if hasattr(st.session_state, 'historico_contexto'):
                st.session_state.historico_contexto = []
            st.success("Chat limpo!")
            st.rerun()
    
    with col3:
        if st.button("🔁 Trocar Agente", key="trocar_agente_chat"):
            st.session_state.agente_selecionado = None
            st.session_state.messages = []
            st.session_state.historico_contexto = []
            st.rerun()
    
    # Mostrar se há histórico carregado
    if hasattr(st.session_state, 'historico_contexto') and st.session_state.historico_contexto:
        st.info(f"📖 Usando histórico anterior com {len(st.session_state.historico_contexto)} mensagens como contexto")
    
    # Modal para seleção de histórico
    if st.session_state.show_historico:
        with st.expander("📚 Selecionar Histórico de Conversa", expanded=True):
            conversas_anteriores = obter_conversas(agente['_id'])
            
            if conversas_anteriores:
                for i, conversa in enumerate(conversas_anteriores[:10]):  # Últimas 10 conversas
                    col_hist1, col_hist2, col_hist3 = st.columns([3, 1, 1])
                    
                    with col_hist1:
                        # CORREÇÃO: Usar get() para evitar KeyError
                        data_display = conversa.get('data_formatada', conversa.get('data', 'Data desconhecida'))
                        mensagens_count = len(conversa.get('mensagens', []))
                        st.write(f"**{data_display}** - {mensagens_count} mensagens")
                    
                    with col_hist2:
                        if st.button("👀 Visualizar", key=f"ver_{i}"):
                            st.session_state.conversa_visualizada = conversa.get('mensagens', [])
                    
                    with col_hist3:
                        if st.button("📥 Usar", key=f"usar_{i}"):
                            st.session_state.messages = conversa.get('mensagens', [])
                            st.session_state.historico_contexto = conversa.get('mensagens', [])
                            st.session_state.show_historico = False
                            st.success(f"✅ Histórico carregado: {len(conversa.get('mensagens', []))} mensagens")
                            st.rerun()
                
                # Visualizar conversa selecionada
                if hasattr(st.session_state, 'conversa_visualizada'):
                    st.subheader("👀 Visualização do Histórico")
                    for msg in st.session_state.conversa_visualizada[-6:]:  # Últimas 6 mensagens
                        with st.chat_message(msg.get("role", "user")):
                            st.markdown(msg.get("content", ""))
                    
                    if st.button("Fechar Visualização", key="fechar_visualizacao"):
                        st.session_state.conversa_visualizada = None
                        st.rerun()
            else:
                st.info("Nenhuma conversa anterior encontrada")
    
    # Mostrar informações de herança se aplicável
    if 'agente_mae_id' in agente and agente['agente_mae_id']:
        agente_original = obter_agente(agente['_id'])
        if agente_original and agente_original.get('herdar_elementos'):
            st.info(f"🔗 Este agente herda {len(agente_original['herdar_elementos'])} elementos do agente mãe")
    
    # Controles de segmentos na sidebar do chat
    st.sidebar.subheader("🔧 Configurações do Agente")
    st.sidebar.write("Selecione quais bases de conhecimento usar:")
    
    segmentos_disponiveis = {
        "Prompt do Sistema": "system_prompt",
        "Brand Guidelines": "base_conhecimento", 
        "Comentários do Cliente": "comments",
        "Planejamento": "planejamento"
    }
    
    segmentos_selecionados = []
    for nome, chave in segmentos_disponiveis.items():
        if st.sidebar.checkbox(nome, value=chave in st.session_state.segmentos_selecionados, key=f"seg_{chave}"):
            segmentos_selecionados.append(chave)
    
    st.session_state.segmentos_selecionados = segmentos_selecionados
    
    # Exibir status dos segmentos
    if segmentos_selecionados:
        st.sidebar.success(f"✅ Usando {len(segmentos_selecionados)} segmento(s)")
    else:
        st.sidebar.warning("⚠️ Nenhum segmento selecionado")
    
    # Indicador de posição na conversa
    if len(st.session_state.messages) > 4:
        st.caption(f"📄 Conversa com {len(st.session_state.messages)} mensagens")
    
    # CORREÇÃO: Exibir histórico de mensagens DENTRO do contexto correto
    # Verificar se messages existe e é iterável
    if hasattr(st.session_state, 'messages') and st.session_state.messages:
        for message in st.session_state.messages:
            # Verificar se message é um dicionário e tem a chave 'role'
            if isinstance(message, dict) and "role" in message:
                with st.chat_message(message["role"]):
                    st.markdown(message.get("content", ""))
            else:
                # Se a estrutura não for a esperada, pular esta mensagem
                continue
    else:
        # Se não houver mensagens, mostrar estado vazio
        st.info("💬 Inicie uma conversa digitando uma mensagem abaixo!")
    
    # Input do usuário
    if prompt := st.chat_input("Digite sua mensagem..."):
        # Adicionar mensagem do usuário ao histórico
        st.session_state.messages.append({"role": "user", "content": prompt})
        with st.chat_message("user"):
            st.markdown(prompt)
        
        # Construir contexto com segmentos selecionados
        contexto = construir_contexto(
            agente, 
            st.session_state.segmentos_selecionados, 
            st.session_state.messages
        )
        
        # Gerar resposta
        with st.chat_message("assistant"):
            with st.spinner('Pensando...'):
                try:
                    resposta = gerar_resposta_modelo(
                        contexto, 
                        st.session_state.modelo_chat,
                        contexto
                    )
                    st.markdown(resposta)
                    
                    # Adicionar ao histórico
                    st.session_state.messages.append({"role": "assistant", "content": resposta})
                    
                    # Salvar conversa com segmentos utilizados
                    salvar_conversa(
                        agente['_id'], 
                        st.session_state.messages,
                        st.session_state.segmentos_selecionados
                    )
                    
                except Exception as e:
                    st.error(f"Erro ao gerar resposta: {str(e)}")

# --- ABA: GERENCIAMENTO DE AGENTES (MODIFICADA PARA SQUADS) ---
with tab_mapping["⚙️ Gerenciar Agentes"]:
    st.header("Gerenciamento de Agentes")
    
    # Verificar autenticação apenas para gerenciamento
    current_user = get_current_user()
    current_squad = get_current_squad()
    
    if current_squad not in ["admin", "Syngenta", "SME", "Enterprise"]:
        st.warning("Acesso restrito a usuários autorizados")
    else:
        # Para admin, verificar senha adicional
        if current_squad == "admin":
            if not check_admin_password():
                st.warning("Digite a senha de administrador")
            else:
                st.write(f'Bem-vindo administrador!')
        else:
            st.write(f'Bem-vindo {current_user.get("nome", "Usuário")} do squad {current_squad}!')
            
        # Subabas para gerenciamento
        sub_tab1, sub_tab2, sub_tab3 = st.tabs(["Criar Agente", "Editar Agente", "Gerenciar Agentes"])
        
        with sub_tab1:
            st.subheader("Criar Novo Agente")
            
            with st.form("form_criar_agente"):
                nome_agente = st.text_input("Nome do Agente:")
                
                # Seleção de categoria - AGORA COM MONITORAMENTO
                categoria = st.selectbox(
                    "Categoria:",
                    ["Social", "SEO", "Conteúdo", "Monitoramento"],
                    help="Organize o agente por área de atuação"
                )
                
                # NOVO: Seleção de squad permitido
                squad_permitido = st.selectbox(
                    "Squad Permitido:",
                    ["Todos", "Syngenta", "SME", "Enterprise"],
                    help="Selecione qual squad pode ver e usar este agente"
                )
                
                # Configurações específicas para agentes de monitoramento
                if categoria == "Monitoramento":
                    st.info("🔍 **Agente de Monitoramento**: Este agente será usado apenas na aba de Monitoramento de Redes e terá uma estrutura simplificada.")
                    
                    # Para monitoramento, apenas base de conhecimento
                    base_conhecimento = st.text_area(
                        "Base de Conhecimento para Monitoramento:", 
                        height=300,
                        placeholder="""Cole aqui a base de conhecimento específica para monitoramento de redes sociais.

PERSONALIDADE: Especialista técnico do agronegócio com habilidade social - "Especialista que fala como gente"

TOM DE VOZ:
- Técnico, confiável e seguro, mas acessível
- Evita exageros e promessas vazias
- Sempre embasado em fatos e ciência
- Frases curtas e diretas, mais simpáticas
- Toque de leveza e ironia pontual quando o contexto permite

PRODUTOS SYN:
- Fortenza: Tratamento de sementes inseticida para Cerrado
- Verdatis: Inseticida com tecnologia PLINAZOLIN
- Megafol: Bioativador natural
- Miravis Duo: Fungicida para controle de manchas foliares

DIRETRIZES:
- NÃO inventar informações técnicas
- Sempre basear respostas em fatos
- Manter tom profissional mas acessível
- Adaptar resposta ao tipo de pergunta""",
                        help="Esta base será usada exclusivamente para monitoramento de redes sociais"
                    )
                    
                    # Campos específicos ocultos para monitoramento
                    system_prompt = ""
                    comments = ""
                    planejamento = ""
                    criar_como_filho = False
                    agente_mae_id = None
                    herdar_elementos = []
                    
                else:
                    # Para outras categorias, manter estrutura original
                    criar_como_filho = st.checkbox("Criar como agente filho (herdar elementos)")
                    
                    agente_mae_id = None
                    herdar_elementos = []
                    
                    if criar_como_filho:
                        # Listar TODOS os agentes disponíveis para herança (exceto monitoramento)
                        agentes_mae = listar_agentes_para_heranca()
                        agentes_mae = [agente for agente in agentes_mae if agente.get('categoria') != 'Monitoramento']
                        
                        if agentes_mae:
                            agente_mae_options = {f"{agente['nome']} ({agente.get('categoria', 'Social')})": agente['_id'] for agente in agentes_mae}
                            agente_mae_selecionado = st.selectbox(
                                "Agente Mãe:",
                                list(agente_mae_options.keys()),
                                help="Selecione o agente do qual este agente irá herdar elementos"
                            )
                            agente_mae_id = agente_mae_options[agente_mae_selecionado]
                            
                            st.subheader("Elementos para Herdar")
                            herdar_elementos = st.multiselect(
                                "Selecione os elementos a herdar do agente mãe:",
                                ["system_prompt", "base_conhecimento", "comments", "planejamento"],
                                help="Estes elementos serão herdados do agente mãe se não preenchidos abaixo"
                            )
                        else:
                            st.info("Nenhum agente disponível para herança. Crie primeiro um agente mãe.")
                    
                    system_prompt = st.text_area("Prompt de Sistema:", height=150, 
                                                placeholder="Ex: Você é um assistente especializado em...",
                                                help="Deixe vazio se for herdar do agente mãe")
                    base_conhecimento = st.text_area("Brand Guidelines:", height=200,
                                                   placeholder="Cole aqui informações, diretrizes, dados...",
                                                   help="Deixe vazio se for herdar do agente mãe")
                    comments = st.text_area("Comentários do cliente:", height=200,
                                                   placeholder="Cole aqui os comentários de ajuste do cliente (Se houver)",
                                                   help="Deixe vazio se for herdar do agente mãe")
                    planejamento = st.text_area("Planejamento:", height=200,
                                               placeholder="Estratégias, planejamentos, cronogramas...",
                                               help="Deixe vazio se for herdar do agente mãe")
                
                submitted = st.form_submit_button("Criar Agente")
                if submitted:
                    if nome_agente:
                        agente_id = criar_agente(
                            nome_agente, 
                            system_prompt, 
                            base_conhecimento, 
                            comments, 
                            planejamento,
                            categoria,
                            squad_permitido,  # Novo campo
                            agente_mae_id if criar_como_filho else None,
                            herdar_elementos if criar_como_filho else []
                        )
                        st.success(f"Agente '{nome_agente}' criado com sucesso na categoria {categoria} para o squad {squad_permitido}!")
                    else:
                        st.error("Nome é obrigatório!")
        
        with sub_tab2:
            st.subheader("Editar Agente Existente")
            
            agentes = listar_agentes()
            if agentes:
                agente_options = {agente['nome']: agente for agente in agentes}
                agente_selecionado_nome = st.selectbox("Selecione o agente para editar:", 
                                                     list(agente_options.keys()))
                
                if agente_selecionado_nome:
                    agente = agente_options[agente_selecionado_nome]
                    
                    with st.form("form_editar_agente"):
                        novo_nome = st.text_input("Nome do Agente:", value=agente['nome'])
                        
                        # Categoria - AGORA COM MONITORAMENTO
                        categorias_disponiveis = ["Social", "SEO", "Conteúdo", "Monitoramento"]
                        if agente.get('categoria') in categorias_disponiveis:
                            index_categoria = categorias_disponiveis.index(agente.get('categoria', 'Social'))
                        else:
                            index_categoria = 0
                            
                        nova_categoria = st.selectbox(
                            "Categoria:",
                            categorias_disponiveis,
                            index=index_categoria,
                            help="Organize o agente por área de atuação"
                        )
                        
                        # NOVO: Squad permitido
                        squads_disponiveis = ["Todos", "Syngenta", "SME", "Enterprise"]
                        squad_atual = agente.get('squad_permitido', 'Todos')
                        if squad_atual in squads_disponiveis:
                            index_squad = squads_disponiveis.index(squad_atual)
                        else:
                            index_squad = 0
                            
                        novo_squad_permitido = st.selectbox(
                            "Squad Permitido:",
                            squads_disponiveis,
                            index=index_squad,
                            help="Selecione qual squad pode ver e usar este agente"
                        )
                        
                        # Interface diferente para agentes de monitoramento
                        if nova_categoria == "Monitoramento":
                            st.info("🔍 **Agente de Monitoramento**: Este agente será usado apenas na aba de Monitoramento de Redes.")
                            
                            # Para monitoramento, apenas base de conhecimento
                            nova_base = st.text_area(
                                "Base de Conhecimento para Monitoramento:", 
                                value=agente.get('base_conhecimento', ''),
                                height=300,
                                help="Esta base será usada exclusivamente para monitoramento de redes sociais"
                            )
                            
                            # Campos específicos ocultos para monitoramento
                            novo_prompt = ""
                            nova_comment = ""
                            novo_planejamento = ""
                            agente_mae_id = None
                            herdar_elementos = []
                            
                            # Remover herança se existir
                            if agente.get('agente_mae_id'):
                                st.warning("⚠️ Agentes de monitoramento não suportam herança. A herança será removida.")
                            
                        else:
                            # Para outras categorias, manter estrutura original
                            
                            # Informações de herança (apenas se não for monitoramento)
                            if agente.get('agente_mae_id'):
                                agente_mae = obter_agente(agente['agente_mae_id'])
                                if agente_mae:
                                    st.info(f"🔗 Este agente é filho de: {agente_mae['nome']}")
                                    st.write(f"Elementos herdados: {', '.join(agente.get('herdar_elementos', []))}")
                            
                            # Opção para tornar independente
                            if agente.get('agente_mae_id'):
                                tornar_independente = st.checkbox("Tornar agente independente (remover herança)")
                                if tornar_independente:
                                    agente_mae_id = None
                                    herdar_elementos = []
                                else:
                                    agente_mae_id = agente.get('agente_mae_id')
                                    herdar_elementos = agente.get('herdar_elementos', [])
                            else:
                                agente_mae_id = None
                                herdar_elementos = []
                                # Opção para adicionar herança
                                adicionar_heranca = st.checkbox("Adicionar herança de agente mãe")
                                if adicionar_heranca:
                                    # Listar TODOS os agentes disponíveis para herança (excluindo o próprio e monitoramento)
                                    agentes_mae = listar_agentes_para_heranca(agente['_id'])
                                    agentes_mae = [agente_mae for agente_mae in agentes_mae if agente_mae.get('categoria') != 'Monitoramento']
                                    
                                    if agentes_mae:
                                        agente_mae_options = {f"{agente_mae['nome']} ({agente_mae.get('categoria', 'Social')})": agente_mae['_id'] for agente_mae in agentes_mae}
                                        if agente_mae_options:
                                            agente_mae_selecionado = st.selectbox(
                                                "Agente Mãe:",
                                                list(agente_mae_options.keys()),
                                                help="Selecione o agente do qual este agente irá herdar elementos"
                                            )
                                            agente_mae_id = agente_mae_options[agente_mae_selecionado]
                                            herdar_elementos = st.multiselect(
                                                "Elementos para herdar:",
                                                ["system_prompt", "base_conhecimento", "comments", "planejamento"],
                                                default=herdar_elementos
                                            )
                                        else:
                                            st.info("Nenhum agente disponível para herança.")
                                    else:
                                        st.info("Nenhum agente disponível para herança.")
                            
                            novo_prompt = st.text_area("Prompt de Sistema:", value=agente['system_prompt'], height=150)
                            nova_base = st.text_area("Brand Guidelines:", value=agente.get('base_conhecimento', ''), height=200)
                            nova_comment = st.text_area("Comentários:", value=agente.get('comments', ''), height=200)
                            novo_planejamento = st.text_area("Planejamento:", value=agente.get('planejamento', ''), height=200)
                        
                        submitted = st.form_submit_button("Atualizar Agente")
                        if submitted:
                            if novo_nome:
                                atualizar_agente(
                                    agente['_id'], 
                                    novo_nome, 
                                    novo_prompt, 
                                    nova_base, 
                                    nova_comment, 
                                    novo_planejamento,
                                    nova_categoria,
                                    novo_squad_permitido,  # Novo campo
                                    agente_mae_id,
                                    herdar_elementos
                                )
                                st.success(f"Agente '{novo_nome}' atualizado com sucesso!")
                                st.rerun()
                            else:
                                st.error("Nome é obrigatório!")
            else:
                st.info("Nenhum agente criado ainda.")
        
        with sub_tab3:
            st.subheader("Gerenciar Agentes")
            
            # Mostrar informações do usuário atual
            current_squad = get_current_squad()
            if current_squad == "admin":
                st.info("👑 Modo Administrador: Visualizando todos os agentes do sistema")
            else:
                st.info(f"👤 Visualizando agentes do squad {current_squad} e squad 'Todos'")
            
            # Filtros por categoria - AGORA COM MONITORAMENTO
            categorias = ["Todos", "Social", "SEO", "Conteúdo", "Monitoramento"]
            categoria_filtro = st.selectbox("Filtrar por categoria:", categorias)
            
            agentes = listar_agentes()
            
            # Aplicar filtro
            if categoria_filtro != "Todos":
                agentes = [agente for agente in agentes if agente.get('categoria') == categoria_filtro]
            
            if agentes:
                for i, agente in enumerate(agentes):
                    with st.expander(f"{agente['nome']} - {agente.get('categoria', 'Social')} - Squad: {agente.get('squad_permitido', 'Todos')} - Criado em {agente['data_criacao'].strftime('%d/%m/%Y')}"):
                        
                        # Mostrar proprietário se for admin
                        owner_info = ""
                        if current_squad == "admin" and agente.get('criado_por'):
                            owner_info = f" | 👤 {agente['criado_por']}"
                            st.write(f"**Proprietário:** {agente['criado_por']}")
                            st.write(f"**Squad do Criador:** {agente.get('criado_por_squad', 'N/A')}")
                        
                        # Mostrar informações específicas por categoria
                        if agente.get('categoria') == 'Monitoramento':
                            st.info("🔍 **Agente de Monitoramento** - Usado apenas na aba de Monitoramento de Redes")
                            
                            if agente.get('base_conhecimento'):
                                st.write(f"**Base de Conhecimento:** {agente['base_conhecimento'][:200]}...")
                            else:
                                st.warning("⚠️ Base de conhecimento não configurada")
                            
                            # Agentes de monitoramento não mostram outros campos
                            st.write("**System Prompt:** (Não utilizado em monitoramento)")
                            st.write("**Comentários:** (Não utilizado em monitoramento)")
                            st.write("**Planejamento:** (Não utilizado em monitoramento)")
                            
                        else:
                            # Para outras categorias, mostrar estrutura completa
                            if agente.get('agente_mae_id'):
                                agente_mae = obter_agente(agente['agente_mae_id'])
                                if agente_mae:
                                    st.write(f"**🔗 Herda de:** {agente_mae['nome']}")
                                    st.write(f"**Elementos herdados:** {', '.join(agente.get('herdar_elementos', []))}")
                            
                            st.write(f"**Prompt de Sistema:** {agente['system_prompt'][:100]}..." if agente['system_prompt'] else "**Prompt de Sistema:** (herdado ou vazio)")
                            if agente.get('base_conhecimento'):
                                st.write(f"**Brand Guidelines:** {agente['base_conhecimento'][:200]}...")
                            if agente.get('comments'):
                                st.write(f"**Comentários do cliente:** {agente['comments'][:200]}...")
                            if agente.get('planejamento'):
                                st.write(f"**Planejamento:** {agente['planejamento'][:200]}...")
                        
                        col1, col2 = st.columns(2)
                        with col1:
                            if st.button("Selecionar para Chat", key=f"select_{i}"):
                                agente_completo = obter_agente_com_heranca(agente['_id'])
                                st.session_state.agente_selecionado = agente_completo
                                st.session_state.messages = []
                                st.success(f"Agente '{agente['nome']}' selecionado!")
                                st.rerun()
                        with col2:
                            if st.button("Desativar", key=f"delete_{i}"):
                                desativar_agente(agente['_id'])
                                st.success(f"Agente '{agente['nome']}' desativado!")
                                st.rerun()
            else:
                st.info("Nenhum agente encontrado para esta categoria.")

if "📋 Briefing" in tab_mapping:
    with tab_mapping["📋 Briefing"]:
        st.header("📋 Gerador de Briefings - SYN")
        st.markdown("Digite o conteúdo da célula do calendário para gerar um briefing completo no padrão SYN.")
        
        # Abas para diferentes modos de operação
        tab1, tab2 = st.tabs(["Briefing Individual", "Processamento em Lote (CSV)"])
        
        with tab1:
            st.markdown("### Digite o conteúdo da célula do calendário")

            content_input = st.text_area(
                "Conteúdo da célula:",
                placeholder="Ex: megafol - série - potencial máximo, todo o tempo",
                height=100,
                help="Cole aqui o conteúdo exato da célula do calendário do Sheets",
                key="individual_content"
            )

            # Campos opcionais para ajuste
            col1, col2 = st.columns(2)

            with col1:
                data_input = st.date_input("Data prevista:", value=datetime.datetime.now(), key="individual_date")

            with col2:
                formato_principal = st.selectbox(
                    "Formato principal:",
                    ["Reels + capa", "Carrossel + stories", "Blog + redes", "Vídeo + stories", "Multiplataforma"],
                    key="individual_format"
                )

            generate_btn = st.button("Gerar Briefing Individual", type="primary", key="individual_btn")

            # Processamento e exibição do briefing individual
            if generate_btn and content_input:
                with st.spinner("Analisando conteúdo e gerando briefing..."):
                    # Extrair informações do produto
                    product, culture, action = extract_product_info(content_input)
                    
                    if product and product in PRODUCT_DESCRIPTIONS:
                        # Gerar briefing completo
                        briefing = generate_briefing(content_input, product, culture, action, data_input, formato_principal)
                        
                        # Exibir briefing
                        st.markdown("## Briefing Gerado")
                        st.text(briefing)
                        
                        # Botão de download
                        st.download_button(
                            label="Baixar Briefing",
                            data=briefing,
                            file_name=f"briefing_{product}_{data_input.strftime('%Y%m%d')}.txt",
                            mime="text/plain",
                            key="individual_download"
                        )
                        
                        # Informações extras
                        with st.expander("Informações Extraídas"):
                            st.write(f"Produto: {product}")
                            st.write(f"Cultura: {culture}")
                            st.write(f"Ação: {action}")
                            st.write(f"Data: {data_input.strftime('%d/%m/%Y')}")
                            st.write(f"Formato principal: {formato_principal}")
                            st.write(f"Descrição: {PRODUCT_DESCRIPTIONS[product]}")
                            
                    elif product:
                        st.warning(f"Produto '{product}' não encontrado no dicionário. Verifique a grafia.")
                        st.info("Produtos disponíveis: " + ", ".join(list(PRODUCT_DESCRIPTIONS.keys())[:10]) + "...")
                    else:
                        st.error("Não foi possível identificar um produto no conteúdo. Tente formatos como:")
                        st.code("""
                        megafol - série - potencial máximo, todo o tempo
                        verdavis - soja - depoimento produtor
                        engeo pleno s - milho - controle percevejo
                        miravis duo - algodão - reforço preventivo
                        """)

        with tab2:
            st.markdown("### Processamento em Lote via CSV")
            
            st.info("""
            Faça upload de um arquivo CSV exportado do Google Sheets.
            O sistema irá processar cada linha a partir da segunda linha (ignorando cabeçalhos)
            e gerar briefings apenas para as linhas que contêm produtos reconhecidos.
            """)
            
            uploaded_file = st.file_uploader(
                "Escolha o arquivo CSV", 
                type=['csv'],
                help="Selecione o arquivo CSV exportado do Google Sheets"
            )
            
            if uploaded_file is not None:
                try:
                    # Ler o CSV
                    df = pd.read_csv(uploaded_file)
                    st.success(f"CSV carregado com sucesso! {len(df)} linhas encontradas.")
                    
                    # Mostrar prévia do arquivo
                    with st.expander("Visualizar primeiras linhas do CSV"):
                        st.dataframe(df.head())
                    
                    # Configurações para processamento em lote
                    st.markdown("### Configurações do Processamento em Lote")
                    col1, col2 = st.columns(2)
                    
                    with col1:
                        data_padrao = st.date_input(
                            "Data padrão para todos os briefings:",
                            value=datetime.datetime.now(),
                            key="batch_date"
                        )
                    
                    with col2:
                        formato_padrao = st.selectbox(
                            "Formato principal padrão:",
                            ["Reels + capa", "Carrossel + stories", "Blog + redes", "Vídeo + stories", "Multiplataforma"],
                            key="batch_format"
                        )
                    
                    # Identificar coluna com conteúdo
                    colunas = df.columns.tolist()
                    coluna_conteudo = st.selectbox(
                        "Selecione a coluna que contém o conteúdo das células:",
                        colunas,
                        help="Selecione a coluna que contém os textos das células do calendário"
                    )
                    
                    processar_lote = st.button("Processar CSV e Gerar Briefings", type="primary", key="batch_btn")
                    
                    if processar_lote:
                        briefings_gerados = []
                        linhas_processadas = 0
                        linhas_com_produto = 0
                        
                        progress_bar = st.progress(0)
                        status_text = st.empty()
                        
                        for index, row in df.iterrows():
                            linhas_processadas += 1
                            progress_bar.progress(linhas_processadas / len(df))
                            status_text.text(f"Processando linha {linhas_processadas} de {len(df)}...")
                            
                            # Pular a primeira linha (cabeçalhos)
                            if index == 0:
                                continue
                            
                            # Obter conteúdo da célula
                            content = str(row[coluna_conteudo]) if pd.notna(row[coluna_conteudo]) else ""
                            
                            if content:
                                # Extrair informações do produto
                                product, culture, action = extract_product_info(content)
                                
                                if product and product in PRODUCT_DESCRIPTIONS:
                                    linhas_com_produto += 1
                                    # Gerar briefing
                                    briefing = generate_briefing(
                                        content, 
                                        product, 
                                        culture, 
                                        action, 
                                        data_padrao, 
                                        formato_padrao
                                    )
                                    
                                    briefings_gerados.append({
                                        'linha': index + 1,
                                        'produto': product,
                                        'conteudo': content,
                                        'briefing': briefing,
                                        'arquivo': f"briefing_{product}_{index+1}.txt"
                                    })
                        
                        progress_bar.empty()
                        status_text.empty()
                        
                        # Resultados do processamento
                        st.success(f"Processamento concluído! {linhas_com_produto} briefings gerados de {linhas_processadas-1} linhas processadas.")
                        
                        if briefings_gerados:
                            # Exibir resumo
                            st.markdown("### Briefings Gerados")
                            resumo_df = pd.DataFrame([{
                                'Linha': b['linha'],
                                'Produto': b['produto'],
                                'Conteúdo': b['conteudo'][:50] + '...' if len(b['conteudo']) > 50 else b['conteudo']
                            } for b in briefings_gerados])
                            
                            st.dataframe(resumo_df)
                            
                            # Criar arquivo ZIP com todos os briefings
                            import zipfile
                            from io import BytesIO
                            
                            zip_buffer = BytesIO()
                            with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
                                for briefing_info in briefings_gerados:
                                    zip_file.writestr(
                                        briefing_info['arquivo'], 
                                        briefing_info['briefing']
                                    )
                            
                            zip_buffer.seek(0)
                            
                            # Botão para download do ZIP
                            st.download_button(
                                label="📥 Baixar Todos os Briefings (ZIP)",
                                data=zip_buffer,
                                file_name="briefings_syn.zip",
                                mime="application/zip",
                                key="batch_download_zip"
                            )
                            
                            # Também permitir download individual
                            st.markdown("---")
                            st.markdown("### Download Individual")
                            
                            for briefing_info in briefings_gerados:
                                col1, col2 = st.columns([3, 1])
                                with col1:
                                    st.text(f"Linha {briefing_info['linha']}: {briefing_info['produto']} - {briefing_info['conteudo'][:30]}...")
                                with col2:
                                    st.download_button(
                                        label="📄 Baixar",
                                        data=briefing_info['briefing'],
                                        file_name=briefing_info['arquivo'],
                                        mime="text/plain",
                                        key=f"download_{briefing_info['linha']}"
                                    )
                        else:
                            st.warning("Nenhum briefing foi gerado. Verifique se o CSV contém produtos reconhecidos.")
                            st.info("Produtos reconhecidos: " + ", ".join(list(PRODUCT_DESCRIPTIONS.keys())[:15]) + "...")
                            
                except Exception as e:
                    st.error(f"Erro ao processar o arquivo CSV: {str(e)}")

        # Seção de exemplos
        with st.expander("Exemplos de Conteúdo", expanded=True):
            st.markdown("""
            Formatos Reconhecidos:

            Padrão: PRODUTO - CULTURA - AÇÃO ou PRODUTO - AÇÃO

            Exemplos:
            - megafol - série - potencial máximo, todo o tempo
            - verdavis - milho - resultados do produto
            - engeo pleno s - soja - resultados GTEC
            - miravis duo - algodão - depoimento produtor
            - axial - trigo - reforço pós-emergente
            - manejo limpo - importância manejo antecipado
            - certano HF - a jornada de certano
            - elestal neo - soja - depoimento de produtor
            - fortenza - a jornada da semente mais forte - EP 01
            - reverb - vídeo conceito
            """)

        # Lista de produtos reconhecidos
        with st.expander("Produtos Reconhecidos"):
            col1, col2, col3 = st.columns(3)
            products = list(PRODUCT_DESCRIPTIONS.keys())
            
            with col1:
                for product in products[:10]:
                    st.write(f"• {product}")
            
            with col2:
                for product in products[10:20]:
                    st.write(f"• {product}")
            
            with col3:
                for product in products[20:]:
                    st.write(f"• {product}")

        # Rodapé
        st.markdown("---")
        st.caption("Ferramenta de geração automática de briefings - Padrão SYN. Digite o conteúdo da célula do calendário para gerar briefings completos.")

def criar_analisadores_especialistas(contexto_agente, contexto_global):
    """Cria prompts especializados para cada área de análise"""
    
    analisadores = {
        'ortografia': {
            'nome': '🔤 Especialista em Ortografia e Gramática',
            'prompt': f"""
{contexto_global}

## FUNÇÃO: ESPECIALISTA EM ORTOGRAFIA E GRAMÁTICA PORTUGUÊS BR

**Sua tarefa:** Analisar EXCLUSIVAMENTE aspectos ortográficos e gramaticais.

### CRITÉRIOS DE ANÁLISE:
1. **Ortografia** - Erros de escrita
2. **Gramática** - Concordância, regência, colocação
3. **Pontuação** - Uso de vírgulas, pontos, etc.
4. **Acentuação** - Erros de acentuação
5. **Padrão Culto** - Conformidade com norma culta

### FORMATO DE RESPOSTA OBRIGATÓRIO:

## 🔤 RELATÓRIO ORTOGRÁFICO

### ✅ ACERTOS
- [Itens corretos]

### ❌ ERROS IDENTIFICADOS
- [Lista específica de erros com correções]


### 💡 SUGESTÕES DE MELHORIA
- [Recomendações específicas]
"""
        },
        'lexico': {
            'nome': '📚 Especialista em Léxico e Vocabulário',
            'prompt': f"""
{contexto_agente}
{contexto_global}

## FUNÇÃO: ESPECIALISTA EM LÉXICO E VOCABULÁRIO

**Sua tarefa:** Analisar EXCLUSIVAMENTE aspectos lexicais e de vocabulário.

### CRITÉRIOS DE ANÁLISE:
1. **Variedade Lexical** - Riqueza de vocabulário
2. **Precisão Semântica** - Uso adequado das palavras
3. **Repetição** - Palavras ou expressões repetidas em excesso
4. **Jargões** - Uso inadequado de termos técnicos
5. **Clareza** - Facilidade de compreensão

### FORMATO DE RESPOSTA OBRIGATÓRIO:

## 📚 RELATÓRIO LEXICAL

### ✅ VOCABULÁRIO ADEQUADO
- [Pontos fortes do vocabulário]

### ⚠️ ASPECTOS A MELHORAR
- [Problemas lexicais identificados]

### 🔄 SUGESTÕES DE SINÔNIMOS
- [Palavras para substituir]

"""
        },
        'branding': {
            'nome': '🎨 Especialista em Branding e Identidade',
            'prompt': f"""
{contexto_agente}
{contexto_global}

## FUNÇÃO: ESPECIALISTA EM BRANDING E IDENTIDADE

**Sua tarefa:** Analisar EXCLUSIVAMENTE conformidade com diretrizes de branding.

### CRITÉRIOS DE ANÁLISE:
1. **Tom de Voz** - Alinhamento com personalidade da marca
2. **Mensagem Central** - Consistência da mensagem
3. **Valores da Marca** - Reflexo dos valores organizacionais
4. **Público-Alvo** - Adequação ao público pretendido
5. **Diferenciação** - Elementos únicos da marca

### FORMATO DE RESPOSTA OBRIGATÓRIO:

## 🎨 RELATÓRIO DE BRANDING

### ✅ ALINHAMENTOS
- [Elementos que seguem as diretrizes]

### ❌ DESVIOS IDENTIFICADOS
- [Elementos fora do padrão da marca]


### 💡 RECOMENDAÇÕES ESTRATÉGICAS
- [Sugestões para melhor alinhamento]
"""
        
        
        }
    }
    
    return analisadores

def executar_analise_especializada(texto, nome_arquivo, analisadores):
    """Executa análise com múltiplos especialistas"""
    
    resultados = {}
    
    for area, config in analisadores.items():
        with st.spinner(f"Executando {config['nome']}..."):
            try:
                prompt_completo = f"""
{config['prompt']}

###BEGIN TEXTO PARA ANÁLISE###
**Arquivo:** {nome_arquivo}
**Conteúdo:**
{texto[:8000]}
###END TEXTO PARA ANÁLISE###

Por favor, forneça sua análise no formato solicitado.
"""
                
                resposta = modelo_texto.generate_content(prompt_completo)
                resultados[area] = {
                    'nome': config['nome'],
                    'analise': resposta.text,
                }
                
            except Exception as e:
                resultados[area] = {
                    'nome': config['nome'],
                    'analise': f"❌ Erro na análise: {str(e)}",
                    'score': 0
                }
    
    return resultados

def extrair_score(texto_analise):
    """Extrai score numérico do texto de análise"""
    import re
    padrao = r'SCORE.*?\[(\d+)(?:/10)?\]'
    correspondencias = re.findall(padrao, texto_analise, re.IGNORECASE)
    if correspondencias:
        return int(correspondencias[0])
    return 5  # Score padrão se não encontrar

def gerar_relatorio_consolidado(resultados_especialistas, nome_arquivo):
    """Gera relatório consolidado a partir das análises especializadas"""
    
   
    
    relatorio = f"""
# 📊 RELATÓRIO CONSOLIDADO DE VALIDAÇÃO

**Documento:** {nome_arquivo}
**Data da Análise:** {datetime.datetime.now().strftime('%d/%m/%Y %H:%M')}

"""
    
    # Adicionar scores individuais
    for area, resultado in resultados_especialistas.items():
        emoji = "✅" if resultado['score'] >= 8 else "⚠️" if resultado['score'] >= 6 else "❌"
        relatorio += f"- {emoji} **{resultado['nome']}:** {resultado['score']}/10\n"
    
    relatorio += "\n## 📋 ANÁLISES DETALHADAS POR ESPECIALISTA\n"
    
    # Adicionar análises detalhadas
    for area, resultado in resultados_especialistas.items():
        relatorio += f"\n### {resultado['nome']}\n"
        relatorio += f"{resultado['analise']}\n"
        relatorio += "---\n"
    
    # Resumo executivo
    relatorio += f"""
## 🚀 RESUMO EXECUTIVO


### 🎯 PRÓXIMOS PASSOS RECOMENDADOS:
"""
    
    # Recomendações baseadas nos scores
    areas_baixas = [area for area, resultado in resultados_especialistas.items() if resultado['score'] < 6]
    if areas_baixas:
        relatorio += f"- **Prioridade:** Focar em {', '.join(areas_baixas)}\n"
    
    areas_medianas = [area for area, resultado in resultados_especialistas.items() if 6 <= resultado['score'] < 8]
    if areas_medianas:
        relatorio += f"- **Otimização:** Melhorar {', '.join(areas_medianas)}\n"
    
    relatorio += "- **Manutenção:** Manter as áreas com scores altos\n"
    
    return relatorio

# --- FUNÇÕES ORIGINAIS MANTIDAS ---

def criar_prompt_validacao_preciso(texto, nome_arquivo, contexto_agente):
    """Cria um prompt de validação muito mais preciso para evitar falsos positivos"""
    
    prompt = f"""
{contexto_agente}

###BEGIN TEXTO PARA VALIDAÇÃO###
**Arquivo:** {nome_arquivo}
**Conteúdo:**
{texto[:12000]}
###END TEXTO PARA VALIDAÇÃO###

## FORMATO DE RESPOSTA OBRIGATÓRIO:

### ✅ CONFORMIDADE COM DIRETRIZES
- [Itens que estão alinhados com as diretrizes de branding]

**INCONSISTÊNCIAS COM BRANDING:**
- [Só liste desvios REAIS das diretrizes de branding]

### 💡 TEXTO REVISADO
- [Sugestões para aprimorar]

### 📊 STATUS FINAL
**Documento:** [Aprovado/Necessita ajustes/Reprovado]
**Principais ações necessárias:** [Lista resumida]
"""
    return prompt

def analisar_documento_por_slides(doc, contexto_agente):
    """Analisa documento slide por slide com alta precisão"""
    
    resultados = []
    
    for i, slide in enumerate(doc['slides']):
        with st.spinner(f"Analisando slide {i+1}..."):
            try:
                prompt_slide = f"""
{contexto_agente}

## ANÁLISE POR SLIDE - PRECISÃO ABSOLUTA

###BEGIN TEXTO PARA VALIDAÇÃO###
**SLIDE {i+1}:**
{slide['conteudo'][:2000]}
###END TEXTO PARA VALIDAÇÃO###

**ANÁLISE DO SLIDE {i+1}:**

### ✅ Pontos Fortes:
[O que está bom neste slide]

### ⚠️ Problemas REAIS:
- [Lista CURTA de problemas]

### 💡 Sugestões Específicas:
[Melhorias para ESTE slide específico]

Considere que slides que são introdutórios ou apenas de títulos não precisam de tanto rigor de branding

**STATUS:** [✔️ Aprovado / ⚠️ Ajustes Menores / ❌ Problemas Sérios]
"""
                
                resposta = modelo_texto.generate_content(prompt_slide)
                resultados.append({
                    'slide_num': i+1,
                    'analise': resposta.text,
                    'tem_alteracoes': '❌' in resposta.text or '⚠️' in resposta.text
                })
                
            except Exception as e:
                resultados.append({
                    'slide_num': i+1,
                    'analise': f"❌ Erro na análise do slide: {str(e)}",
                    'tem_alteracoes': False
                })
    
    # Construir relatório consolidado
    relatorio = f"# 📊 RELATÓRIO DE VALIDAÇÃO - {doc['nome']}\n\n"
    relatorio += f"**Total de Slides:** {len(doc['slides'])}\n"
    relatorio += f"**Slides com Alterações:** {sum(1 for r in resultados if r['tem_alteracoes'])}\n\n"
    
    # Slides que precisam de atenção
    slides_com_problemas = [r for r in resultados if r['tem_alteracoes']]
    if slides_com_problemas:
        relatorio += "## 🚨 SLIDES QUE PRECISAM DE ATENÇÃO:\n\n"
        for resultado in slides_com_problemas:
            relatorio += f"### 📋 Slide {resultado['slide_num']}\n"
            relatorio += f"{resultado['analise']}\n\n"
    
    # Resumo executivo
    relatorio += "## 📈 RESUMO EXECUTIVO\n\n"
    if slides_com_problemas:
        relatorio += f"**⚠️ {len(slides_com_problemas)} slide(s) necessitam de ajustes**\n"
        relatorio += f"**✅ {len(doc['slides']) - len(slides_com_problemas)} slide(s) estão adequados**\n"
    else:
        relatorio += "**🎉 Todos os slides estão em conformidade com as diretrizes!**\n"
    
    return relatorio

def extract_text_from_pdf_com_slides(arquivo_pdf):
    """Extrai texto de PDF com informação de páginas"""
    try:
        import PyPDF2
        pdf_reader = PyPDF2.PdfReader(arquivo_pdf)
        slides_info = []
        
        for pagina_num, pagina in enumerate(pdf_reader.pages):
            texto = pagina.extract_text()
            slides_info.append({
                'numero': pagina_num + 1,
                'conteudo': texto,
                'tipo': 'página'
            })
        
        texto_completo = "\n\n".join([f"--- PÁGINA {s['numero']} ---\n{s['conteudo']}" for s in slides_info])
        return texto_completo, slides_info
        
    except Exception as e:
        return f"Erro na extração PDF: {str(e)}", []

def extract_text_from_pptx_com_slides(arquivo_pptx):
    """Extrai texto de PPTX com informação de slides"""
    try:
        from pptx import Presentation
        import io
        
        prs = Presentation(io.BytesIO(arquivo_pptx.read()))
        slides_info = []
        
        for slide_num, slide in enumerate(prs.slides):
            texto_slide = f"--- SLIDE {slide_num + 1} ---\n"
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texto_slide += shape.text + "\n"
            
            slides_info.append({
                'numero': slide_num + 1,
                'conteudo': texto_slide,
                'tipo': 'slide'
            })
        
        texto_completo = "\n\n".join([s['conteudo'] for s in slides_info])
        return texto_completo, slides_info
        
    except Exception as e:
        return f"Erro na extração PPTX: {str(e)}", []

def extrair_texto_arquivo(arquivo):
    """Extrai texto de arquivos TXT e DOCX"""
    try:
        if arquivo.type == "text/plain":
            return str(arquivo.read(), "utf-8")
        elif arquivo.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            import docx
            import io
            doc = docx.Document(io.BytesIO(arquivo.read()))
            texto = ""
            for para in doc.paragraphs:
                texto += para.text + "\n"
            return texto
        else:
            return f"Tipo não suportado: {arquivo.type}"
    except Exception as e:
        return f"Erro na extração: {str(e)}"

def extract_text_from_pdf(pdf_path):
    """
    Extract text from a PDF file using multiple methods for better coverage
    """
    text = ""

    # Method 1: Try with pdfplumber (better for some PDFs)
    try:
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text
    except Exception as e:
        print(f"pdfplumber failed for {pdf_path}: {e}")

    # Method 2: Fallback to PyPDF2 if pdfplumber didn't extract much text
    if len(text.strip()) < 100:  # If very little text was extracted
        try:
            with open(pdf_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text 
        except Exception as e:
            print(f"PyPDF2 also failed for {pdf_path}: {e}")

    return text

def criar_analisadores_imagem(contexto_agente, contexto_global):
    """Cria analisadores especializados para imagens"""
    
    analisadores = {
        'composicao_visual': {
            'nome': '🎨 Especialista em Composição Visual',
            'prompt': f"""
{contexto_agente}
{contexto_global}

## FUNÇÃO: ESPECIALISTA EM COMPOSIÇÃO VISUAL

**Sua tarefa:** Analisar EXCLUSIVAMENTE a composição visual da imagem.

### CRITÉRIOS DE ANÁLISE:
1. **Balanceamento** - Distribuição equilibrada dos elementos
2. **Hierarquia Visual** - Foco e pontos de atenção
3. **Espaçamento** - Uso adequado do espaço
4. **Proporções** - Relação entre elementos visuais
5. **Harmonia** - Conjunto visual coeso

### FORMATO DE RESPOSTA OBRIGATÓRIO:

## 🎨 RELATÓRIO DE COMPOSIÇÃO VISUAL

### ✅ PONTOS FORTES DA COMPOSIÇÃO
- [Elementos bem compostos]

### ⚠️ PROBLEMAS DE COMPOSIÇÃO
- [Issues de organização visual]

### 📊 SCORE COMPOSIÇÃO: [X/10]

### 💡 SUGESTÕES DE MELHORIA VISUAL
- [Recomendações para melhor composição]
"""
        },
        'cores_branding': {
            'nome': '🌈 Especialista em Cores e Branding',
            'prompt': f"""
{contexto_agente}
{contexto_global}

## FUNÇÃO: ESPECIALISTA EM CORES E BRANDING

**Sua tarefa:** Analisar EXCLUSIVAMENTE cores e alinhamento com branding.

### CRITÉRIOS DE ANÁLISE:
1. **Paleta de Cores** - Cores utilizadas na imagem
2. **Contraste** - Legibilidade e visibilidade
3. **Consistência** - Coerência com identidade visual
4. **Psicologia das Cores** - Efeito emocional das cores
5. **Acessibilidade** - Visibilidade para diferentes usuários

### FORMATO DE RESPOSTA OBRIGATÓRIO:

## 🌈 RELATÓRIO DE CORES E BRANDING

### ✅ CORES ALINHADAS
- [Cores que seguem as diretrizes]

### ❌ PROBLEMAS DE COR
- [Cores fora do padrão]


### 🎯 RECOMENDAÇÕES DE COR
- [Sugestões para paleta de cores]
"""
        },
        'tipografia_texto': {
            'nome': '🔤 Especialista em Tipografia e Texto',
            'prompt': f"""
{contexto_agente}
{contexto_global}

## FUNÇÃO: ESPECIALISTA EM TIPOGRAFIA E TEXTO

**Sua tarefa:** Analisar EXCLUSIVAMENTE tipografia e elementos textuais.

### CRITÉRIOS DE ANÁLISE:
1. **Legibilidade** - Facilidade de leitura do texto
2. **Hierarquia Tipográfica** - Tamanhos e pesos de fonte
3. **Alinhamento** - Organização do texto na imagem
4. **Consistência** - Uso uniforme de fontes
5. **Mensagem Textual** - Conteúdo das palavras

### FORMATO DE RESPOSTA OBRIGATÓRIO:

## 🔤 RELATÓRIO DE TIPOGRAFIA

### ✅ ACERTOS TIPOGRÁFICOS
- [Elementos textuais bem executados]

### ⚠️ PROBLEMAS DE TEXTO
- [Problemas com tipografia e texto - Sejam erros visuais, ortográficos ou lexicais]


### ✏️ SUGESTÕES TIPOGRÁFICAS
- [Melhorias para texto e fontes]
"""
        },
        'elementos_marca': {
            'nome': '🏷️ Especialista em Elementos de Marca',
            'prompt': f"""
{contexto_agente}
{contexto_global}

## FUNÇÃO: ESPECIALISTA EM ELEMENTOS DE MARCA

**Sua tarefa:** Analisar EXCLUSIVAMENTE elementos de identidade visual da marca.

### CRITÉRIOS DE ANÁLISE:
1. **Logo e Identidade** - Uso correto da marca
2. **Elementos Gráficos** - Ícones, padrões, ilustrações
3. **Fotografia** - Estilo e tratamento de imagens
4. **Consistência Visual** - Coerência com guidelines
5. **Diferenciação** - Elementos únicos da marca

### FORMATO DE RESPOSTA OBRIGATÓRIO:

## 🏷️ RELATÓRIO DE ELEMENTOS DE MARCA

### ✅ ELEMENTOS CORRETOS
- [Elementos alinhados com a marca]

### ❌ ELEMENTOS INCORRETOS
- [Elementos fora do padrão]


### 🎨 RECOMENDAÇÕES DE MARCA
- [Sugestões para identidade visual]
"""
        },
        'impacto_comunicacao': {
            'nome': '🎯 Especialista em Impacto e Comunicação',
            'prompt': f"""
{contexto_agente}
{contexto_global}

## FUNÇÃO: ESPECIALISTA EM IMPACTO E COMUNICAÇÃO

**Sua tarefa:** Analisar EXCLUSIVAMENTE impacto visual e comunicação.

### CRITÉRIOS DE ANÁLISE:
1. **Mensagem Central** - Clareza da comunicação
2. **Apelo Emocional** - Conexão com o público
3. **Chamada para Ação** - Efetividade persuasiva
4. **Originalidade** - Diferenciação criativa
5. **Memorabilidade** - Capacidade de ser lembrado

### FORMATO DE RESPOSTA OBRIGATÓRIO:

## 🎯 RELATÓRIO DE IMPACTO

### ✅ PONTOS DE IMPACTO
- [Elementos comunicativos eficazes]

### 📉 OPORTUNIDADES DE MELHORIA
- [Áreas para aumentar impacto]


### 🚀 ESTRATÉGIAS DE COMUNICAÇÃO
- [Técnicas para melhor comunicação]
"""
        }
    }
    
    return analisadores

def criar_analisadores_video(contexto_agente, contexto_global, contexto_video_especifico):
        """Cria analisadores especializados para vídeos - VERSÃO COMPLETA COM 6 ESPECIALISTAS"""
        
        analisadores = {
            'narrativa_estrutura': {
                'nome': '📖 Especialista em Narrativa e Estrutura',
                'prompt': f"""
    {contexto_agente}
    {contexto_global}
    {contexto_video_especifico}
    
    ## FUNÇÃO: ESPECIALISTA EM NARRATIVA E ESTRUTURA
    
    **Sua tarefa:** Analisar EXCLUSIVAMENTE a estrutura narrativa do vídeo.
    
    ### CRITÉRIOS DE ANÁLISE:
    1. **Arco Narrativo** - Desenvolvimento da história
    2. **Ritmo** - Velocidade e fluidez da narrativa
    3. **Estrutura** - Organização do conteúdo
    4. **Transições** - Conexão entre cenas/ideias
    5. **Clímax e Resolução** - Ponto alto e conclusão
    
    ### FORMATO DE RESPOSTA OBRIGATÓRIO:
    
    ## 📖 RELATÓRIO DE NARRATIVA
    
    ### ✅ PONTOS FORTES DA NARRATIVA
    - [Elementos narrativos bem executados]
    
    ### ⚠️ PROBLEMAS DE ESTRUTURA
    - [Issues na organização do conteúdo]
    
    ### 📊 SCORE NARRATIVA: [X/10]
    
    ### 💡 SUGESTÕES NARRATIVAS
    - [Melhorias para estrutura e ritmo]
    """
            },
            'qualidade_audio': {
                'nome': '🔊 Especialista em Qualidade de Áudio',
                'prompt': f"""
    {contexto_agente}
    {contexto_global}
    {contexto_video_especifico}
    
    ## FUNÇÃO: ESPECIALISTA EM QUALIDADE DE ÁUDIO
    
    **Sua tarefa:** Analisar EXCLUSIVAMENTE aspectos de áudio do vídeo.
    
    ### CRITÉRIOS DE ANÁLISE:
    1. **Clareza Vocal** - Inteligibilidade da fala
    2. **Qualidade Técnica** - Ruído, distorção, equilíbrio
    3. **Trilha Sonora** - Música e efeitos sonoros
    4. **Sincronização** - Relação áudio-vídeo
    5. **Mixagem** - Balanceamento de elementos sonoros
    
    ### FORMATO DE RESPOSTA OBRIGATÓRIO:
    
    ## 🔊 RELATÓRIO DE ÁUDIO
    
    ### ✅ ACERTOS DE ÁUDIO
    - [Elementos sonoros bem executados]
    
    ### ❌ PROBLEMAS DE ÁUDIO
    - [Issues técnicos e de qualidade]
    
    ### 📊 SCORE ÁUDIO: [X/10]
    
    ### 🎧 RECOMENDAÇÕES DE ÁUDIO
    - [Sugestões para melhor qualidade sonora]
    """
            },
            'visual_cinematografia': {
                'nome': '🎥 Especialista em Visual e Cinematografia',
                'prompt': f"""
    {contexto_agente}
    {contexto_global}
    {contexto_video_especifico}
    
    ## FUNÇÃO: ESPECIALISTA EM VISUAL E CINEMATOGRAFIA
    
    **Sua tarefa:** Analisar EXCLUSIVAMENTE aspectos visuais do vídeo.
    
    ### CRITÉRIOS DE ANÁLISE:
    1. **Enquadramento** - Composição de cenas
    2. **Iluminação** - Uso da luz e sombras
    3. **Movimento de Câmera** - Dinâmica visual
    
    ### FORMATO DE RESPOSTA OBRIGATÓRIO:
    
    ## 🎥 RELATÓRIO VISUAL
    
    ### ✅ PONTOS FORTES VISUAIS
    - [Elementos visuais bem executados]
    
    ### ⚠️ PROBLEMAS VISUAIS
    - [Issues de qualidade visual]
    
    ### 📊 SCORE VISUAL: [X/10]
    
    ### 🌟 SUGESTÕES VISUAIS
    - [Melhorias para cinematografia]
    """
            },
            'branding_consistencia': {
                'nome': '🏢 Especialista em Branding e Consistência',
                'prompt': f"""
    {contexto_agente}
    {contexto_global}
    {contexto_video_especifico}
    
    ## FUNÇÃO: ESPECIALISTA EM BRANDING E CONSISTÊNCIA
    
    **Sua tarefa:** Analisar EXCLUSIVAMENTE alinhamento com branding.
    
    ### CRITÉRIOS DE ANÁLISE:
    1. **Identidade Visual** - Cores, logos, elementos da marca
    2. **Tom de Voz** - Personalidade da comunicação
    3. **Mensagem Central** - Alinhamento com valores
    4. **Público-Alvo** - Adequação ao destinatário
    
    ### FORMATO DE RESPOSTA OBRIGATÓRIO:
    
    ## 🏢 RELATÓRIO DE BRANDING
    
    ### ✅ ALINHAMENTOS DE MARCA
    - [Elementos que seguem as diretrizes]
    
    ### ❌ DESVIOS DE MARCA
    - [Elementos fora do padrão]
    
    
    ### 🎯 RECOMENDAÇÕES DE MARCA
    - [Sugestões para melhor alinhamento]
    """
            },
            'engajamento_eficacia': {
                'nome': '📈 Especialista em Engajamento e Eficácia',
                'prompt': f"""
    {contexto_agente}
    {contexto_global}
    {contexto_video_especifico}
    
    ## FUNÇÃO: ESPECIALISTA EM ENGAJAMENTO E EFICÁCIA
    
    **Sua tarefa:** Analisar EXCLUSIVAMENTE potencial de engajamento e eficácia comunicativa.
    
    ### CRITÉRIOS DE ANÁLISE:
    1. **Hook Inicial** - Capacidade de prender atenção
    2. **Retenção** - Manutenção do interesse
    3. **Chamada para Ação** - Clareza e persuasão
    4. **Emoção** - Conexão emocional com o público
    5. **Compartilhamento** - Potencial viral
    
    ### FORMATO DE RESPOSTA OBRIGATÓRIO:
    
    ## 📈 RELATÓRIO DE ENGAJAMENTO
    
    ### ✅ PONTOS FORTES DE ENGAJAMENTO
    - [Elementos que engajam o público]
    
    ### 📉 OPORTUNIDADES DE MELHORIA
    - [Áreas para aumentar engajamento]
    
    
    ### 🚀 ESTRATÉGIAS DE ENGAJAMENTO
    - [Técnicas para melhor conexão]
    """
            },
            'sincronizacao_audio_legendas': {
                'nome': '🎯 Especialista em Sincronização Áudio-Legendas',
                'prompt': f"""
    {contexto_agente}
    {contexto_global}
    {contexto_video_especifico}
    
    ## FUNÇÃO: ESPECIALISTA EM SINCRONIZAÇÃO ÁUDIO-LEGENDAS
    
    **Sua tarefa:** Analisar EXCLUSIVAMENTE sincronização entre áudio e legendas.
    
    ### CRITÉRIOS DE ANÁLISE:
    1. **Timing** - Sincronização precisa
    2. **Legibilidade** - Clareza das legendas

    
    ### FORMATO DE RESPOSTA OBRIGATÓRIO:
    
    ## 🎯 RELATÓRIO DE SINCRONIZAÇÃO
    
    ### Time stamps específicos das ocorrências de erros entre o que foi falado e o que está escrito nas legendas
    ### Verificação se a legenda em si está escrita corretamente
    

    """
            }
        }
        
        return analisadores

def executar_analise_imagem_especializada(uploaded_image, nome_imagem, analisadores):
    """Executa análise especializada para imagens com múltiplos especialistas"""
    
    resultados = {}
    
    for area, config in analisadores.items():
        with st.spinner(f"Executando {config['nome']}..."):
            try:
                prompt_completo = f"""
{config['prompt']}

###BEGIN IMAGEM PARA ANÁLISE###
**Arquivo:** {nome_imagem}
**Análise solicitada para:** {config['nome']}
###END IMAGEM PARA ANÁLISE###

Por favor, forneça sua análise especializada no formato solicitado.
"""
                
                # Processar imagem com o especialista específico
                response = modelo_vision.generate_content([
                    prompt_completo,
                    {"mime_type": "image/jpeg", "data": uploaded_image.getvalue()}
                ])
                
                resultados[area] = {
                    'nome': config['nome'],
                    'analise': response.text,
                    'score': extrair_score(response.text)
                }
                
            except Exception as e:
                resultados[area] = {
                    'nome': config['nome'],
                    'analise': f"❌ Erro na análise: {str(e)}",
                    'score': 0
                }
    
    return resultados

def executar_analise_video_especializada(uploaded_video, nome_video, analisadores):
    """Executa análise especializada para vídeos com múltiplos especialistas"""
    
    resultados = {}
    
    for area, config in analisadores.items():
        with st.spinner(f"Executando {config['nome']}..."):
            try:
                prompt_completo = f"""
{config['prompt']}

###BEGIN VÍDEO PARA ANÁLISE###
**Arquivo:** {nome_video}
**Análise solicitada para:** {config['nome']}
###END VÍDEO PARA ANÁLISE###

Por favor, forneça sua análise especializada no formato solicitado.
"""
                
                # Processar vídeo com o especialista específico
                video_bytes = uploaded_video.getvalue()
                
                if len(video_bytes) < 200 * 1024 * 1024:
                    response = modelo_vision.generate_content([
                        prompt_completo,
                        {"mime_type": uploaded_video.type, "data": video_bytes}
                    ])
                else:
                    response = modelo_vision.generate_content([
                        prompt_completo,
                        {"mime_type": uploaded_video.type, "data": video_bytes}
                    ])
                
                resultados[area] = {
                    'nome': config['nome'],
                    'analise': response.text,
                    'score': extrair_score(response.text)
                }
                
            except Exception as e:
                resultados[area] = {
                    'nome': config['nome'],
                    'analise': f"❌ Erro na análise: {str(e)}",
                    'score': 0
                }
    
    return resultados

def gerar_relatorio_imagem_consolidado(resultados_especialistas, nome_imagem, dimensoes):
    """Gera relatório consolidado para imagens"""

    
    relatorio = f"""
# 🖼️ RELATÓRIO CONSOLIDADO DE IMAGEM

**Arquivo:** {nome_imagem}
**Dimensões:** {dimensoes}

**Data da Análise:** {datetime.datetime.now().strftime('%d/%m/%Y %H:%M')}

## 🎖️ SCORES POR ÁREA ESPECIALIZADA
"""
    
    # Adicionar scores individuais

    
    relatorio += "\n## 📋 ANÁLISES DETALHADAS POR ESPECIALISTA\n"
    
    # Adicionar análises detalhadas
    for area, resultado in resultados_especialistas.items():
        relatorio += f"\n### {resultado['nome']}\n"
        relatorio += f"{resultado['analise']}\n"
        relatorio += "---\n"
    
    # Resumo executivo
    relatorio += f"""
## 🚀 RESUMO EXECUTIVO - IMAGEM



### 🎯 PRÓXIMOS PASSOS RECOMENDADOS:
"""
    

    
    return relatorio

def gerar_relatorio_video_consolidado(resultados_especialistas, nome_video, tipo_video):
    """Gera relatório consolidado para vídeos"""
    
   
    
    relatorio = f"""
# 🎬 RELATÓRIO CONSOLIDADO DE VÍDEO

**Arquivo:** {nome_video}
**Formato:** {tipo_video}
**Data da Análise:** {datetime.datetime.now().strftime('%d/%m/%Y %H:%M')}

## 🎖️ SCORES POR ÁREA ESPECIALIZADA
"""
    
    
    
    relatorio += "\n## 📋 ANÁLISES DETALHADAS POR ESPECIALISTA\n"
    
    # Adicionar análises detalhadas
    for area, resultado in resultados_especialistas.items():
        relatorio += f"\n### {resultado['nome']}\n"
        relatorio += f"{resultado['analise']}\n"
        relatorio += "---\n"
    
    # Resumo executivo
    relatorio += f"""
## 🚀 RESUMO EXECUTIVO - VÍDEO


### 🎯 PRÓXIMOS PASSOS RECOMENDADOS:
"""
    
    # Recomendações baseadas nos scores
    areas_baixas = [area for area, resultado in resultados_especialistas.items() if resultado['score'] < 6]
    if areas_baixas:
        nomes_areas = [resultados_especialistas[area]['nome'] for area in areas_baixas]
        relatorio += f"- **Prioridade Máxima:** Focar em {', '.join(nomes_areas)}\n"
    
    areas_medianas = [area for area, resultado in resultados_especialistas.items() if 6 <= resultado['score'] < 8]
    if areas_medianas:
        nomes_areas = [resultados_especialistas[area]['nome'] for area in areas_medianas]
        relatorio += f"- **Otimização Necessária:** Melhorar {', '.join(nomes_areas)}\n"
    
    areas_altas = [area for area, resultado in resultados_especialistas.items() if resultado['score'] >= 8]
    if areas_altas:
        nomes_areas = [resultados_especialistas[area]['nome'] for area in areas_altas]
        relatorio += f"- **Manutenção:** Manter a excelência em {', '.join(nomes_areas)}\n"
    
    return relatorio

# --- FUNÇÕES DE ANÁLISE DE TEXTO (MANTIDAS) ---

def criar_analisadores_texto(contexto_agente, contexto_global):
    """Cria prompts especializados para cada área de análise de texto"""
    
    analisadores = {
        'ortografia': {
            'nome': '🔤 Especialista em Ortografia e Gramática',
            'prompt': f"""
{contexto_global}

## FUNÇÃO: ESPECIALISTA EM ORTOGRAFIA E GRAMÁTICA PORTUGUÊS BR

**Sua tarefa:** Analisar EXCLUSIVAMENTE aspectos ortográficos e gramaticais.

### CRITÉRIOS DE ANÁLISE:
1. **Ortografia** - Erros de escrita
2. **Gramática** - Concordância, regência, colocação
3. **Pontuação** - Uso de vírgulas, pontos, etc.
4. **Acentuação** - Erros de acentuação
5. **Padrão Culto** - Conformidade com norma culta

### FORMATO DE RESPOSTA OBRIGATÓRIO:

## 🔤 RELATÓRIO ORTOGRÁFICO

### ✅ ACERTOS
- [Itens corretos]

### ❌ ERROS IDENTIFICADOS
- [Lista específica de erros com correções]

### 📊 SCORE ORTOGRÁFICO: [X/10]

### 💡 SUGESTÕES DE MELHORIA
- [Recomendações específicas]
"""
        },
        'lexico': {
            'nome': '📚 Especialista em Léxico e Vocabulário',
            'prompt': f"""
{contexto_agente}
{contexto_global}

## FUNÇÃO: ESPECIALISTA EM LÉXICO E VOCABULÁRIO

**Sua tarefa:** Analisar EXCLUSIVAMENTE aspectos lexicais e de vocabulário.

### CRITÉRIOS DE ANÁLISE:
1. **Variedade Lexical** - Riqueza de vocabulário
2. **Precisão Semântica** - Uso adequado das palavras
3. **Repetição** - Palavras ou expressões repetidas em excesso
4. **Jargões** - Uso inadequado de termos técnicos
5. **Clareza** - Facilidade de compreensão

### FORMATO DE RESPOSTA OBRIGATÓRIO:

## 📚 RELATÓRIO LEXICAL

### ✅ VOCABULÁRIO ADEQUADO
- [Pontos fortes do vocabulário]

### ⚠️ ASPECTOS A MELHORAR
- [Problemas lexicais identificados]

### 🔄 SUGESTÕES DE SINÔNIMOS
- [Palavras para substituir]

### 📊 SCORE LEXICAL: [X/10]
"""
        },
        'branding': {
            'nome': '🎨 Especialista em Branding e Identidade',
            'prompt': f"""
{contexto_agente}
{contexto_global}

## FUNÇÃO: ESPECIALISTA EM BRANDING E IDENTIDADE

**Sua tarefa:** Analisar EXCLUSIVAMENTE conformidade com diretrizes de branding.

### CRITÉRIOS DE ANÁLISE:
1. **Tom de Voz** - Alinhamento com personalidade da marca
2. **Mensagem Central** - Consistência da mensagem
3. **Valores da Marca** - Reflexo dos valores organizacionais
4. **Público-Alvo** - Adequação ao público pretendido
5. **Diferenciação** - Elementos únicos da marca

### FORMATO DE RESPOSTA OBRIGATÓRIO:

## 🎨 RELATÓRIO DE BRANDING

### ✅ ALINHAMENTOS
- [Elementos que seguem as diretrizes]

### ❌ DESVIOS IDENTIFICADOS
- [Elementos fora do padrão da marca]

### 📊 SCORE BRANDING: [X/10]

### 💡 RECOMENDAÇÕES ESTRATÉGICAS
- [Sugestões para melhor alinhamento]
"""
        },
        'estrutura': {
            'nome': '📋 Especialista em Estrutura e Formatação',
            'prompt': f"""
{contexto_agente}
{contexto_global}

## FUNÇÃO: ESPECIALISTA EM ESTRUTURA E FORMATAÇÃO

**Sua tarefa:** Analisar EXCLUSIVAMENTE estrutura e organização do conteúdo.

### CRITÉRIOS DE ANÁLISE:
1. **Organização** - Estrutura lógica e sequência
2. **Hierarquia** - Uso adequado de títulos e subtítulos
3. **Coesão** - Ligação entre ideias e parágrafos
4. **Formatação** - Consistência visual
5. **Objetividade** - Clareza na apresentação das ideias

### FORMATO DE RESPOSTA OBRIGATÓRIO:

## 📋 RELATÓRIO ESTRUTURAL

### ✅ ESTRUTURA ADEQUADA
- [Elementos bem organizados]

### ⚠️ PROBLEMAS ESTRUTURAIS
- [Issues de organização identificados]

### 📊 SCORE ESTRUTURAL: [X/10]

### 🏗️ SUGESTÕES DE REORGANIZAÇÃO
- [Melhorias na estrutura]
"""
        }
        
    }
    
    return analisadores

def executar_analise_texto_especializada(texto, nome_arquivo, analisadores):
    """Executa análise com múltiplos especialistas para texto"""
    
    resultados = {}
    
    for area, config in analisadores.items():
        with st.spinner(f"Executando {config['nome']}..."):
            try:
                prompt_completo = f"""
{config['prompt']}

###BEGIN TEXTO PARA ANÁLISE###
**Arquivo:** {nome_arquivo}
**Conteúdo:**
{texto[:8000]}
###END TEXTO PARA ANÁLISE###

Por favor, forneça sua análise no formato solicitado.
"""
                
                resposta = modelo_texto.generate_content(prompt_completo)
                resultados[area] = {
                    'nome': config['nome'],
                    'analise': resposta.text,
                    'score': extrair_score(resposta.text)
                }
                
            except Exception as e:
                resultados[area] = {
                    'nome': config['nome'],
                    'analise': f"❌ Erro na análise: {str(e)}",
                    'score': 0
                }
    
    return resultados

def gerar_relatorio_texto_consolidado(resultados_especialistas, nome_arquivo):
    """Gera relatório consolidado a partir das análises especializadas de texto"""

  
    
    relatorio = f"""
# 📊 RELATÓRIO CONSOLIDADO DE VALIDAÇÃO

**Documento:** {nome_arquivo}

**Data da Análise:** {datetime.datetime.now().strftime('%d/%m/%Y %H:%M')}

## 🎖️ SCORES POR ÁREA
"""
    
  
    
    relatorio += "\n## 📋 ANÁLISES DETALHADAS POR ESPECIALISTA\n"
    
    # Adicionar análises detalhadas
    for area, resultado in resultados_especialistas.items():
        relatorio += f"\n### {resultado['nome']}\n"
        relatorio += f"{resultado['analise']}\n"
        relatorio += "---\n"
    
    # Resumo executivo
    relatorio += f"""
## 🚀 RESUMO EXECUTIVO



### 🎯 PRÓXIMOS PASSOS RECOMENDADOS:
"""
    
   
    
    relatorio += "- **Manutenção:** Manter as áreas com scores altos\n"
    
    return relatorio

def extrair_score(texto_analise):
    """Extrai score numérico do texto de análise"""
    import re
    padrao = r'SCORE.*?\[(\d+)(?:/10)?\]'
    correspondencias = re.findall(padrao, texto_analise, re.IGNORECASE)
    if correspondencias:
        return int(correspondencias[0])
    return 5  # Score padrão se não encontrar

# --- FUNÇÕES ORIGINAIS MANTIDAS ---

def criar_prompt_validacao_preciso(texto, nome_arquivo, contexto_agente):
    """Cria um prompt de validação muito mais preciso para evitar falsos positivos"""
    
    prompt = f"""
{contexto_agente}

###BEGIN TEXTO PARA VALIDAÇÃO###
**Arquivo:** {nome_arquivo}
**Conteúdo:**
{texto[:12000]}
###END TEXTO PARA VALIDAÇÃO###

## FORMATO DE RESPOSTA OBRIGATÓRIO:

### ✅ CONFORMIDADE COM DIRETRIZES
- [Itens que estão alinhados com as diretrizes de branding]

**INCONSISTÊNCIAS COM BRANDING:**
- [Só liste desvios REAIS das diretrizes de branding]

### 💡 TEXTO REVISADO
- [Sugestões para aprimorar]

### 📊 STATUS FINAL
**Documento:** [Aprovado/Necessita ajustes/Reprovado]
**Principais ações necessárias:** [Lista resumida]
"""
    return prompt

def analisar_documento_por_slides(doc, contexto_agente):
    """Analisa documento slide por slide com alta precisão"""
    
    resultados = []
    
    for i, slide in enumerate(doc['slides']):
        with st.spinner(f"Analisando slide {i+1}..."):
            try:
                prompt_slide = f"""
{contexto_agente}

## ANÁLISE POR SLIDE - PRECISÃO ABSOLUTA

###BEGIN TEXTO PARA VALIDAÇÃO###
**SLIDE {i+1}:**
{slide['conteudo'][:2000]}
###END TEXTO PARA VALIDAÇÃO###

**ANÁLISE DO SLIDE {i+1}:**

### ✅ Pontos Fortes:
[O que está bom neste slide]

### ⚠️ Problemas REAIS:
- [Lista CURTA de problemas]

### 💡 Sugestões Específicas:
[Melhorias para ESTE slide específico]

Considere que slides que são introdutórios ou apenas de títulos não precisam de tanto rigor de branding

**STATUS:** [✔️ Aprovado / ⚠️ Ajustes Menores / ❌ Problemas Sérios]
"""
                
                resposta = modelo_texto.generate_content(prompt_slide)
                resultados.append({
                    'slide_num': i+1,
                    'analise': resposta.text,
                    'tem_alteracoes': '❌' in resposta.text or '⚠️' in resposta.text
                })
                
            except Exception as e:
                resultados.append({
                    'slide_num': i+1,
                    'analise': f"❌ Erro na análise do slide: {str(e)}",
                    'tem_alteracoes': False
                })
    
    # Construir relatório consolidado
    relatorio = f"# 📊 RELATÓRIO DE VALIDAÇÃO - {doc['nome']}\n\n"
    relatorio += f"**Total de Slides:** {len(doc['slides'])}\n"
    relatorio += f"**Slides com Alterações:** {sum(1 for r in resultados if r['tem_alteracoes'])}\n\n"
    
    # Slides que precisam de atenção
    slides_com_problemas = [r for r in resultados if r['tem_alteracoes']]
    if slides_com_problemas:
        relatorio += "## 🚨 SLIDES QUE PRECISAM DE ATENÇÃO:\n\n"
        for resultado in slides_com_problemas:
            relatorio += f"### 📋 Slide {resultado['slide_num']}\n"
            relatorio += f"{resultado['analise']}\n\n"
    
    # Resumo executivo
    relatorio += "## 📈 RESUMO EXECUTIVO\n\n"
    if slides_com_problemas:
        relatorio += f"**⚠️ {len(slides_com_problemas)} slide(s) necessitam de ajustes**\n"
        relatorio += f"**✅ {len(doc['slides']) - len(slides_com_problemas)} slide(s) estão adequados**\n"
    else:
        relatorio += "**🎉 Todos os slides estão em conformidade com as diretrizes!**\n"
    
    return relatorio

def extract_text_from_pdf_com_slides(arquivo_pdf):
    """Extrai texto de PDF com informação de páginas"""
    try:
        import PyPDF2
        pdf_reader = PyPDF2.PdfReader(arquivo_pdf)
        slides_info = []
        
        for pagina_num, pagina in enumerate(pdf_reader.pages):
            texto = pagina.extract_text()
            slides_info.append({
                'numero': pagina_num + 1,
                'conteudo': texto,
                'tipo': 'página'
            })
        
        texto_completo = "\n\n".join([f"--- PÁGINA {s['numero']} ---\n{s['conteudo']}" for s in slides_info])
        return texto_completo, slides_info
        
    except Exception as e:
        return f"Erro na extração PDF: {str(e)}", []

def extract_text_from_pptx_com_slides(arquivo_pptx):
    """Extrai texto de PPTX com informação de slides"""
    try:
        from pptx import Presentation
        import io
        
        prs = Presentation(io.BytesIO(arquivo_pptx.read()))
        slides_info = []
        
        for slide_num, slide in enumerate(prs.slides):
            texto_slide = f"--- SLIDE {slide_num + 1} ---\n"
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texto_slide += shape.text + "\n"
            
            slides_info.append({
                'numero': slide_num + 1,
                'conteudo': texto_slide,
                'tipo': 'slide'
            })
        
        texto_completo = "\n\n".join([s['conteudo'] for s in slides_info])
        return texto_completo, slides_info
        
    except Exception as e:
        return f"Erro na extração PPTX: {str(e)}", []

def extrair_texto_arquivo(arquivo):
    """Extrai texto de arquivos TXT e DOCX"""
    try:
        if arquivo.type == "text/plain":
            return str(arquivo.read(), "utf-8")
        elif arquivo.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            import docx
            import io
            doc = docx.Document(io.BytesIO(arquivo.read()))
            texto = ""
            for para in doc.paragraphs:
                texto += para.text + "\n"
            return texto
        else:
            return f"Tipo não suportado: {arquivo.type}"
    except Exception as e:
        return f"Erro na extração: {str(e)}"

def extract_text_from_pdf(pdf_path):
    """
    Extract text from a PDF file using multiple methods for better coverage
    """
    text = ""

    # Method 1: Try with pdfplumber (better for some PDFs)
    try:
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text
    except Exception as e:
        print(f"pdfplumber failed for {pdf_path}: {e}")

    # Method 2: Fallback to PyPDF2 if pdfplumber didn't extract much text
    if len(text.strip()) < 100:  # If very little text was extracted
        try:
            with open(pdf_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text 
        except Exception as e:
            print(f"PyPDF2 also failed for {pdf_path}: {e}")

    return text

# --- INICIALIZAÇÃO DE SESSION_STATE ---
if 'analise_especializada_texto' not in st.session_state:
    st.session_state.analise_especializada_texto = True

if 'analise_especializada_imagem' not in st.session_state:
    st.session_state.analise_especializada_imagem = True

if 'analise_especializada_video' not in st.session_state:
    st.session_state.analise_especializada_video = True

if 'analisadores_selecionados_texto' not in st.session_state:
    st.session_state.analisadores_selecionados_texto = ['ortografia', 'lexico', 'branding']

if 'analisadores_selecionados_imagem' not in st.session_state:
    st.session_state.analisadores_selecionados_imagem = ['composicao_visual', 'cores_branding', 'tipografia_texto', 'elementos_marca']

if 'analisadores_selecionados_video' not in st.session_state:
    st.session_state.analisadores_selecionados_video = ['narrativa_estrutura', 'qualidade_audio', 'visual_cinematografia', 'branding_consistencia']

if 'analise_detalhada' not in st.session_state:
    st.session_state.analise_detalhada = True

if 'validacao_triggered' not in st.session_state:
    st.session_state.validacao_triggered = False

if 'todos_textos' not in st.session_state:
    st.session_state.todos_textos = []

if 'resultados_analise_imagem' not in st.session_state:
    st.session_state.resultados_analise_imagem = []

if 'resultados_analise_video' not in st.session_state:
    st.session_state.resultados_analise_video = []

# --- NOVAS FUNÇÕES PARA COMENTÁRIOS EM PDF ---
from pypdf import PdfReader, PdfWriter
from pypdf.annotations import Text
import io

def extrair_comentarios_analise(texto_analise):
    """Extrai os comentários principais do texto de análise da LLM"""
    comentarios = []
    
    # Padrões para extrair comentários
    padroes = [
        r'❌\s*(.*?)(?=\n|$)',
        r'⚠️\s*(.*?)(?=\n|$)',
        r'###\s*❌\s*(.*?)(?=###|\n\n|$)',
        r'###\s*⚠️\s*(.*?)(?=###|\n\n|$)',
        r'PROBLEMAS.*?\n(.*?)(?=###|\n\n|$)',
        r'ALTERAÇÕES.*?\n(.*?)(?=###|\n\n|$)',
        r'DESVIOS.*?\n(.*?)(?=###|\n\n|$)'
    ]
    
    for padrao in padroes:
        matches = re.findall(padrao, texto_analise, re.IGNORECASE | re.DOTALL)
        for match in matches:
            if isinstance(match, tuple):
                match = match[0]
            comentario = match.strip()
            if comentario and len(comentario) > 10:  # Filtra comentários muito curtos
                comentarios.append(comentario)
    
    # Se não encontrou padrões específicos, extrai parágrafos que contenham palavras-chave
    if not comentarios:
        linhas = texto_analise.split('\n')
        for linha in linhas:
            linha = linha.strip()
            if any(palavra in linha.lower() for palavra in ['erro', 'problema', 'ajuste', 'corrigir', 'melhorar', 'sugestão', 'recomendação']):
                if len(linha) > 20 and not linha.startswith('#'):
                    comentarios.append(linha)
    
    return comentarios[:10]  # Limita a 10 comentários

def adicionar_comentarios_pdf(arquivo_pdf_original, comentarios, nome_documento):
    """Adiciona comentários como anotações no PDF"""
    try:
        # Ler o PDF original
        reader = PdfReader(io.BytesIO(arquivo_pdf_original.getvalue()))
        writer = PdfWriter()
        
        # Copiar todas as páginas
        for page in reader.pages:
            writer.add_page(page)
        
        # Adicionar comentários como anotações
        for i, comentario in enumerate(comentarios):
            if i >= 5:  # Limita a 5 comentários para não sobrecarregar
                break
                
            # Calcular posição (distribui os comentários verticalmente)
            y_pos = 750 - (i * 100)
            
            # Criar anotação de texto
            annotation = Text(
                text=f"📝 Comentário {i+1}: {comentario[:200]}...",  # Limita o texto
                rect=(50, y_pos, 400, y_pos + 20),
                open=False
            )
            
            # Adicionar anotação à primeira página
            writer.add_annotation(page_number=0, annotation=annotation)
        
        # Salvar PDF com comentários
        pdf_com_comentarios = io.BytesIO()
        writer.write(pdf_com_comentarios)
        pdf_com_comentarios.seek(0)
        
        return pdf_com_comentarios
        
    except Exception as e:
        st.error(f"❌ Erro ao adicionar comentários ao PDF: {str(e)}")
        return None


def criar_relatorio_comentarios(comentarios, nome_documento, contexto_analise):
    """Cria um relatório de comentários em formato de texto"""
    relatorio = f"""
# 📋 RELATÓRIO DE COMENTÁRIOS - {nome_documento}

**Data da Análise:** {datetime.datetime.now().strftime('%d/%m/%Y %H:%M')}
**Total de Comentários:** {len(comentarios)}

## 🎯 CONTEXTO DA ANÁLISE
{contexto_analise[:500]}...

## 📝 COMENTÁRIOS E SUGESTÕES

"""
    
    for i, comentario in enumerate(comentarios, 1):
        relatorio += f"### 🔍 Comentário {i}\n{comentario}\n\n"
    
    relatorio += """
## 📊 RESUMO EXECUTIVO

**Próximos Passos Recomendados:**
1. Revisar os comentários no PDF anotado
2. Implementar as correções sugeridas
3. Validar conformidade com diretrizes de branding
4. Realizar revisão final do documento

---
*Relatório gerado automaticamente pelo Sistema de Validação Unificada*
"""
    
    return relatorio
# --- FUNÇÕES PARA VALIDAÇÃO DE TEXTO EM IMAGEM ---

def gerar_relatorio_texto_imagem_consolidado(resultados):
    """Gera relatório consolidado no formato específico para texto em imagem"""
    
    relatorio = f"""
# 📝 RELATÓRIO DE VALIDAÇÃO DE TEXTO EM IMAGEM

**Data da Análise:** {datetime.datetime.now().strftime('%d/%m/%Y %H:%M')}
**Total de Imagens Analisadas:** {len(resultados)}

## 📋 ANÁLISE INDIVIDUAL POR ARTE
"""
    
    for resultado in resultados:
        relatorio += f"\n{resultado['analise']}\n"
    
    # Resumo final em formato de tabela
    relatorio += "\n\n## 📌 RESUMO FINAL\n"
    relatorio += "Arte\tErros encontrados?\tObservações\n"
    relatorio += "---\t---\t---\n"
    
    for resultado in resultados:
        status_text = {
            "Correto": "❌ Não",
            "Ajustes sugeridos": "⚠️ Sugestões apenas",
            "Com erros": "✅ Sim",
            "Erro": "❌ Erro na análise"
        }.get(resultado['status'], "❓ Desconhecido")
        
        relatorio += f"Arte {resultado['indice']}\t{status_text}\t{resultado['status']}\n"
    
    relatorio += f"""
    
**🔍 LEGENDA:**
✅ = Correto
⚠️ = Ajustes sugeridos (não são erros, apenas melhorias)
❌ = Sem erros
❌ = Erro na análise (problema técnico)

---
Relatório gerado automaticamente pelo Sistema de Validação de Texto em Imagem
"""
    
    return relatorio

# --- ABA: VALIDAÇÃO UNIFICADA (COMPLETA) ---
with tab_mapping["✅ Validação Unificada"]:
    st.header("✅ Validação Unificada de Conteúdo")
    
    if not st.session_state.get('agente_selecionado'):
        st.info("Selecione um agente primeiro na aba de Chat")
    else:
        agente = st.session_state.agente_selecionado
        st.subheader(f"Validação com: {agente.get('nome', 'Agente')}")
        
        # Container de contexto global
        st.markdown("---")
        st.subheader("🎯 Contexto para Análise")
        
        contexto_global = st.text_area(
            "**✍️ Contexto adicional para todas as análises:**", 
            height=120, 
            key="contexto_global_validacao",
            placeholder="Forneça contexto adicional que será aplicado a TODAS as análises (texto, documentos, imagens e vídeos)..."
        )
        
        # Subabas para diferentes tipos de validação - AGORA COM VALIDAÇÃO DE TEXTO EM IMAGEM
        subtab_imagem, subtab_texto, subtab_video, subtab_texto_imagem = st.tabs(
            ["🖼️ Validação de Imagem", "📄 Validação de Documentos", "🎬 Validação de Vídeo", "📝 Validação de Texto em Imagem"]
        )
        
        # --- SUBTAB: VALIDAÇÃO DE TEXTO EM IMAGEM ---
        with subtab_texto_imagem:
            st.subheader("📝 Validação de Texto em Imagem")
            
            
            # Upload de múltiplas imagens
            st.markdown("### 📤 Upload de Imagens com Texto")
            
            uploaded_images_texto = st.file_uploader(
                "Carregue uma ou mais imagens para análise de texto",
                type=["jpg", "jpeg", "png", "webp", "gif", "bmp"],
                accept_multiple_files=True,
                key="image_text_upload",
                help="Arquivos de imagem contendo texto para validação"
            )
            
            # Botão para limpar análises anteriores
            if st.button("🗑️ Limpar Análises Anteriores", key="limpar_texto_imagem"):
                if 'resultados_texto_imagem' in st.session_state:
                    del st.session_state.resultados_texto_imagem
                st.rerun()
            
            if uploaded_images_texto:
                st.success(f"✅ {len(uploaded_images_texto)} imagem(ns) carregada(s) para análise de texto")
                
                # Exibir miniaturas das imagens
                st.markdown("### 🖼️ Imagens Carregadas")
                cols = st.columns(min(4, len(uploaded_images_texto)))
                
                for idx, img in enumerate(uploaded_images_texto):
                    with cols[idx % 4]:
                        # Abrir imagem para mostrar miniatura
                        image = Image.open(img)
                        st.image(image, use_container_width=True, caption=f"Arte {idx+1}")
                        st.caption(f"📏 {image.width}x{image.height}px")
                
                # Botão para iniciar análise
                if st.button("🔍 Validar Texto em Todas as Imagens", type="primary", key="validar_texto_imagens"):
                    
                    resultados = []
                    progress_bar = st.progress(0)
                    status_text = st.empty()
                    
                    for idx, uploaded_image in enumerate(uploaded_images_texto):
                        status_text.text(f"📊 Analisando texto na imagem {idx+1} de {len(uploaded_images_texto)}...")
                        progress_bar.progress((idx + 1) / len(uploaded_images_texto))
                        
                        with st.spinner(f'Processando "Arte {idx+1}"...'):
                            try:
                                # Criar prompt específico para análise de texto em imagem
                                prompt_texto_imagem = f"""
                                {contexto_global if contexto_global else ''}
                                
                                ## ANÁLISE DE TEXTO EM IMAGEM
                                
                                **INSTRUÇÕES:**
                                1. Transcreva e analise TODO o texto visível na imagem
                                2. Foque em: ortografia, gramática, clareza e adequação
                                3. Use emojis para indicar o status
                                
                                **FORMATO DE RESPOSTA OBRIGATÓRIO:**
                                
                                ## Arte {idx+1} – [Título do texto extraído ou descrição da imagem]
                                
                                **Texto:**
                                "[Texto extraído da imagem]"
                                
                                **Correções:**
                                [✅/⚠️/❌] [Descrição da análise]
                                
                                🔍 [Observação opcional: sugestões de estilo ou melhoria]
                                
                                ---
                                """
                                
                                # Usar modelo de visão para análise
                                response = modelo_vision.generate_content([
                                    prompt_texto_imagem,
                                    {"mime_type": uploaded_image.type, "data": uploaded_image.getvalue()}
                                ])
                                
                                # Processar resposta
                                analise = response.text
                                
                                # Determinar status baseado na resposta
                                if "❌" in analise:
                                    status = "Com erros"
                                elif "⚠️" in analise:
                                    status = "Ajustes sugeridos"
                                else:
                                    status = "Correto"
                                
                                resultados.append({
                                    'indice': idx + 1,
                                    'nome': uploaded_image.name,
                                    'analise': analise,
                                    'status': status,
                                    'imagem': uploaded_image
                                })
                                
                            except Exception as e:
                                st.error(f"❌ Erro ao processar imagem {uploaded_image.name}: {str(e)}")
                                resultados.append({
                                    'indice': idx + 1,
                                    'nome': uploaded_image.name,
                                    'analise': f"❌ Erro na análise: {str(e)}",
                                    'status': "Erro",
                                    'imagem': uploaded_image
                                })
                    
                    progress_bar.empty()
                    status_text.empty()
                    
                    # Armazenar resultados na sessão
                    st.session_state.resultados_texto_imagem = resultados
                    
                    # Gerar relatório consolidado
                    relatorio_consolidado = gerar_relatorio_texto_imagem_consolidado(resultados)
                    
                    # Exibir resultados
                    st.markdown("---")
                    st.subheader("📋 Relatório de Validação de Texto em Imagens")
                    
                    # Exibir análises individuais
                    for resultado in resultados:
                        with st.expander(f"🖼️ Arte {resultado['indice']} - {resultado['status']}", expanded=True):
                            col_img, col_text = st.columns([1, 2])
                            
                            with col_img:
                                image = Image.open(resultado['imagem'])
                                st.image(image, use_container_width=True, caption=f"Arte {resultado['indice']}")
                            
                            with col_text:
                                st.markdown(resultado['analise'])
                    
                    # Exibir resumo final
                    st.markdown("---")
                    st.subheader("📌 Resumo Final")
                    
                    # Criar tabela de resumo
                    resumo_data = []
                    for resultado in resultados:
                        emoji = {
                            "Correto": "✅",
                            "Ajustes sugeridos": "⚠️", 
                            "Com erros": "❌",
                            "Erro": "❌"
                        }.get(resultado['status'], "❓")
                        
                        resumo_data.append({
                            "Arte": resultado['indice'],
                            "Status": emoji,
                            "Erros encontrados?": "❌ Não" if resultado['status'] == "Correto" else "✅ Sim" if resultado['status'] == "Com erros" else "⚠️ Sugestões",
                            "Observações": resultado['status']
                        })
                    
                    # Mostrar tabela
                    import pandas as pd
                    df_resumo = pd.DataFrame(resumo_data)
                    st.table(df_resumo)
                    
                    # Botão de download
                    st.download_button(
                        "📥 Baixar Relatório Completo (TXT)",
                        data=relatorio_consolidado,
                        file_name=f"relatorio_texto_imagens_{datetime.datetime.now().strftime('%Y%m%d_%H%M')}.txt",
                        mime="text/plain",
                        key="download_relatorio_texto_imagem"
                    )
            
            # Mostrar análises anteriores se existirem
            elif 'resultados_texto_imagem' in st.session_state and st.session_state.resultados_texto_imagem:
                st.info("📋 Análises anteriores encontradas. Carregue novas imagens para nova análise ou use o botão 'Limpar Análises'.")
                
                resultados = st.session_state.resultados_texto_imagem
                
                for resultado in resultados:
                    with st.expander(f"🖼️ Arte {resultado['indice']} - {resultado['status']} (Análise Anterior)", expanded=False):
                        st.markdown(resultado['analise'])
            
            else:
                # Instruções de uso
                st.info("""
                **📋 Como usar a Validação de Texto em Imagem:**
                
                1. **Carregue imagens** contendo texto para análise
                2. **Clique em "Validar Texto em Todas as Imagens"**
                3. **Revise** o relatório detalhado
                4. **Baixe** os resultados para referência
                
                **🎯 O que é analisado:**
                - ✅ Ortografia e acentuação
                - ✅ Concordância verbal e nominal
                - ✅ Clareza e compreensão do texto
                - ✅ Adequação ao contexto (se fornecido)
                - ✅ Sugestões de melhoria de estilo
                
                **📊 Formato do relatório:**
                - Análise individual por imagem
                - Texto extraído entre aspas
                - Correções específicas com emojis
                - Observações opcionais de estilo
                - Resumo final em tabela
                """)
        
        # --- SUBTAB: VALIDAÇÃO DE DOCUMENTOS E TEXTO ---
        with subtab_texto:
            st.subheader("📄 Validação de Documentos e Texto")
            
            # Configurações de exportação PDF
            with st.expander("📤 Configurações de Exportação PDF", expanded=True):
                col_export1, col_export2 = st.columns(2)
                
                with col_export1:
                    incluir_comentarios_pdf = st.checkbox(
                        "Incluir comentários no PDF",
                        value=True,
                        help="Adiciona os comentários da análise como anotações no PDF original"
                    )
                    
                    gerar_relatorio_completo = st.checkbox(
                        "Gerar relatório completo",
                        value=True,
                        help="Cria um arquivo de texto com todos os comentários e análises"
                    )
                
                with col_export2:
                    limitar_comentarios = st.slider(
                        "Máximo de comentários por PDF:",
                        min_value=1,
                        max_value=10,
                        value=5,
                        help="Limita o número de comentários adicionados ao PDF"
                    )
            
            # Botão para limpar análises de texto
            if st.button("🗑️ Limpar Análises de Texto", key="limpar_analises_texto"):
                st.session_state.validacao_triggered = False
                st.session_state.todos_textos = []
                st.session_state.resultados_pdf = {}
                st.rerun()
            
            # Container principal com duas colunas
            col_entrada, col_saida = st.columns([1, 1])
            
            with col_entrada:
                st.markdown("### 📥 Entrada de Conteúdo")
                
                # Opção 1: Texto direto
                texto_input = st.text_area(
                    "**✍️ Digite o texto para validação:**", 
                    height=150, 
                    key="texto_validacao",
                    placeholder="Cole aqui o texto que deseja validar..."
                )
                
                # Opção 2: Upload de múltiplos arquivos
                st.markdown("### 📎 Ou carregue arquivos")
                
                arquivos_documentos = st.file_uploader(
                    "**Documentos suportados:** PDF, PPTX, TXT, DOCX",
                    type=['pdf', 'pptx', 'txt', 'docx'],
                    accept_multiple_files=True,
                    key="arquivos_documentos_validacao"
                )
                
                # Configurações de análise
                with st.expander("⚙️ Configurações de Análise de Texto"):
                    analise_especializada = st.checkbox(
                        "Análise especializada por áreas (recomendado)",
                        value=st.session_state.analise_especializada_texto,
                        help="Usa múltiplos especialistas para análise mais precisa"
                    )
                    
                    analisadores_selecionados = st.multiselect(
                        "Especialistas de texto a incluir:",
                        options=['ortografia', 'lexico', 'branding', 'estrutura', 'engajamento'],
                        default=st.session_state.analisadores_selecionados_texto,
                        format_func=lambda x: {
                            'ortografia': '🔤 Ortografia e Gramática',
                            'lexico': '📚 Léxico e Vocabulário', 
                            'branding': '🎨 Branding e Identidade',
                            'estrutura': '📋 Estrutura e Formatação',
                            'engajamento': '🎯 Engajamento e Persuasão'
                        }[x]
                    )
                    
                    analise_detalhada = st.checkbox(
                        "Análise detalhada por slide/página",
                        value=st.session_state.analise_detalhada
                    )
                
                # Botão de validação
                if st.button("✅ Validar Conteúdo de Texto", type="primary", key="validate_documents", use_container_width=True):
                    st.session_state.validacao_triggered = True
                    st.session_state.analise_especializada_texto = analise_especializada
                    st.session_state.analise_detalhada = analise_detalhada
                    st.session_state.analisadores_selecionados_texto = analisadores_selecionados
            
            with col_saida:
                st.markdown("### 📊 Resultados de Texto")
                
                if st.session_state.validacao_triggered:
                    # Processar todos os conteúdos
                    todos_textos = []
                    arquivos_processados = []
                    resultados_pdf = {}  # Armazena resultados para exportação PDF
                    
                    # Adicionar texto manual se existir
                    if texto_input and texto_input.strip():
                        todos_textos.append({
                            'nome': 'Texto_Manual',
                            'conteudo': texto_input,
                            'tipo': 'texto_direto',
                            'tamanho': len(texto_input),
                            'slides': []
                        })
                    
                    # Processar arquivos uploadados
                    if arquivos_documentos:
                        for arquivo in arquivos_documentos:
                            with st.spinner(f"Processando {arquivo.name}..."):
                                try:
                                    if arquivo.type == "application/pdf":
                                        texto_extraido, slides_info = extract_text_from_pdf_com_slides(arquivo)
                                        # Guardar o arquivo PDF original para possível anotação
                                        arquivo_original = arquivo
                                    elif arquivo.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                                        texto_extraido, slides_info = extract_text_from_pptx_com_slides(arquivo)
                                        arquivo_original = None
                                    elif arquivo.type in ["text/plain", "application/vnd.openxmlformats-officedocument.wordprocessingml.document"]:
                                        texto_extraido = extrair_texto_arquivo(arquivo)
                                        slides_info = []
                                        arquivo_original = None
                                    else:
                                        st.warning(f"Tipo de arquivo não suportado: {arquivo.name}")
                                        continue
                                    
                                    if texto_extraido and texto_extraido.strip():
                                        doc_info = {
                                            'nome': arquivo.name,
                                            'conteudo': texto_extraido,
                                            'slides': slides_info,
                                            'tipo': arquivo.type,
                                            'tamanho': len(texto_extraido),
                                            'arquivo_original': arquivo_original
                                        }
                                        todos_textos.append(doc_info)
                                        arquivos_processados.append(arquivo.name)
                                    
                                except Exception as e:
                                    st.error(f"❌ Erro ao processar {arquivo.name}: {str(e)}")
                    
                    # Verificar se há conteúdo para validar
                    if not todos_textos:
                        st.warning("⚠️ Nenhum conteúdo válido encontrado para validação.")
                    else:
                        st.success(f"✅ {len(todos_textos)} documento(s) processado(s) com sucesso!")
                        
                        # Exibir estatísticas rápidas
                        col_docs, col_palavras, col_chars = st.columns(3)
                        with col_docs:
                            st.metric("📄 Documentos", len(todos_textos))
                        with col_palavras:
                            total_palavras = sum(len(doc['conteudo'].split()) for doc in todos_textos)
                            st.metric("📝 Palavras", total_palavras)
                        with col_chars:
                            total_chars = sum(doc['tamanho'] for doc in todos_textos)
                            st.metric("🔤 Caracteres", f"{total_chars:,}")
                        
                        # Análise individual por documento
                        st.markdown("---")
                        st.subheader("📋 Análise Individual por Documento")
                        
                        for doc in todos_textos:
                            with st.expander(f"📄 {doc['nome']} - {doc['tamanho']} chars", expanded=True):
                                # Informações básicas do documento
                                col_info1, col_info2 = st.columns(2)
                                with col_info1:
                                    st.write(f"**Tipo:** {doc['tipo']}")
                                    st.write(f"**Tamanho:** {doc['tamanho']} caracteres")
                                with col_info2:
                                    if doc['slides']:
                                        st.write(f"**Slides/Páginas:** {len(doc['slides'])}")
                                    else:
                                        st.write("**Estrutura:** Texto simples")
                                
                                # Contexto aplicado
                                if contexto_global and contexto_global.strip():
                                    st.info(f"**🎯 Contexto Aplicado:** {contexto_global}")
                                
                                # Análise de branding
                                with st.spinner(f"Analisando {doc['nome']}..."):
                                    try:
                                        # Construir contexto do agente
                                        contexto_agente = ""
                                        if "base_conhecimento" in agente:
                                            contexto_agente = f"""
                                            ###BEGIN DIRETRIZES DE BRANDING DO AGENTE:###
                                            {agente['base_conhecimento']}
                                            ###END DIRETRIZES DE BRANDING DO AGENTE###
                                            """
                                        
                                        # Adicionar contexto global se fornecido
                                        contexto_completo = contexto_agente
                                        if contexto_global and contexto_global.strip():
                                            contexto_completo += f"""
                                            ###BEGIN CONTEXTO ADICIONAL DO USUARIO###
                                            {contexto_global}
                                            ###END CONTEXTO ADICIONAL DO USUARIO###
                                            """
                                        
                                        # Escolher método de análise
                                        if st.session_state.analise_especializada_texto:
                                            # ANÁLISE ESPECIALIZADA POR MÚLTIPLOS ESPECIALISTAS
                                            st.info("🎯 **Executando análise especializada por múltiplos especialistas...**")
                                            
                                            # Criar analisadores especialistas
                                            analisadores_config = criar_analisadores_texto(contexto_completo, "")
                                            
                                            # Filtrar apenas os selecionados
                                            analisadores_filtrados = {k: v for k, v in analisadores_config.items() 
                                                                     if k in st.session_state.analisadores_selecionados_texto}
                                            
                                            # Executar análises especializadas
                                            resultados_especialistas = executar_analise_texto_especializada(
                                                doc['conteudo'], 
                                                doc['nome'], 
                                                analisadores_filtrados
                                            )
                                            
                                            # Gerar relatório consolidado
                                            relatorio_consolidado = gerar_relatorio_texto_consolidado(
                                                resultados_especialistas, 
                                                doc['nome']
                                            )
                                            
                                            st.markdown(relatorio_consolidado, unsafe_allow_html=True)
                                            
                                            # EXTRAIR COMENTÁRIOS PARA PDF
                                            if incluir_comentarios_pdf and doc['tipo'] == "application/pdf" and doc.get('arquivo_original'):
                                                comentarios = extrair_comentarios_analise(relatorio_consolidado)
                                                if comentarios:
                                                    with st.spinner("📝 Adicionando comentários ao PDF..."):
                                                        pdf_com_comentarios = adicionar_comentarios_pdf(
                                                            doc['arquivo_original'],
                                                            comentarios[:limitar_comentarios],
                                                            doc['nome']
                                                        )
                                                        
                                                        if pdf_com_comentarios:
                                                            # Armazenar para download posterior
                                                            resultados_pdf[doc['nome']] = {
                                                                'pdf_com_comentarios': pdf_com_comentarios,
                                                                'comentarios': comentarios,
                                                                'relatorio': relatorio_consolidado
                                                            }
                                                            
                                                            # Botão de download imediato
                                                            st.download_button(
                                                                label="📥 Baixar PDF com Comentários",
                                                                data=pdf_com_comentarios.getvalue(),
                                                                file_name=f"comentarios_{doc['nome']}",
                                                                mime="application/pdf",
                                                                key=f"download_pdf_{doc['nome']}"
                                                            )
                                            
                                        elif st.session_state.analise_detalhada and doc['slides']:
                                            # Análise detalhada por slide (método antigo)
                                            resultado_analise = analisar_documento_por_slides(doc, contexto_completo)
                                            st.markdown(resultado_analise)
                                            
                                            # EXTRAIR COMENTÁRIOS PARA PDF
                                            if incluir_comentarios_pdf and doc['tipo'] == "application/pdf" and doc.get('arquivo_original'):
                                                comentarios = extrair_comentarios_analise(resultado_analise)
                                                if comentarios:
                                                    with st.spinner("📝 Adicionando comentários ao PDF..."):
                                                        pdf_com_comentarios = adicionar_comentarios_pdf(
                                                            doc['arquivo_original'],
                                                            comentarios[:limitar_comentarios],
                                                            doc['nome']
                                                        )
                                                        
                                                        if pdf_com_comentarios:
                                                            resultados_pdf[doc['nome']] = {
                                                                'pdf_com_comentarios': pdf_com_comentarios,
                                                                'comentarios': comentarios,
                                                                'relatorio': resultado_analise
                                                            }
                                                            
                                                            st.download_button(
                                                                label="📥 Baixar PDF com Comentários",
                                                                data=pdf_com_comentarios.getvalue(),
                                                                file_name=f"comentarios_{doc['nome']}",
                                                                mime="application/pdf",
                                                                key=f"download_pdf_{doc['nome']}"
                                                            )
                                            
                                        else:
                                            # Análise geral do documento (método antigo)
                                            prompt_analise = criar_prompt_validacao_preciso(doc['conteudo'], doc['nome'], contexto_completo)
                                            resposta = modelo_texto.generate_content(prompt_analise)
                                            st.markdown(resposta.text)
                                            
                                            # EXTRAIR COMENTÁRIOS PARA PDF
                                            if incluir_comentarios_pdf and doc['tipo'] == "application/pdf" and doc.get('arquivo_original'):
                                                comentarios = extrair_comentarios_analise(resposta.text)
                                                if comentarios:
                                                    with st.spinner("📝 Adicionando comentários ao PDF..."):
                                                        pdf_com_comentarios = adicionar_comentarios_pdf(
                                                            doc['arquivo_original'],
                                                            comentarios[:limitar_comentarios],
                                                            doc['nome']
                                                        )
                                                        
                                                        if pdf_com_comentarios:
                                                            resultados_pdf[doc['nome']] = {
                                                                'pdf_com_comentarios': pdf_com_comentarios,
                                                                'comentarios': comentarios,
                                                                'relatorio': resposta.text
                                                            }
                                                            
                                                            st.download_button(
                                                                label="📥 Baixar PDF com Comentários",
                                                                data=pdf_com_comentarios.getvalue(),
                                                                file_name=f"comentarios_{doc['nome']}",
                                                                mime="application/pdf",
                                                                key=f"download_pdf_{doc['nome']}"
                                                            )
                                        
                                    except Exception as e:
                                        st.error(f"❌ Erro na análise de {doc['nome']}: {str(e)}")
                        
                        # Armazenar na sessão
                        st.session_state.todos_textos = todos_textos
                        st.session_state.resultados_pdf = resultados_pdf
                        
                        # DOWNLOADS CONSOLIDADOS
                        if resultados_pdf or gerar_relatorio_completo:
                            st.markdown("---")
                            st.subheader("📦 Downloads Consolidados")
                            
                            # Download de todos os PDFs com comentários
                            if resultados_pdf and incluir_comentarios_pdf:
                                col_dl1, col_dl2 = st.columns(2)
                                
                                with col_dl1:
                                    # Criar ZIP com todos os PDFs comentados
                                    import zipfile
                                    from io import BytesIO
                                    
                                    zip_buffer = BytesIO()
                                    with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
                                        for nome_doc, resultado in resultados_pdf.items():
                                            pdf_data = resultado['pdf_com_comentarios'].getvalue()
                                            zip_file.writestr(f"comentarios_{nome_doc}", pdf_data)
                                    
                                    zip_buffer.seek(0)
                                    
                                    st.download_button(
                                        "📚 Baixar Todos os PDFs com Comentários (ZIP)",
                                        data=zip_buffer.getvalue(),
                                        file_name=f"pdfs_com_comentarios_{datetime.datetime.now().strftime('%Y%m%d_%H%M')}.zip",
                                        mime="application/zip",
                                        key="download_zip_pdfs"
                                    )
                                
                                with col_dl2:
                                    # Relatório completo com todos os comentários
                                    if gerar_relatorio_completo:
                                        relatorio_completo = f"""
# 📋 RELATÓRIO COMPLETO DE VALIDAÇÃO

**Data:** {datetime.datetime.now().strftime('%d/%m/%Y %H:%M')}
**Agente:** {agente.get('nome', 'N/A')}
**Total de Documentos:** {len(todos_textos)}
**Contexto Aplicado:** {contexto_global if contexto_global else 'Nenhum contexto adicional'}

## DOCUMENTOS ANALISADOS:
"""
                                        
                                        for doc in todos_textos:
                                            relatorio_completo += f"\n### 📄 {doc['nome']}\n"
                                            if doc['nome'] in resultados_pdf:
                                                resultado = resultados_pdf[doc['nome']]
                                                relatorio_completo += f"**Comentários extraídos:** {len(resultado['comentarios'])}\n\n"
                                                for i, comentario in enumerate(resultado['comentarios'][:limitar_comentarios], 1):
                                                    relatorio_completo += f"**Comentário {i}:** {comentario}\n\n"
                                            relatorio_completo += "---\n"
                                        
                                        st.download_button(
                                            "📝 Baixar Relatório Completo (TXT)",
                                            data=relatorio_completo,
                                            file_name=f"relatorio_completo_{datetime.datetime.now().strftime('%Y%m%d_%H%M')}.txt",
                                            mime="text/plain",
                                            key="download_relatorio_completo"
                                        )
                            
                            # Download individual de relatórios de comentários
                            if gerar_relatorio_completo:
                                st.markdown("### 📄 Relatórios Individuais de Comentários")
                                
                                for nome_doc, resultado in resultados_pdf.items():
                                    col_rel1, col_rel2 = st.columns([3, 1])
                                    
                                    with col_rel1:
                                        st.write(f"**{nome_doc}** - {len(resultado['comentarios'])} comentários")
                                    
                                    with col_rel2:
                                        relatorio_individual = criar_relatorio_comentarios(
                                            resultado['comentarios'],
                                            nome_doc,
                                            resultado['relatorio'][:500]  # Contexto resumido
                                        )
                                        
                                        st.download_button(
                                            "📋 Baixar Relatório",
                                            data=relatorio_individual,
                                            file_name=f"relatorio_comentarios_{nome_doc.split('.')[0]}.txt",
                                            mime="text/plain",
                                            key=f"download_relatorio_{nome_doc}"
                                        )
                
                else:
                    st.info("Digite texto ou carregue arquivos para validar")
        
        # --- SUBTAB: VALIDAÇÃO DE IMAGEM ---
        with subtab_imagem:
            st.subheader("🖼️ Validação de Imagem")
            
            # Botão para limpar análises de imagem
            if st.button("🗑️ Limpar Análises de Imagem", key="limpar_analises_imagem"):
                st.session_state.resultados_analise_imagem = []
                st.rerun()
            
            uploaded_images = st.file_uploader(
                "Carregue uma ou mais imagens para análise", 
                type=["jpg", "jpeg", "png", "webp"], 
                key="image_upload_validacao",
                accept_multiple_files=True
            )
            
            # Configurações de análise de imagem
            with st.expander("⚙️ Configurações de Análise de Imagem"):
                analise_especializada_imagem = st.checkbox(
                    "Análise especializada por áreas (recomendado)",
                    value=st.session_state.analise_especializada_imagem,
                    help="Usa múltiplos especialistas visuais para análise mais precisa",
                    key="analise_especializada_imagem_check"
                )
                
                analisadores_selecionados_imagem = st.multiselect(
                    "Especialistas de imagem a incluir:",
                    options=['composicao_visual', 'cores_branding', 'tipografia_texto', 'elementos_marca', 'impacto_comunicacao'],
                    default=st.session_state.analisadores_selecionados_imagem,
                    format_func=lambda x: {
                        'composicao_visual': '🎨 Composição Visual',
                        'cores_branding': '🌈 Cores e Branding', 
                        'tipografia_texto': '🔤 Tipografia e Texto',
                        'elementos_marca': '🏷️ Elementos de Marca',
                        'impacto_comunicacao': '🎯 Impacto e Comunicação'
                    }[x],
                    key="analisadores_imagem_select"
                )
            
            if uploaded_images:
                st.success(f"✅ {len(uploaded_images)} imagem(ns) carregada(s)")
                
                # Botão para validar todas as imagens
                if st.button("🔍 Validar Todas as Imagens", type="primary", key="validar_imagens_multiplas"):
                    
                    # Lista para armazenar resultados
                    resultados_analise = []
                    
                    # Loop através de cada imagem
                    for idx, uploaded_image in enumerate(uploaded_images):
                        with st.spinner(f'Analisando imagem {idx+1} de {len(uploaded_images)}: {uploaded_image.name}...'):
                            try:
                                # Criar container para cada imagem
                                with st.container():
                                    st.markdown("---")
                                    col_img, col_info = st.columns([2, 1])
                                    
                                    with col_img:
                                        # Exibir imagem
                                        image = Image.open(uploaded_image)
                                        st.image(image, use_container_width=True, caption=f"Imagem {idx+1}: {uploaded_image.name}")
                                    
                                    with col_info:
                                        # Informações da imagem
                                        st.metric("📐 Dimensões", f"{image.width} x {image.height}")
                                        st.metric("📊 Formato", uploaded_image.type)
                                        st.metric("📁 Tamanho", f"{uploaded_image.size / 1024:.1f} KB")
                                    
                                    # Contexto aplicado
                                    if contexto_global and contexto_global.strip():
                                        st.info(f"**🎯 Contexto Aplicado:** {contexto_global}")
                                    
                                    # Análise individual
                                    with st.expander(f"📋 Análise Detalhada - Imagem {idx+1}", expanded=True):
                                        try:
                                            # Construir contexto com base de conhecimento do agente
                                            contexto_agente = ""
                                            if "base_conhecimento" in agente:
                                                contexto_agente = f"""
                                                ###BEGIN DIRETRIZES DE BRANDING DO AGENTE:###
                                                {agente['base_conhecimento']}
                                                ###END DIRETRIZES DE BRANDING DO AGENTE###
                                                """
                                            
                                            # Adicionar contexto global se fornecido
                                            contexto_completo = contexto_agente
                                            if contexto_global and contexto_global.strip():
                                                contexto_completo += f"""
                                                ###BEGIN CONTEXTO ADICIONAL DO USUARIO###
                                                {contexto_global}
                                                ###END CONTEXTO ADICIONAL DO USUARIO###
                                                """
                                            
                                            # Escolher método de análise
                                            if st.session_state.analise_especializada_imagem:
                                                # ANÁLISE ESPECIALIZADA POR MÚLTIPLOS ESPECIALISTAS VISUAIS
                                                st.info("🎯 **Executando análise especializada por múltiplos especialistas visuais...**")
                                                
                                                # Criar analisadores especialistas
                                                analisadores_config = criar_analisadores_imagem(contexto_completo, "")
                                                
                                                # Filtrar apenas os selecionados
                                                analisadores_filtrados = {k: v for k, v in analisadores_config.items() 
                                                                         if k in st.session_state.analisadores_selecionados_imagem}
                                                
                                                # Executar análises especializadas
                                                resultados_especialistas = executar_analise_imagem_especializada(
                                                    uploaded_image, 
                                                    uploaded_image.name, 
                                                    analisadores_filtrados
                                                )
                                                
                                                # Gerar relatório consolidado
                                                relatorio_consolidado = gerar_relatorio_imagem_consolidado(
                                                    resultados_especialistas, 
                                                    uploaded_image.name,
                                                    f"{image.width}x{image.height}"
                                                )
                                                
                                                st.markdown(relatorio_consolidado, unsafe_allow_html=True)
                                                
                                                # Armazenar resultado
                                                resultados_analise.append({
                                                    'nome': uploaded_image.name,
                                                    'indice': idx,
                                                    'analise': relatorio_consolidado,
                                                    'dimensoes': f"{image.width}x{image.height}",
                                                    'tamanho': uploaded_image.size
                                                })
                                                
                                            else:
                                                # Análise geral da imagem (método antigo)
                                                prompt_analise = f"""
                                                {contexto_completo}
                                                
                                                Analise esta imagem e verifique o alinhamento com as diretrizes de branding.
                                                
                                                Forneça a análise em formato claro:
                                                
                                                ## 🖼️ RELATÓRIO DE ALINHAMENTO - IMAGEM {idx+1}
                                                
                                                **Arquivo:** {uploaded_image.name}
                                                **Dimensões:** {image.width} x {image.height}
                                                
                                                ### 🎯 RESUMO DA IMAGEM
                                                [Avaliação geral de conformidade visual e textual]
                                                
                                                ### ✅ ELEMENTOS ALINHADOS 
                                                [Itens visuais e textuais que seguem as diretrizes]
                                                
                                                ### ⚠️ ELEMENTOS FORA DO PADRÃO
                                                [Itens visuais e textuais que não seguem as diretrizes]
                                                
                                                ### 💡 RECOMENDAÇÕES
                                                [Sugestões para melhorar o alinhamento visual e textual]
                                                
                                                ### 🎨 ASPECTOS TÉCNICOS
                                                [Composição, cores, tipografia, etc.]
                                                """
                                                
                                                # Processar imagem
                                                response = modelo_vision.generate_content([
                                                    prompt_analise,
                                                    {"mime_type": "image/jpeg", "data": uploaded_image.getvalue()}
                                                ])
                                                
                                                st.markdown(response.text)
                                                
                                                # Armazenar resultado
                                                resultados_analise.append({
                                                    'nome': uploaded_image.name,
                                                    'indice': idx,
                                                    'analise': response.text,
                                                    'dimensoes': f"{image.width}x{image.height}",
                                                    'tamanho': uploaded_image.size
                                                })
                                            
                                        except Exception as e:
                                            st.error(f"❌ Erro ao processar imagem {uploaded_image.name}: {str(e)}")
                                
                                # Separador visual entre imagens
                                if idx < len(uploaded_images) - 1:
                                    st.markdown("---")
                                    
                            except Exception as e:
                                st.error(f"❌ Erro ao carregar imagem {uploaded_image.name}: {str(e)}")
                    
                    # Armazenar na sessão
                    st.session_state.resultados_analise_imagem = resultados_analise
                    
                    # Resumo executivo
                    st.markdown("---")
                    st.subheader("📋 Resumo Executivo de Imagens")
                    
                    col_resumo1, col_resumo2, col_resumo3 = st.columns(3)
                    with col_resumo1:
                        st.metric("📊 Total de Imagens", len(uploaded_images))
                    with col_resumo2:
                        st.metric("✅ Análises Concluídas", len(resultados_analise))
                    with col_resumo3:
                        st.metric("🖼️ Processadas", len(uploaded_images))
                    
                    # Contexto aplicado no resumo
                    if contexto_global and contexto_global.strip():
                        st.info(f"**🎯 Contexto Aplicado em Todas as Análises:** {contexto_global}")
                    
                    # Botão para download do relatório consolidado
                    if st.button("📥 Exportar Relatório Completo de Imagens", key="exportar_relatorio_imagens"):
                        relatorio = f"""
                        # RELATÓRIO DE VALIDAÇÃO DE IMAGENS
                        
                        **Agente:** {agente.get('nome', 'N/A')}
                        **Data:** {datetime.datetime.now().strftime('%d/%m/%Y %H:%M')}
                        **Total de Imagens:** {len(uploaded_images)}
                        **Contexto Aplicado:** {contexto_global if contexto_global else 'Nenhum contexto adicional'}
                        **Método de Análise:** {'Especializada por Múltiplos Especialistas' if st.session_state.analise_especializada_imagem else 'Tradicional'}
                        
                        ## RESUMO EXECUTIVO
                        {chr(10).join([f"{idx+1}. {img.name}" for idx, img in enumerate(uploaded_images)])}
                        
                        ## ANÁLISES INDIVIDUAIS
                        {chr(10).join([f'### {res["nome"]} {chr(10)}{res["analise"]}' for res in resultados_analise])}
                        """
                        
                        st.download_button(
                            "💾 Baixar Relatório em TXT",
                            data=relatorio,
                            file_name=f"relatorio_validacao_imagens_{datetime.datetime.now().strftime('%Y%m%d_%H%M')}.txt",
                            mime="text/plain"
                        )
            
            # Mostrar análises existentes da sessão
            elif st.session_state.resultados_analise_imagem:
                st.info("📋 Análises anteriores encontradas. Use o botão 'Limpar Análises' para recomeçar.")
                
                for resultado in st.session_state.resultados_analise_imagem:
                    with st.expander(f"🖼️ {resultado['nome']} - Análise Salva", expanded=False):
                        st.markdown(resultado['analise'])
            
            else:
                st.info("📁 Carregue uma ou mais imagens para iniciar a validação de branding")
        
        # --- SUBTAB: VALIDAÇÃO DE VÍDEO ---
        with subtab_video:
            st.subheader("🎬 Validação de Vídeo")
            
            # Botão para limpar análises de vídeo
            if st.button("🗑️ Limpar Análises de Vídeo", key="limpar_analises_video"):
                st.session_state.resultados_analise_video = []
                st.rerun()
            
            # Container principal
            col_upload, col_config = st.columns([2, 1])
            
            with col_upload:
                uploaded_videos = st.file_uploader(
                    "Carregue um ou mais vídeos para análise",
                    type=["mp4", "mpeg", "mov", "avi", "flv", "mpg", "webm", "wmv", "3gpp"],
                    key="video_upload_validacao",
                    accept_multiple_files=True
                )
            
            with col_config:
                st.markdown("### ⚙️ Configurações de Vídeo")
                contexto_video_especifico = st.text_area(
                    "**🎯 Contexto específico para vídeos:**", 
                    height=120, 
                    key="video_context_especifico",
                    placeholder="Contexto adicional específico para análise de vídeos (opcional)..."
                )
                
                analise_especializada_video = st.checkbox(
                    "Análise especializada por áreas (recomendado)",
                    value=True,  # Sempre ativo por padrão
                    help="Usa múltiplos especialistas em vídeo para análise mais precisa",
                    key="analise_especializada_video_check"
                )
                
                # Definir todos os especialistas disponíveis
                todos_analisadores_video = ['narrativa_estrutura', 'qualidade_audio', 'visual_cinematografia', 'branding_consistencia', 'engajamento_eficacia', 'sincronizacao_audio_legendas']
                
                # SEMPRE selecionar todos os especialistas por padrão
                analisadores_selecionados_video = st.multiselect(
                    "Especialistas de vídeo a incluir:",
                    options=todos_analisadores_video,
                    default=todos_analisadores_video,  # TODOS selecionados por padrão
                    format_func=lambda x: {
                        'narrativa_estrutura': '📖 Narrativa e Estrutura',
                        'qualidade_audio': '🔊 Qualidade de Áudio', 
                        'visual_cinematografia': '🎥 Visual e Cinematografia',
                        'sincronizacao_audio_legendas': '🎯 Sincronização Áudio-Legendas',
                        'branding_consistencia': '🏢 Branding e Consistência',
                        'engajamento_eficacia': '📈 Engajamento e Eficácia'
                    }[x],
                    key="analisadores_video_select"
                )
                
                # Botão para selecionar automaticamente todos os especialistas
                if st.button("✅ Selecionar Todos os Especialistas", key="select_all_video_analysts"):
                    st.session_state.analisadores_selecionados_video = todos_analisadores_video
                    st.rerun()
            
            if uploaded_videos:
                st.success(f"✅ {len(uploaded_videos)} vídeo(s) carregado(s)")
                
                # Contexto aplicado
                if contexto_global and contexto_global.strip():
                    st.info(f"**🎯 Contexto Global Aplicado:** {contexto_global}")
                if contexto_video_especifico and contexto_video_especifico.strip():
                    st.info(f"**🎯 Contexto Específico Aplicado:** {contexto_video_especifico}")
                
                # Exibir informações dos vídeos
                st.markdown("### 📊 Informações dos Vídeos")
                
                for idx, video in enumerate(uploaded_videos):
                    col_vid, col_info, col_actions = st.columns([2, 2, 1])
                    
                    with col_vid:
                        st.write(f"**{idx+1}. {video.name}**")
                        st.caption(f"Tipo: {video.type} | Tamanho: {video.size / (1024*1024):.1f} MB")
                    
                    with col_info:
                        st.write("📏 Duração: A ser detectada")
                        st.write("🎞️ Resolução: A ser detectada")
                    
                    with col_actions:
                        if st.button("🔍 Preview", key=f"preview_{idx}"):
                            st.video(video, format=f"video/{video.type.split('/')[-1]}")
                
                # Botão para validar todos os vídeos
                if st.button("🎬 Validar Todos os Vídeos", type="primary", key="validar_videos_multiplas"):
                    
                    resultados_video = []
                    
                    for idx, uploaded_video in enumerate(uploaded_videos):
                        with st.spinner(f'Analisando vídeo {idx+1} de {len(uploaded_videos)}: {uploaded_video.name}...'):
                            try:
                                # Container para cada vídeo
                                with st.container():
                                    st.markdown("---")
                                    
                                    # Header do vídeo
                                    col_header, col_stats = st.columns([3, 1])
                                    
                                    with col_header:
                                        st.subheader(f"🎬 {uploaded_video.name}")
                                    
                                    with col_stats:
                                        st.metric("📊 Status", "Processando")
                                    
                                    # Contexto aplicado para este vídeo
                                    if contexto_global and contexto_global.strip():
                                        st.info(f"**🎯 Contexto Aplicado:** {contexto_global}")
                                    if contexto_video_especifico and contexto_video_especifico.strip():
                                        st.info(f"**🎯 Contexto Específico:** {contexto_video_especifico}")
                                    
                                    # Preview do vídeo
                                    with st.expander("👀 Preview do Vídeo", expanded=False):
                                        st.video(uploaded_video, format=f"video/{uploaded_video.type.split('/')[-1]}")
                                    
                                    # Análise detalhada
                                    with st.expander(f"📋 Análise Completa - {uploaded_video.name}", expanded=True):
                                        try:
                                            # Construir contexto com base de conhecimento do agente
                                            contexto_agente = ""
                                            if "base_conhecimento" in agente:
                                                contexto_agente = f"""
                                                ###BEGIN DIRETRIZES DE BRANDING DO AGENTE:###
                                                {agente['base_conhecimento']}
                                                ###END DIRETRIZES DE BRANDING DO AGENTE###
                                                """
                                            
                                            # Adicionar contexto global se fornecido
                                            contexto_completo = contexto_agente
                                            if contexto_global and contexto_global.strip():
                                                contexto_completo += f"""
                                                ###BEGIN CONTEXTO GLOBAL DO USUARIO###
                                                {contexto_global}
                                                ###END CONTEXTO GLOBAL DO USUARIO###
                                                """
                                            
                                            # Adicionar contexto específico de vídeo se fornecido
                                            if contexto_video_especifico and contexto_video_especifico.strip():
                                                contexto_completo += f"""
                                                ###BEGIN CONTEXTO ESPECÍFICO PARA VÍDEOS###
                                                {contexto_video_especifico}
                                                ###END CONTEXTO ESPECÍFICO PARA VÍDEOS###
                                                """
                                            
                                            # SEMPRE usar análise especializada com TODOS os especialistas selecionados
                                            st.info("🎯 **Executando análise especializada por TODOS os especialistas de vídeo...**")
                                            
                                            # Atualizar session state com os analisadores selecionados
                                            st.session_state.analisadores_selecionados_video = analisadores_selecionados_video
                                            
                                            # Verificar se há especialistas selecionados
                                            if not analisadores_selecionados_video:
                                                st.warning("⚠️ Nenhum especialista selecionado. Selecionando todos automaticamente.")
                                                analisadores_selecionados_video = todos_analisadores_video
                                                st.session_state.analisadores_selecionados_video = todos_analisadores_video
                                            
                                            # Criar analisadores especialistas
                                            analisadores_config = criar_analisadores_video(contexto_agente, contexto_global, contexto_video_especifico)
                                            
                                            # Usar SEMPRE todos os especialistas selecionados
                                            analisadores_filtrados = {k: v for k, v in analisadores_config.items() 
                                                                     if k in analisadores_selecionados_video}
                                            
                                            # Mostrar quais especialistas estão sendo executados
                                            st.success(f"**Especialistas ativos:** {len(analisadores_filtrados)}")
                                            for analista_key in analisadores_filtrados.keys():
                                                emoji_nome = {
                                                    'narrativa_estrutura': '📖 Narrativa e Estrutura',
                                                    'qualidade_audio': '🔊 Qualidade de Áudio',
                                                    'visual_cinematografia': '🎥 Visual e Cinematografia', 
                                                    'sincronizacao_audio_legendas': '🎯 Sincronização Áudio-Legendas',
                                                    'branding_consistencia': '🏢 Branding e Consistência',
                                                    'engajamento_eficacia': '📈 Engajamento e Eficácia'
                                                }.get(analista_key, analista_key)
                                                st.write(f"  - {emoji_nome}")
                                            
                                            # Executar análises especializadas
                                            resultados_especialistas = executar_analise_video_especializada(
                                                uploaded_video, 
                                                uploaded_video.name, 
                                                analisadores_filtrados
                                            )
                                            
                                            # Gerar relatório consolidado
                                            relatorio_consolidado = gerar_relatorio_video_consolidado(
                                                resultados_especialistas, 
                                                uploaded_video.name,
                                                uploaded_video.type
                                            )
                                            
                                            st.markdown(relatorio_consolidado, unsafe_allow_html=True)
                                            
                                            # Armazenar resultado
                                            resultados_video.append({
                                                'nome': uploaded_video.name,
                                                'indice': idx,
                                                'analise': relatorio_consolidado,
                                                'tipo': uploaded_video.type,
                                                'tamanho': uploaded_video.size,
                                                'especialistas_utilizados': list(analisadores_filtrados.keys())
                                            })
                                            
                                        except Exception as e:
                                            st.error(f"❌ Erro ao processar vídeo {uploaded_video.name}: {str(e)}")
                                            resultados_video.append({
                                                'nome': uploaded_video.name,
                                                'indice': idx,
                                                'analise': f"Erro na análise: {str(e)}",
                                                'tipo': uploaded_video.type,
                                                'tamanho': uploaded_video.size,
                                                'especialistas_utilizados': []
                                            })
                                    
                            except Exception as e:
                                st.error(f"❌ Erro ao processar vídeo {uploaded_video.name}: {str(e)}")
                    
                    # Armazenar resultados na sessão
                    st.session_state.resultados_analise_video = resultados_video
                    
                    # Resumo executivo dos vídeos
                    st.markdown("---")
                    st.subheader("📋 Resumo Executivo - Vídeos")
                    
                    col_vid1, col_vid2, col_vid3 = st.columns(3)
                    with col_vid1:
                        st.metric("🎬 Total de Vídeos", len(uploaded_videos))
                    with col_vid2:
                        st.metric("✅ Análises Concluídas", len(resultados_video))
                    with col_vid3:
                        total_especialistas = sum(len(res.get('especialistas_utilizados', [])) for res in resultados_video)
                        st.metric("🎯 Especialistas Executados", total_especialistas)
                    
                    # Contexto aplicado no resumo
                    if contexto_global and contexto_global.strip():
                        st.info(f"**🎯 Contexto Global Aplicado:** {contexto_global}")
                    if contexto_video_especifico and contexto_video_especifico.strip():
                        st.info(f"**🎯 Contexto Específico Aplicado:** {contexto_video_especifico}")
                    
                    # Mostrar especialistas utilizados
                    st.info(f"**🔧 Especialistas utilizados na análise:** {', '.join([analisadores_config[k]['nome'] for k in analisadores_selecionados_video if k in analisadores_config])}")
                    
                    # Botão para download do relatório
                    if st.button("📥 Exportar Relatório de Vídeos", key="exportar_relatorio_videos"):
                        relatorio_videos = f"""
                        # RELATÓRIO DE VALIDAÇÃO DE VÍDEOS
                        
                        **Agente:** {agente.get('nome', 'N/A')}
                        **Data:** {datetime.datetime.now().strftime('%d/%m/%Y %H:%M')}
                        **Total de Vídeos:** {len(uploaded_videos)}
                        **Contexto Global:** {contexto_global if contexto_global else 'Nenhum'}
                        **Contexto Específico:** {contexto_video_especifico if contexto_video_especifico else 'Nenhum'}
                        **Método de Análise:** Análise Especializada por Múltiplos Especialistas
                        **Especialistas Utilizados:** {', '.join(analisadores_selecionados_video)}
                        
                        ## VÍDEOS ANALISADOS:
                        {chr(10).join([f"{idx+1}. {vid.name} ({vid.type}) - {vid.size/(1024*1024):.1f} MB" for idx, vid in enumerate(uploaded_videos)])}
                        
                        ## ANÁLISES INDIVIDUAIS:
                        {chr(10).join([f'### {res["nome"]} {chr(10)}{res["analise"]}' for res in resultados_video])}
                        """
                        
                        st.download_button(
                            "💾 Baixar Relatório em TXT",
                            data=relatorio_videos,
                            file_name=f"relatorio_validacao_videos_{datetime.datetime.now().strftime('%Y%m%d_%H%M')}.txt",
                            mime="text/plain"
                        )
            
            # Mostrar análises existentes da sessão
            elif st.session_state.get('resultados_analise_video'):
                st.info("📋 Análises anteriores encontradas. Use o botão 'Limpar Análises' para recomeçar.")
                
                for resultado in st.session_state.resultados_analise_video:
                    with st.expander(f"🎬 {resultado['nome']} - Análise Salva", expanded=False):
                        st.markdown(resultado['analise'])
                        if resultado.get('especialistas_utilizados'):
                            st.caption(f"**Especialistas utilizados:** {', '.join(resultado['especialistas_utilizados'])}")
            
            else:
                st.info("🎬 Carregue um ou mais vídeos para iniciar a validação")
                
# --- ABA: GERAÇÃO DE CONTEÚDO ---
with tab_mapping["✨ Geração de Conteúdo"]:
    st.header("✨ Geração de Conteúdo com Múltiplos Insumos")
    
    # Conexão com MongoDB para briefings
    try:
        client2 = MongoClient("mongodb+srv://gustavoromao3345:RqWFPNOJQfInAW1N@cluster0.5iilj.mongodb.net/auto_doc?retryWrites=true&w=majority&ssl=true&ssl_cert_reqs=CERT_NONE&tlsAllowInvalidCertificates=true")
        db_briefings = client2['briefings_Broto_Tecnologia']
        collection_briefings = db_briefings['briefings']
        mongo_connected_conteudo = True
    except Exception as e:
        st.error(f"Erro na conexão com MongoDB: {str(e)}")
        mongo_connected_conteudo = False

    # Função para extrair texto de diferentes tipos de arquivo
    def extrair_texto_arquivo(arquivo):
        """Extrai texto de diferentes formatos de arquivo"""
        try:
            extensao = arquivo.name.split('.')[-1].lower()
            
            if extensao == 'pdf':
                return extrair_texto_pdf(arquivo)
            elif extensao == 'txt':
                return extrair_texto_txt(arquivo)
            elif extensao in ['pptx', 'ppt']:
                return extrair_texto_pptx(arquivo)
            elif extensao in ['docx', 'doc']:
                return extrair_texto_docx(arquivo)
            else:
                return f"Formato {extensao} não suportado para extração de texto."
                
        except Exception as e:
            return f"Erro ao extrair texto do arquivo {arquivo.name}: {str(e)}"

    def extrair_texto_pdf(arquivo):
        """Extrai texto de arquivos PDF"""
        try:
            import PyPDF2
            pdf_reader = PyPDF2.PdfReader(arquivo)
            texto = ""
            for pagina in pdf_reader.pages:
                texto += pagina.extract_text() + "\n"
            return texto
        except Exception as e:
            return f"Erro na leitura do PDF: {str(e)}"

    def extrair_texto_txt(arquivo):
        """Extrai texto de arquivos TXT"""
        try:
            return arquivo.read().decode('utf-8')
        except:
            try:
                return arquivo.read().decode('latin-1')
            except Exception as e:
                return f"Erro na leitura do TXT: {str(e)}"

    def extrair_texto_pptx(arquivo):
        """Extrai texto de arquivos PowerPoint"""
        try:
            from pptx import Presentation
            import io
            prs = Presentation(io.BytesIO(arquivo.read()))
            texto = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texto += shape.text + "\n"
            return texto
        except Exception as e:
            return f"Erro na leitura do PowerPoint: {str(e)}"

    def extrair_texto_docx(arquivo):
        """Extrai texto de arquivos Word"""
        try:
            import docx
            import io
            doc = docx.Document(io.BytesIO(arquivo.read()))
            texto = ""
            for para in doc.paragraphs:
                texto += para.text + "\n"
            return texto
        except Exception as e:
            return f"Erro na leitura do Word: {str(e)}"

    # Layout principal
    col1, col2 = st.columns([2, 1])
    
    with col1:
        st.subheader("📝 Fontes de Conteúdo")
        
        # Opção 1: Upload de múltiplos arquivos
        st.write("📎 Upload de Arquivos (PDF, TXT, PPTX, DOCX):")
        arquivos_upload = st.file_uploader(
            "Selecione um ou mais arquivos:",
            type=['pdf', 'txt', 'pptx', 'ppt', 'docx', 'doc'],
            accept_multiple_files=True,
            help="Arquivos serão convertidos para texto e usados como base para geração de conteúdo"
        )
        
        # Processar arquivos uploadados
        textos_arquivos = ""
        if arquivos_upload:
            st.success(f"✅ {len(arquivos_upload)} arquivo(s) carregado(s)")
            
            with st.expander("📋 Visualizar Conteúdo dos Arquivos", expanded=False):
                for i, arquivo in enumerate(arquivos_upload):
                    st.write(f"**{arquivo.name}** ({arquivo.size} bytes)")
                    with st.spinner(f"Processando {arquivo.name}..."):
                        texto_extraido = extrair_texto_arquivo(arquivo)
                        textos_arquivos += f"\n\n--- CONTEÚDO DE {arquivo.name.upper()} ---\n{texto_extraido}"
                        
                        # Mostrar preview
                        if len(texto_extraido) > 500:
                            st.text_area(f"Preview - {arquivo.name}", 
                                       value=texto_extraido[:500] + "...", 
                                       height=100,
                                       key=f"preview_{i}")
                        else:
                            st.text_area(f"Preview - {arquivo.name}", 
                                       value=texto_extraido, 
                                       height=100,
                                       key=f"preview_{i}")
        
        # Opção 2: Upload de imagem para geração de legenda
        st.write("🖼️ Gerar Legenda para Imagem:")
        imagem_upload = st.file_uploader(
            "Selecione uma imagem:",
            type=['jpg', 'jpeg', 'png', 'gif', 'bmp', 'webp'],
            help="A legenda será gerada com base na imagem e no contexto do agente selecionado"
        )
        
        # Mostrar preview da imagem se carregada
        if imagem_upload:
            col_img1, col_img2 = st.columns([1, 2])
            with col_img1:
                st.image(imagem_upload, caption="Imagem Carregada", use_container_width=True)
                st.write(f"**Arquivo:** {imagem_upload.name}")
                st.write(f"**Tamanho:** {imagem_upload.size / 1024:.1f} KB")
            
            with col_img2:
                # Configurações específicas para legenda de imagem
                st.subheader("Configurações da Legenda")
                
                estilo_legenda = st.selectbox(
                    "Estilo da Legenda:",
                    ["Descritiva", "Criativa", "Técnica", "Comercial", "Emocional", "Storytelling"],
                    help="Escolha o estilo da legenda a ser gerada"
                )
                
                comprimento_legenda = st.select_slider(
                    "Comprimento da Legenda:",
                    options=["Curta", "Média", "Longa"],
                    value="Média"
                )
                
                incluir_hashtags = st.checkbox("Incluir hashtags relevantes", value=True)
                
                # Botão para gerar legenda individual
                if st.button("📝 Gerar Legenda para esta Imagem", use_container_width=True):
                    if not st.session_state.agente_selecionado:
                        st.error("❌ Selecione um agente primeiro para usar seu contexto na geração da legenda")
                    else:
                        with st.spinner("Analisando imagem e gerando legenda..."):
                            try:
                                # Preparar contexto do agente
                                contexto_agente = ""
                                if st.session_state.agente_selecionado:
                                    agente = st.session_state.agente_selecionado
                                    contexto_agente = construir_contexto(agente, st.session_state.segmentos_selecionados)
                                
                                # Usar modelo de visão para analisar a imagem
                                prompt_legenda = f"""
                                {contexto_agente}
                                
                                ## ANÁLISE DE IMAGEM PARA GERAÇÃO DE LEGENDA:
                                
                                **ESTILO SOLICITADO:** {estilo_legenda}
                                **COMPRIMENTO:** {comprimento_legenda}
                                **INCLUIR HASHTAGS:** {incluir_hashtags}
                                
                                ## TAREFA:
                                Analise esta imagem e gere uma legenda que:
                                
                                1. **Descreva** accuratamente o conteúdo visual
                                2. **Contextualize** com base no conhecimento do agente selecionado
                                3. **Engaje** o público-alvo apropriado
                                4. **Siga** o estilo {estilo_legenda.lower()}
                                5. **Tenha** comprimento {comprimento_legenda.lower()}
                                { "6. **Inclua** hashtags relevantes ao final" if incluir_hashtags else "" }
                                
                                Seja criativo mas mantenha a precisão factual.
                                """
                                
                                # Usar modelo de visão para gerar legenda
                                modelo_visao = genai.GenerativeModel('gemini-2.0-flash')
                                resposta_legenda = modelo_visao.generate_content([
                                    prompt_legenda,
                                    {"mime_type": imagem_upload.type, "data": imagem_upload.getvalue()}
                                ])
                                
                                legenda_gerada = resposta_legenda.text
                                
                                # Mostrar resultado
                                st.success("✅ Legenda gerada com sucesso!")
                                st.subheader("Legenda Gerada:")
                                st.write(legenda_gerada)
                                
                                # Botão para copiar legenda
                                st.download_button(
                                    "📋 Copiar Legenda",
                                    data=legenda_gerada,
                                    file_name=f"legenda_{imagem_upload.name.split('.')[0]}.txt",
                                    mime="text/plain"
                                )
                                
                                # Salvar no histórico se MongoDB disponível
                                if mongo_connected_conteudo:
                                    try:
                                        historico_legenda = {
                                            "tipo": "legenda_imagem",
                                            "nome_imagem": imagem_upload.name,
                                            "estilo_legenda": estilo_legenda,
                                            "comprimento_legenda": comprimento_legenda,
                                            "legenda_gerada": legenda_gerada,
                                            "agente_utilizado": st.session_state.agente_selecionado.get('nome') if st.session_state.agente_selecionado else "Nenhum",
                                            "data_criacao": datetime.datetime.now()
                                        }
                                        db_briefings['historico_legendas'].insert_one(historico_legenda)
                                        st.success("✅ Legenda salva no histórico!")
                                    except Exception as e:
                                        st.warning(f"Legenda gerada, mas não salva no histórico: {str(e)}")
                                
                            except Exception as e:
                                st.error(f"❌ Erro ao gerar legenda: {str(e)}")
                                st.info("💡 Dica: Verifique se a imagem não está corrompida e tente novamente.")
        
        # Opção 3: Inserir briefing manualmente
        st.write("✍️ Briefing Manual:")
        briefing_manual = st.text_area("Ou cole o briefing completo aqui:", height=150,
                                      placeholder="""Exemplo:
Título: Campanha de Lançamento
Objetivo: Divulgar novo produto
Público-alvo: Empresários...
Pontos-chave: [lista os principais pontos]""")
        
        # Transcrição de áudio/vídeo
        st.write("🎤 Transcrição de Áudio/Video:")
        arquivos_midia = st.file_uploader(
            "Áudios/Vídeos para transcrição:",
            type=['mp3', 'wav', 'mp4', 'mov', 'avi'],
            accept_multiple_files=True,
            help="Arquivos de mídia serão transcritos automaticamente"
        )
        
        transcricoes_texto = ""
        if arquivos_midia:
            st.info(f"🎬 {len(arquivos_midia)} arquivo(s) de mídia carregado(s)")
            if st.button("🔄 Transcrever Todos os Arquivos de Mídia"):
                with st.spinner("Transcrevendo arquivos de mídia..."):
                    for arquivo in arquivos_midia:
                        tipo = "audio" if arquivo.type.startswith('audio') else "video"
                        transcricao = transcrever_audio_video(arquivo, tipo)
                        transcricoes_texto += f"\n\n--- TRANSCRIÇÃO DE {arquivo.name.upper()} ---\n{transcricao}"
                        st.success(f"✅ {arquivo.name} transcrito!")
    
    with col2:
        st.subheader("⚙️ Configurações de Geração")
        
        # Indicador de agente selecionado
        if st.session_state.agente_selecionado:
            st.info(f"🤖 Agente: {st.session_state.agente_selecionado.get('nome', 'N/A')}")
        else:
            st.warning("⚠️ Nenhum agente selecionado")
        
        # Opção para o usuário escolher entre configurações padrão ou prompt personalizado
        modo_geracao = st.radio(
            "Modo de Geração:",
            ["Configurações Padrão", "Prompt Personalizado"],
            help="Escolha entre usar configurações pré-definidas ou escrever seu próprio prompt"
        )
        
        if modo_geracao == "Configurações Padrão":
            # Configurações básicas (versão simplificada)
            tipo_conteudo = st.selectbox("Tipo de Conteúdo:", 
                                       ["Post Social", "Artigo Blog", "Email Marketing", 
                                        "Landing Page", "Script Vídeo", "Relatório Técnico",
                                        "Press Release", "Newsletter", "Case Study"])
            
            tom_voz = st.selectbox("Tom de Voz:", 
                                  ["Formal", "Informal", "Persuasivo", "Educativo", 
                                   "Inspirador", "Técnico", "Jornalístico"])
            
            palavras_chave = st.text_input("Palavras-chave (opcional):",
                                          placeholder="separadas por vírgula")
            
            numero_palavras = st.slider("Número de Palavras:", 100, 3000, 800)
            
            # Configurações avançadas simplificadas
            with st.expander("🔧 Configurações Avançadas"):
                usar_contexto_agente = st.checkbox("Usar contexto do agente selecionado", 
                                                 value=bool(st.session_state.agente_selecionado))
                
                incluir_cta = st.checkbox("Incluir Call-to-Action", value=True)
                
                formato_saida = st.selectbox("Formato de Saída:", 
                                           ["Texto Simples", "Markdown", "HTML Básico"])
        
        else:  # Prompt Personalizado
            st.info("💡 Escreva seu próprio prompt de geração. Use {contexto} para incluir automaticamente todas as fontes de conteúdo.")
            prompt_personalizado = st.text_area(
                "Seu Prompt Personalizado:",
                height=200,
                placeholder="""Exemplo:
Com base no contexto fornecido, crie um artigo detalhado que:

1. Explique os conceitos principais de forma clara
2. Destaque os benefícios para o público-alvo
3. Inclua exemplos práticos de aplicação
4. Mantenha um tom {tom} e acessível

Contexto: {contexto}

Gere o conteúdo em formato {formato} com aproximadamente {palavras} palavras."""
            )
            
            # Variáveis que o usuário pode usar no prompt personalizado
            col_var1, col_var2, col_var3 = st.columns(3)
            with col_var1:
                tom_personalizado = st.selectbox("Tom:", 
                                               ["formal", "informal", "persuasivo", "educativo"], 
                                               key="tom_personalizado")
            with col_var2:
                formato_personalizado = st.selectbox("Formato:", 
                                                   ["texto simples", "markdown", "HTML básico"], 
                                                   key="formato_personalizado")
            with col_var3:
                palavras_personalizado = st.slider("Palavras:", 100, 3000, 800, key="palavras_personalizado")
            
            usar_contexto_agente = st.checkbox("Usar contexto do agente selecionado", 
                                             value=bool(st.session_state.agente_selecionado),
                                             key="contexto_personalizado")

    # Área de instruções específicas (apenas para modo padrão)
    if modo_geracao == "Configurações Padrão":
        st.subheader("🎯 Instruções Específicas")
        instrucoes_especificas = st.text_area(
            "Diretrizes adicionais para geração:",
            placeholder="""Exemplos:
- Focar nos benefícios para o usuário final
- Incluir estatísticas quando possível
- Manter linguagem acessível
- Evitar jargões técnicos excessivos
- Seguir estrutura: problema → solução → benefícios""",
            height=100
        )

    # Botão para gerar conteúdo
    if st.button("🚀 Gerar Conteúdo com Todos os Insumos", type="primary", use_container_width=True):
        # Verificar se há pelo menos uma fonte de conteúdo
        tem_conteudo = (arquivos_upload or 
                       briefing_manual or 
                       ('briefing_data' in locals() and briefing_data) or
                       arquivos_midia)
        
        if not tem_conteudo:
            st.error("❌ Por favor, forneça pelo menos uma fonte de conteúdo (arquivos, briefing ou mídia)")
        elif modo_geracao == "Prompt Personalizado" and not prompt_personalizado:
            st.error("❌ Por favor, escreva um prompt personalizado para geração")
        else:
            with st.spinner("Processando todos os insumos e gerando conteúdo..."):
                try:
                    # Construir o contexto combinado de todas as fontes
                    contexto_completo = "## FONTES DE CONTEÚDO COMBINADAS:\n\n"
                    
                    # Adicionar conteúdo dos arquivos uploadados
                    if textos_arquivos:
                        contexto_completo += "### CONTEÚDO DOS ARQUIVOS:\n" + textos_arquivos + "\n\n"
                    
                    # Adicionar briefing do banco ou manual
                    if briefing_manual:
                        contexto_completo += "### BRIEFING MANUAL:\n" + briefing_manual + "\n\n"
                    elif 'briefing_data' in locals() and briefing_data:
                        contexto_completo += "### BRIEFING DO BANCO:\n" + briefing_data['conteudo'] + "\n\n"
                    
                    # Adicionar transcrições
                    if transcricoes_texto:
                        contexto_completo += "### TRANSCRIÇÕES DE MÍDIA:\n" + transcricoes_texto + "\n\n"
                    
                    # Adicionar contexto do agente se selecionado
                    contexto_agente = ""
                    if usar_contexto_agente and st.session_state.agente_selecionado:
                        agente = st.session_state.agente_selecionado
                        contexto_agente = construir_contexto(agente, st.session_state.segmentos_selecionados)
                    
                    # Construir prompt final baseado no modo selecionado
                    if modo_geracao == "Configurações Padrão":
                        prompt_final = f"""
                        {contexto_agente}
                        
                        ## INSTRUÇÕES PARA GERAÇÃO DE CONTEÚDO:
                        
                        **TIPO DE CONTEÚDO:** {tipo_conteudo}
                        **TOM DE VOZ:** {tom_voz}
                        **PALAVRAS-CHAVE:** {palavras_chave if palavras_chave else 'Não especificadas'}
                        **NÚMERO DE PALAVRAS:** {numero_palavras} (±10%)
                        **INCLUIR CALL-TO-ACTION:** {incluir_cta}
                        
                        **INSTRUÇÕES ESPECÍFICAS:**
                        {instrucoes_especificas if instrucoes_especificas else 'Nenhuma instrução específica fornecida.'}
                        
                        ## FONTES E REFERÊNCIAS:
                        {contexto_completo}
                        
                        ## TAREFA:
                        Com base em TODAS as fontes fornecidas acima, gere um conteúdo do tipo {tipo_conteudo} que:
                        
                        1. **Síntese Eficiente:** Combine e sintetize informações de todas as fontes
                        2. **Coerência:** Mantenha consistência com as informações originais
                        3. **Valor Agregado:** Vá além da simples cópia, agregando insights
                        4. **Engajamento:** Crie conteúdo que engaje o público-alvo
                        5. **Clareza:** Comunique ideias complexas de forma acessível
                        
                        **FORMATO DE SAÍDA:** {formato_saida}
                        
                        Gere um conteúdo completo e profissional.
                        """
                    else:  # Prompt Personalizado
                        # Substituir variáveis no prompt personalizado
                        prompt_processado = prompt_personalizado.replace("{contexto}", contexto_completo)
                        prompt_processado = prompt_processado.replace("{tom}", tom_personalizado)
                        prompt_processado = prompt_processado.replace("{formato}", formato_personalizado)
                        prompt_processado = prompt_processado.replace("{palavras}", str(palavras_personalizado))
                        
                        prompt_final = f"""
                        {contexto_agente}
                        
                        {prompt_processado}
                        """
                    
                    resposta = modelo_texto.generate_content(prompt_final)
                    
                    # Determinar formato de saída baseado no modo
                    if modo_geracao == "Configurações Padrão":
                        formato_output = formato_saida
                    else:
                        formato_output = formato_personalizado
                    
                    # Processar saída baseada no formato selecionado
                    conteudo_gerado = resposta.text
                    
                    if formato_output == "HTML Básico" or formato_output == "HTML básico":
                        # Converter markdown para HTML básico
                        import re
                        conteudo_gerado = re.sub(r'\*\*(.*?)\*\*', r'<strong>\1</strong>', conteudo_gerado)
                        conteudo_gerado = re.sub(r'\*(.*?)\*', r'<em>\1</em>', conteudo_gerado)
                        conteudo_gerado = re.sub(r'### (.*?)\n', r'<h3>\1</h3>', conteudo_gerado)
                        conteudo_gerado = re.sub(r'## (.*?)\n', r'<h2>\1</h2>', conteudo_gerado)
                        conteudo_gerado = re.sub(r'# (.*?)\n', r'<h1>\1</h1>', conteudo_gerado)
                        conteudo_gerado = conteudo_gerado.replace('\n', '<br>')
                    
                    st.subheader("📄 Conteúdo Gerado")
                    
                    if formato_output == "HTML Básico" or formato_output == "HTML básico":
                        st.components.v1.html(conteudo_gerado, height=400, scrolling=True)
                    else:
                        st.markdown(conteudo_gerado)
                    
                    # Estatísticas
                    palavras_count = len(conteudo_gerado.split())
                    col_stat1, col_stat2, col_stat3 = st.columns(3)
                    with col_stat1:
                        st.metric("Palavras Geradas", palavras_count)
                    with col_stat2:
                        st.metric("Arquivos Processados", len(arquivos_upload) if arquivos_upload else 0)
                    with col_stat3:
                        st.metric("Fontes Utilizadas", 
                                 (1 if arquivos_upload else 0) + 
                                 (1 if briefing_manual or 'briefing_data' in locals() else 0) +
                                 (1 if transcricoes_texto else 0))
                    
                    # Botões de download
                    extensao = ".html" if "HTML" in formato_output else ".md" if "markdown" in formato_output.lower() else ".txt"
                    
                    st.download_button(
                        f"💾 Baixar Conteúdo ({formato_output})",
                        data=conteudo_gerado,
                        file_name=f"conteudo_gerado_{datetime.datetime.now().strftime('%Y%m%d_%H%M')}{extensao}",
                        mime="text/html" if "HTML" in formato_output else "text/plain"
                    )
                    
                    # Salvar no histórico se MongoDB disponível
                    if mongo_connected_conteudo:
                        try:
                            from bson import ObjectId
                            historico_data = {
                                "modo_geracao": modo_geracao,
                                "tipo_conteudo": tipo_conteudo if modo_geracao == "Configurações Padrão" else "Personalizado",
                                "tom_voz": tom_voz if modo_geracao == "Configurações Padrão" else tom_personalizado,
                                "palavras_chave": palavras_chave if modo_geracao == "Configurações Padrão" else "Personalizado",
                                "numero_palavras": numero_palavras if modo_geracao == "Configurações Padrão" else palavras_personalizado,
                                "conteudo_gerado": conteudo_gerado,
                                "fontes_utilizadas": {
                                    "arquivos_upload": [arquivo.name for arquivo in arquivos_upload] if arquivos_upload else [],
                                    "briefing_manual": bool(briefing_manual),
                                    "transcricoes": len(arquivos_midia) if arquivos_midia else 0
                                },
                                "data_criacao": datetime.datetime.now()
                            }
                            db_briefings['historico_geracao'].insert_one(historico_data)
                            st.success("✅ Conteúdo salva no histórico!")
                        except Exception as e:
                            st.warning(f"Conteúdo gerado, mas não salvo no histórico: {str(e)}")
                    
                except Exception as e:
                    st.error(f"❌ Erro ao gerar conteúdo: {str(e)}")
                    st.info("💡 Dica: Verifique se os arquivos não estão corrompidos e tente novamente.")

    # Seção de histórico rápido
    if mongo_connected_conteudo:
        with st.expander("📚 Histórico de Gerações Recentes"):
            try:
                historico = list(db_briefings['historico_geracao'].find().sort("data_criacao", -1).limit(5))
                if historico:
                    for item in historico:
                        st.write(f"**{item['tipo_conteúdo']}** - {item['data_criacao'].strftime('%d/%m/%Y %H:%M')}")
                        st.caption(f"Palavras-chave: {item.get('palavras_chave', 'Nenhuma')} | Tom: {item['tom_voz']}")
                        with st.expander("Ver conteúdo"):
                            st.write(item['conteudo_gerado'][:500] + "..." if len(item['conteudo_gerado']) > 500 else item['conteudo_gerado'])
                else:
                    st.info("Nenhuma geração no histórico")
            except Exception as e:
                st.warning(f"Erro ao carregar histórico: {str(e)}")

        # Histórico de legendas geradas
        with st.expander("🖼️ Histórico de Legendas"):
            try:
                historico_legendas = list(db_briefings['historico_legendas'].find().sort("data_criacao", -1).limit(5))
                if historico_legendas:
                    for item in historico_legendas:
                        st.write(f"**{item['nome_imagem']}** - {item['data_criacao'].strftime('%d/%m/%Y %H:%M')}")
                        st.caption(f"Estilo: {item['estilo_legenda']} | Comprimento: {item['comprimento_legenda']}")
                        st.write(f"*{item['legenda_gerada'][:100]}...*" if len(item['legenda_gerada']) > 100 else item['legenda_gerada'])
                        st.divider()
                else:
                    st.info("Nenhuma legenda no histórico")
            except Exception as e:
                st.warning(f"Erro ao carregar histórico de legendas: {str(e)}")



# --- FUNÇÕES DE REVISÃO ORTOGRÁFICA ---
def revisar_texto_ortografia(texto, agente, segmentos_selecionados, revisao_estilo=True, manter_estrutura=True, explicar_alteracoes=True, modelo_escolhido="Gemini"):
    """
    Realiza revisão ortográfica e gramatical do texto considerando as diretrizes do agente
    """
    
    # Construir o contexto do agente
    contexto_agente = "CONTEXTO DO AGENTE PARA REVISÃO:\n\n"
    
    if "system_prompt" in segmentos_selecionados and agente.get('system_prompt'):
        contexto_agente += f"DIRETRIZES PRINCIPAIS:\n\n"
    
    if "base_conhecimento" in segmentos_selecionados and agente.get('base_conhecimento'):
        contexto_agente += f"BASE DE CONHECIMENTO:\n\n\n"
    
    if "comments" in segmentos_selecionados and agente.get('comments'):
        contexto_agente += f"COMENTÁRIOS E OBSERVAÇÕES:\n\n\n"
    
    if "planejamento" in segmentos_selecionados and agente.get('planejamento'):
        contexto_agente += f"PLANEJAMENTO E ESTRATÉGIA:\n\n\n"
    
    # Construir instruções baseadas nas configurações
    instrucoes_revisao = ""
    
    if revisao_estilo:
        instrucoes_revisao += """
        - Analise e melhore a clareza, coesão e coerência textual
        - Verifique adequação ao tom da marca
        - Elimine vícios de linguagem e redundâncias
        - Simplifique frases muito longas ou complexas
        """
    
    if manter_estrutura:
        instrucoes_revisao += """
        - Mantenha a estrutura geral do texto original
        - Preserve parágrafos e seções quando possível
        - Conserve o fluxo lógico do conteúdo
        """
    
    if explicar_alteracoes:
        instrucoes_revisao += """
        - Inclua justificativa para as principais alterações
        - Explique correções gramaticais importantes
        - Destaque melhorias de estilo significativas
        """
    
    # Construir o prompt para revisão
    prompt_revisao = f"""
    
    TEXTO PARA REVISÃO:
    {texto}
    
    INSTRUÇÕES PARA REVISÃO:
    
    1. **REVISÃO ORTOGRÁFICA E GRAMATICAL:**
       - Corrija erros de ortografia, acentuação e grafia
       - Verifique concordância nominal e verbal
       - Ajuste pontuação (vírgulas, pontos, travessões)
       - Corrija regência verbal e nominal
       - Ajuste colocação pronominal
    
    2. **REVISÃO DE ESTILO E CLAREZA:**
       {instrucoes_revisao}
    
    FORMATO DA RESPOSTA:
    
    ## 📋 TEXTO REVISADO
    [Aqui vai o texto completo revisado, mantendo a estrutura geral quando possível]
    
    ## 🔍 PRINCIPAIS ALTERAÇÕES REALIZADAS
    [Lista das principais correções realizadas com justificativa]
    
    ## 📊 RESUMO DA REVISÃO
    [Resumo dos problemas encontrados e melhorias aplicadas]
    
    **IMPORTANTE:**
    - Seja preciso nas explicações
    - Mantenha o formato markdown para fácil leitura
    - Foque nas correções ortográficas e gramaticais
    """
    
    try:
        resposta = gerar_resposta_modelo(prompt_revisao, modelo_escolhido)
        return resposta
        
    except Exception as e:
        return f"❌ Erro durante a revisão: {str(e)}"

def revisar_documento_por_slides(doc, agente, segmentos_selecionados, revisao_estilo=True, explicar_alteracoes=True, modelo_escolhido="Gemini"):
    """Revisa documento slide por slide com análise detalhada"""
    
    resultados = []
    
    for i, slide in enumerate(doc['slides']):
        with st.spinner(f"Revisando slide {i+1} de {len(doc['slides'])}..."):
            try:
                # Construir contexto do agente para este slide
                contexto_agente = "CONTEXTO DO AGENTE PARA REVISÃO:\n\n"
                
                if "system_prompt" in segmentos_selecionados and agente.get('system_prompt'):
                    contexto_agente += f"DIRETRIZES PRINCIPAIS:\n{agente['system_prompt']}\n\n"
                
                if "base_conhecimento" in segmentos_selecionados and agente.get('base_conhecimento'):
                    contexto_agente += f"BASE DE CONHECIMENTO:\n{agente['base_conhecimento']}\n\n"
                
                prompt_slide = f"""
{contexto_agente}

## REVISÃO ORTOGRÁFICA - SLIDE {i+1}

**CONTEÚDO DO SLIDE {i+1}:**
{slide['conteudo'][:1500]}

**INSTRUÇÕES:**
- Faça uma revisão ortográfica e gramatical detalhada
- Corrija erros de português, acentuação e pontuação
- Mantenha o conteúdo original - apenas corrija ortograficamente e aponte onde as correções foram feitas
- { "Inclua sugestões de melhoria de estilo" if revisao_estilo else "Foque apenas em correções gramaticais" }
- { "Explique as principais alterações" if explicar_alteracoes else "Apenas apresente o texto corrigido" }

**FORMATO DE RESPOSTA:**

### 📋 SLIDE {i+1} - TEXTO REVISADO
[Texto corrigido do slide]

### 🔍 ALTERAÇÕES REALIZADAS
- [Lista das correções com explicação]

### ✅ STATUS
[✔️ Sem erros / ⚠️ Pequenos ajustes / ❌ Correções necessárias]
"""
                
                resposta = gerar_resposta_modelo(prompt_slide, modelo_escolhido)
                resultados.append({
                    'slide_num': i+1,
                    'analise': resposta,
                    'tem_alteracoes': '❌' in resposta or '⚠️' in resposta or 'Correções' in resposta
                })
                
            except Exception as e:
                resultados.append({
                    'slide_num': i+1,
                    'analise': f"❌ Erro na revisão do slide: {str(e)}",
                    'tem_alteracoes': False
                })
    
    # Construir relatório consolidado
    relatorio = f"# 📊 RELATÓRIO DE REVISÃO ORTOGRÁFICA - {doc['nome']}\n\n"
    relatorio += f"**Total de Slides:** {len(doc['slides'])}\n"
    relatorio += f"**Slides com Correções:** {sum(1 for r in resultados if r['tem_alteracoes'])}\n"
    relatorio += f"**Modelo Utilizado:** {modelo_escolhido}\n\n"
    
    # Slides que precisam de atenção
    slides_com_correcoes = [r for r in resultados if r['tem_alteracoes']]
    if slides_com_correcoes:
        relatorio += "## 🚨 SLIDES COM CORREÇÕES:\n\n"
        for resultado in slides_com_correcoes:
            relatorio += f"### 📋 Slide {resultado['slide_num']}\n"
            relatorio += f"{resultado['analise']}\n\n"
    
    # Resumo executivo
    relatorio += "## 📈 RESUMO EXECUTIVO\n\n"
    if slides_com_correcoes:
        relatorio += f"**⚠️ {len(slides_com_correcoes)} slide(s) necessitam de correções**\n"
        relatorio += f"**✅ {len(doc['slides']) - len(slides_com_correcoes)} slide(s) estão corretos**\n"
        
        # Lista resumida de problemas
        relatorio += "\n**📝 PRINCIPAIS TIPOS DE CORREÇÕES:**\n"
        problemas_comuns = []
        for resultado in slides_com_correcoes:
            if "ortográfico" in resultado['analise'].lower():
                problemas_comuns.append("Erros ortográficos")
            if "pontuação" in resultado['analise'].lower():
                problemas_comuns.append("Problemas de pontuação")
            if "concordância" in resultado['analise'].lower():
                problemas_comuns.append("Erros de concordância")
        
        problemas_unicos = list(set(problemas_comuns))
        for problema in problemas_unicos:
            relatorio += f"- {problema}\n"
    else:
        relatorio += "**🎉 Todos os slides estão ortograficamente corretos!**\n"
    
    return relatorio

# --- ABA: REVISÃO ORTOGRÁFICA ---
with tab_mapping["📝 Revisão Ortográfica"]:
    st.header("📝 Revisão Ortográfica e Gramatical")
    
    # Seletor de modelo para revisão
    st.sidebar.subheader("🤖 Modelo para Revisão")
    modelo_revisao = st.sidebar.selectbox(
        "Escolha o modelo:",
        ["Gemini", "Claude"],
        key="modelo_revisao_selector"
    )
    
    if not st.session_state.agente_selecionado:
        st.info("Selecione um agente primeiro na aba de Chat")
    else:
        agente = st.session_state.agente_selecionado
        st.subheader(f"Revisão com: {agente['nome']}")
        
        # Configurações de segmentos para revisão
        st.sidebar.subheader("🔧 Configurações de Revisão")
        st.sidebar.write("Selecione bases para orientar a revisão:")
        
        segmentos_revisao = st.sidebar.multiselect(
            "Bases para revisão:",
            options=["system_prompt", "base_conhecimento", "comments", "planejamento"],
            default=st.session_state.get('segmentos_selecionados', []),
            key="revisao_segmentos"
        )
        
        # Layout em abas para diferentes métodos de entrada
        tab_texto, tab_arquivo = st.tabs(["📝 Texto Direto", "📎 Upload de Arquivos"])
        
        with tab_texto:
            # Layout em colunas para texto direto
            col_original, col_resultado = st.columns(2)
            
            with col_original:
                st.subheader("📄 Texto Original")
                
                texto_para_revisao = st.text_area(
                    "Cole o texto que deseja revisar:",
                    height=400,
                    placeholder="Cole aqui o texto que precisa de revisão ortográfica e gramatical...",
                    help="O texto será analisado considerando as diretrizes do agente selecionado",
                    key="texto_revisao"
                )
                
                # Estatísticas do texto
                if texto_para_revisao:
                    palavras = len(texto_para_revisao.split())
                    caracteres = len(texto_para_revisao)
                    paragrafos = texto_para_revisao.count('\n\n') + 1
                    
                    col_stats1, col_stats2, col_stats3 = st.columns(3)
                    with col_stats1:
                        st.metric("📊 Palavras", palavras)
                    with col_stats2:
                        st.metric("🔤 Caracteres", caracteres)
                    with col_stats3:
                        st.metric("📄 Parágrafos", paragrafos)
                
                # Configurações de revisão
                with st.expander("⚙️ Configurações da Revisão"):
                    revisao_estilo = st.checkbox(
                        "Incluir revisão de estilo",
                        value=True,
                        help="Analisar clareza, coesão e adequação ao tom da marca",
                        key="revisao_estilo"
                    )
                    
                    manter_estrutura = st.checkbox(
                        "Manter estrutura original",
                        value=True,
                        help="Preservar a estrutura geral do texto quando possível",
                        key="manter_estrutura"
                    )
                    
                    explicar_alteracoes = st.checkbox(
                        "Explicar alterações principais",
                        value=True,
                        help="Incluir justificativa para as mudanças mais importantes",
                        key="explicar_alteracoes"
                    )
            
            with col_resultado:
                st.subheader("📋 Resultado da Revisão")
                
                if st.button("🔍 Realizar Revisão Completa", type="primary", key="revisar_texto"):
                    if not texto_para_revisao.strip():
                        st.warning("⚠️ Por favor, cole o texto que deseja revisar.")
                    else:
                        with st.spinner("🔄 Analisando texto e realizando revisão..."):
                            try:
                                resultado = revisar_texto_ortografia(
                                    texto=texto_para_revisao,
                                    agente=agente,
                                    segmentos_selecionados=segmentos_revisao,
                                    revisao_estilo=revisao_estilo,
                                    manter_estrutura=manter_estrutura,
                                    explicar_alteracoes=explicar_alteracoes,
                                    modelo_escolhido=modelo_revisao
                                )
                                
                                st.markdown(resultado)
                                
                                # Opções de download
                                col_dl1, col_dl2, col_dl3 = st.columns(3)
                                
                                with col_dl1:
                                    st.download_button(
                                        "💾 Baixar Relatório Completo",
                                        data=resultado,
                                        file_name=f"relatorio_revisao_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.txt",
                                        mime="text/plain",
                                        key="download_revisao_completo"
                                    )
                                
                                with col_dl2:
                                    # Extrair apenas o texto revisado se disponível
                                    if "## 📋 TEXTO REVISADO" in resultado:
                                        texto_revisado_start = resultado.find("## 📋 TEXTO REVISADO")
                                        texto_revisado_end = resultado.find("##", texto_revisado_start + 1)
                                        texto_revisado = resultado[texto_revisado_start:texto_revisado_end] if texto_revisado_end != -1 else resultado[texto_revisado_start:]
                                        
                                        st.download_button(
                                            "📄 Baixar Texto Revisado",
                                            data=texto_revisado,
                                            file_name=f"texto_revisado_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.txt",
                                            mime="text/plain",
                                            key="download_texto_revisado"
                                        )
                                
                                with col_dl3:
                                    # Extrair apenas as explicações se disponível
                                    if "## 🔍 PRINCIPAIS ALTERAÇÕES REALIZADAS" in resultado:
                                        explicacoes_start = resultado.find("## 🔍 PRINCIPAIS ALTERAÇÕES REALIZADAS")
                                        explicacoes_end = resultado.find("##", explicacoes_start + 1)
                                        explicacoes = resultado[explicacoes_start:explicacoes_end] if explicacoes_end != -1 else resultado[explicacoes_start:]
                                        
                                        st.download_button(
                                            "📝 Baixar Explicações",
                                            data=explicacoes,
                                            file_name=f"explicacoes_revisao_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.txt",
                                            mime="text/plain",
                                            key="download_explicacoes"
                                        )
                                
                            except Exception as e:
                                st.error(f"❌ Erro ao realizar revisão: {str(e)}")
        
        with tab_arquivo:
            st.subheader("📎 Upload de Arquivos para Revisão")
            
            # Upload de múltiplos arquivos
            arquivos_upload = st.file_uploader(
                "Selecione arquivos para revisão:",
                type=['pdf', 'pptx', 'txt', 'docx'],
                accept_multiple_files=True,
                help="Arquivos serão convertidos para texto e revisados ortograficamente",
                key="arquivos_revisao"
            )
            
            # Configurações para arquivos
            with st.expander("⚙️ Configurações da Revisão para Arquivos"):
                analise_por_slide = st.checkbox(
                    "Análise detalhada por slide/página",
                    value=True,
                    help="Analisar cada slide/página individualmente",
                    key="analise_por_slide"
                )
                
                revisao_estilo_arquivos = st.checkbox(
                    "Incluir revisão de estilo",
                    value=True,
                    help="Analisar clareza, coesão e adequação ao tom da marca",
                    key="revisao_estilo_arquivos"
                )
                
                explicar_alteracoes_arquivos = st.checkbox(
                    "Explicar alterações principais",
                    value=True,
                    help="Incluir justificativa para as mudanças mais importantes",
                    key="explicar_alteracoes_arquivos"
                )
            
            if arquivos_upload:
                st.success(f"✅ {len(arquivos_upload)} arquivo(s) carregado(s)")
                
                # Mostrar preview dos arquivos
                with st.expander("📋 Visualizar Arquivos Carregados", expanded=False):
                    for i, arquivo in enumerate(arquivos_upload):
                        st.write(f"**{arquivo.name}** ({arquivo.size} bytes)")
                
                if st.button("🔍 Revisar Todos os Arquivos", type="primary", key="revisar_arquivos"):
                    resultados_completos = []
                    
                    for arquivo in arquivos_upload:
                        with st.spinner(f"Processando {arquivo.name}..."):
                            try:
                                # Extrair texto do arquivo
                                texto_extraido = ""
                                slides_info = []
                                
                                if arquivo.type == "application/pdf":
                                    texto_extraido, slides_info = extract_text_from_pdf_com_slides(arquivo)
                                elif arquivo.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                                    texto_extraido, slides_info = extract_text_from_pptx_com_slides(arquivo)
                                elif arquivo.type == "text/plain":
                                    texto_extraido = extrair_texto_arquivo(arquivo)
                                elif arquivo.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                                    texto_extraido = extrair_texto_arquivo(arquivo)
                                else:
                                    st.warning(f"Tipo de arquivo não suportado: {arquivo.name}")
                                    continue
                                
                                if texto_extraido and len(texto_extraido.strip()) > 0:
                                    doc_info = {
                                        'nome': arquivo.name,
                                        'conteudo': texto_extraido,
                                        'slides': slides_info,
                                        'tipo': arquivo.type
                                    }
                                    
                                    # Escolher o método de revisão baseado nas configurações
                                    if analise_por_slide and slides_info:
                                        # Revisão detalhada por slide
                                        resultado = revisar_documento_por_slides(
                                            doc_info,
                                            agente,
                                            segmentos_revisao,
                                            revisao_estilo_arquivos,
                                            explicar_alteracoes_arquivos,
                                            modelo_revisao
                                        )
                                    else:
                                        # Revisão geral do documento
                                        resultado = revisar_texto_ortografia(
                                            texto=texto_extraido,
                                            agente=agente,
                                            segmentos_selecionados=segmentos_revisao,
                                            revisao_estilo=revisao_estilo_arquivos,
                                            manter_estrutura=True,
                                            explicar_alteracoes=explicar_alteracoes_arquivos,
                                            modelo_escolhido=modelo_revisao
                                        )
                                    
                                    resultados_completos.append({
                                        'nome': arquivo.name,
                                        'texto_original': texto_extraido,
                                        'resultado': resultado,
                                        'tipo': 'por_slide' if (analise_por_slide and slides_info) else 'geral'
                                    })
                                    
                                    # Exibir resultado individual
                                    with st.expander(f"📄 Resultado - {arquivo.name}", expanded=False):
                                        st.markdown(resultado)
                                        
                                        # Estatísticas do arquivo processado
                                        palavras_orig = len(texto_extraido.split())
                                        st.info(f"📊 Arquivo original: {palavras_orig} palavras")
                                        if slides_info:
                                            st.info(f"📑 {len(slides_info)} slides/páginas processados")
                                        
                                else:
                                    st.warning(f"❌ Não foi possível extrair texto do arquivo: {arquivo.name}")
                                
                            except Exception as e:
                                st.error(f"❌ Erro ao processar {arquivo.name}: {str(e)}")
                    
                    # Botão para download de todos os resultados
                    if resultados_completos:
                        st.markdown("---")
                        st.subheader("📦 Download de Todos os Resultados")
                        
                        # Criar relatório consolidado
                        relatorio_consolidado = f"# RELATÓRIO DE REVISÃO ORTOGRÁFICA\n\n"
                        relatorio_consolidado += f"**Data:** {datetime.datetime.now().strftime('%d/%m/%Y %H:%M')}\n"
                        relatorio_consolidado += f"**Agente:** {agente['nome']}\n"
                        relatorio_consolidado += f"**Modelo Utilizado:** {modelo_revisao}\n"
                        relatorio_consolidado += f"**Total de Arquivos:** {len(resultados_completos)}\n\n"
                        
                        for resultado in resultados_completos:
                            relatorio_consolidado += f"## 📄 {resultado['nome']}\n\n"
                            relatorio_consolidado += f"{resultado['resultado']}\n\n"
                            relatorio_consolidado += "---\n\n"
                        
                        st.download_button(
                            "💾 Baixar Relatório Consolidado",
                            data=relatorio_consolidado,
                            file_name=f"relatorio_revisao_arquivos_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.txt",
                            mime="text/plain",
                            key="download_consolidado"
                        )
            
            else:
                st.info("""
                **📎 Como usar o upload de arquivos:**
                
                1. Selecione um ou mais arquivos (PDF, PPTX, TXT, DOCX)
                2. Configure as opções de revisão
                3. Clique em **"Revisar Todos os Arquivos"**
                
                **📋 Formatos suportados:**
                - PDF (documentos, apresentações) - com análise por página
                - PPTX (apresentações PowerPoint) - com análise por slide
                - TXT (arquivos de texto)
                - DOCX (documentos Word)
                
                **🔍 Análise por Slide/Página:**
                - Identifica slides/páginas específicos com problemas
                - Revisão detalhada de cada seção
                - Facilita a localização e correção de erros
                """)
        
        # Seção informativa
        with st.expander("ℹ️ Sobre a Revisão Ortográfica"):
            st.markdown("""
            ### 🎯 O que é Analisado
            
            **🔤 Ortografia:**
            - Erros de grafia e acentuação
            - Uso correto de maiúsculas e minúsculas
            - Escrita de números e datas
            - Concordância nominal e verbal
            
            **📖 Gramática:**
            - Estrutura sintática das frases
            - Uso adequado de preposições
            - Colocação pronominal
            - Regência verbal e nominal
            
            **🔠 Pontuação:**
            - Uso de vírgulas, pontos, dois-pontos
            - Aplicação de travessões e parênteses
            - Pontuação de citações e diálogos
            
            **📝 Estilo e Clareza:**
            - Coesão e coerência textual
            - Adequação ao tom da marca
            - Clareza na comunicação
            - Eliminação de vícios de linguagem
            
            ### 🤖 Modelos Disponíveis
            
            **Gemini:**
            - Análise rápida e eficiente
            - Boa compreensão de contexto
            - Ideal para textos técnicos
            
            **Claude:**
            - Análise mais detalhada e contextual
            - Melhor compreensão de nuances
            - Excelente para textos criativos
            
            ### 💡 Dicas para Melhor Revisão
            
            1. **Texto Completo**: Cole o texto integral para análise detalhada
            2. **Segmentos Relevantes**: Selecione as bases de conhecimento apropriadas
            3. **Contexto Específico**: Use agentes especializados para cada tipo de conteúdo
            4. **Implementação**: Aplique as sugestões sistematicamente
            
            ### 🎨 Benefícios da Revisão Contextual
            
            - **Consistência da Marca**: Mantém o tom e estilo adequados
            - **Qualidade Profissional**: Elimina erros que prejudicam a credibilidade
            - **Otimização de Conteúdo**: Melhora a clareza e impacto da comunicação
            - **Eficiência**: Reduz tempo de revisão manual
            """)
with tab_mapping["Monitoramento de Redes"]:
    st.header("🤖 Agente de Monitoramento")
    st.markdown("**Especialista que fala como gente**")

    def gerar_resposta_agente(pergunta_usuario: str, historico: List[Dict] = None, agente_monitoramento=None, modelo_escolhido="Gemini") -> str:
        """Gera resposta do agente usando RAG e base do agente de monitoramento"""
        
        # Configuração do agente - usa base do agente selecionado ou padrão
        if agente_monitoramento and agente_monitoramento.get('base_conhecimento'):
            system_prompt = agente_monitoramento['base_conhecimento']
        else:
            # Fallback para prompt padrão se não houver agente selecionado
            system_prompt = """
            PERSONALIDADE: Especialista técnico do agronegócio com habilidade social - "Especialista que fala como gente"

            TOM DE VOZ:
            - Tom que encontra um equilíbrio entre institucional e casual, afinal, as respostas estão sendo geradas no ambiente de rede social por parte de um perfil de empresa
            - Evita exageros e promessas vazias
            - Sempre embasado em fatos e ciência
            - Frases curtas e diretas, mais simpáticas

            DIRETRIZES:
            - NÃO inventar informações técnicas
            - Sempre basear respostas em fatos
            - Manter tom profissional mas acessível
            - Adaptar resposta ao tipo de pergunta
            """
        
        # Constrói o prompt final
        prompt_final = f"""
        {system_prompt}
        
        
        PERGUNTA DO USUÁRIO:
        {pergunta_usuario}
        
        HISTÓRICO DA CONVERSA (se aplicável):
        {historico if historico else "Nenhum histórico anterior"}
        
        INSTRUÇÕES FINAIS:
        Adapte seu tom ao tipo de pergunta:
        - Tom que encontra um equilíbrio entre institucional e casual, afinal, as respostas estão sendo geradas no ambiente de rede social por parte de um perfil de empresa
        - Perguntas técnicas: seja preciso e didático
        - Perguntas sociais: seja leve e engajador  
        - Críticas ou problemas: seja construtivo e proativo
        - Forneça respostas breves - 1 a 3 frases
        
        Sua resposta deve refletir a personalidade do "especialista que fala como gente".
        """
        
        try:
            resposta = gerar_resposta_modelo(prompt_final, modelo_escolhido)
            return resposta
        except Exception as e:
            return f"Erro ao gerar resposta: {str(e)}"

    # SELEÇÃO DE AGENTE DE MONITORAMENTO
    st.header("🔧 Configuração do Agente de Monitoramento")
    
    # Seletor de modelo para monitoramento
    st.sidebar.subheader("🤖 Modelo para Monitoramento")
    modelo_monitoramento = st.sidebar.selectbox(
        "Escolha o modelo:",
        ["Gemini", "Claude"],
        key="modelo_monitoramento_selector"
    )
    
    # Carregar apenas agentes de monitoramento
    agentes_monitoramento = [agente for agente in listar_agentes() if agente.get('categoria') == 'Monitoramento']
    
    col_sel1, col_sel2 = st.columns([3, 1])
    
    with col_sel1:
        if agentes_monitoramento:
            # Criar opções para selectbox
            opcoes_agentes = {f"{agente['nome']}": agente for agente in agentes_monitoramento}
            
            agente_selecionado_nome = st.selectbox(
                "Selecione o agente de monitoramento:",
                list(opcoes_agentes.keys()),
                key="seletor_monitoramento"
            )
            
            agente_monitoramento = opcoes_agentes[agente_selecionado_nome]
            
            # Mostrar informações do agente selecionado
            with st.expander("📋 Informações do Agente Selecionado", expanded=False):
                if agente_monitoramento.get('base_conhecimento'):
                    st.text_area(
                        "Base de Conhecimento:",
                        value=agente_monitoramento['base_conhecimento'],
                        height=200,
                        disabled=True
                    )
                else:
                    st.warning("⚠️ Este agente não possui base de conhecimento configurada")
                
                st.write(f"**Criado em:** {agente_monitoramento['data_criacao'].strftime('%d/%m/%Y %H:%M')}")
                # Mostrar proprietário se for admin
                if get_current_user() == "admin" and agente_monitoramento.get('criado_por'):
                    st.write(f"**👤 Proprietário:** {agente_monitoramento['criado_por']}")
        
        else:
            st.error("❌ Nenhum agente de monitoramento encontrado.")
            st.info("💡 Crie um agente de monitoramento na aba 'Gerenciar Agentes' primeiro.")
            agente_monitoramento = None
    
    with col_sel2:
        if st.button("🔄 Atualizar Lista", key="atualizar_monitoramento"):
            st.rerun()

    # Sidebar com informações
    with st.sidebar:
        st.header("ℹ️ Sobre o Monitoramento")
        
        if agente_monitoramento:
            st.success(f"**Agente Ativo:** {agente_monitoramento['nome']}")
        else:
            st.warning("⚠️ Nenhum agente selecionado")
        
        st.markdown("""
        **Personalidade:**
        - 🎯 Técnico mas acessível
        - 💬 Direto mas simpático
        - 🌱 Conhece o campo e a internet
        - 🔬 Baseado em ciência
        
        **Capacidades:**
        - Respostas técnicas baseadas em RAG
        - Engajamento em redes sociais
        - Suporte a produtores
        - Esclarecimento de dúvidas
        """)

        
        if st.button("🔄 Reiniciar Conversa", key="reiniciar_monitoramento"):
            if "messages_monitoramento" in st.session_state:
                st.session_state.messages_monitoramento = []
            st.rerun()

        # Status da conexão
        
        if os.getenv('OPENAI_API_KEY'):
            st.success("✅ OpenAI: Configurado")
        else:
            st.warning("⚠️ OpenAI: Não configurado")

    # Inicializar histórico de mensagens específico para monitoramento
    if "messages_monitoramento" not in st.session_state:
        st.session_state.messages_monitoramento = []

    # Área de chat principal
    st.header("💬 Simulador de Respostas do Agente")

    # Exemplos de perguntas rápidas
    st.subheader("🎯 Exemplos para testar:")
    col1, col2, col3 = st.columns(3)

    with col1:
        if st.button("❓ Dúvida Técnica", use_container_width=True, key="exemplo_tecnico"):
            st.session_state.messages_monitoramento.append({"role": "user", "content": "Esse produto serve pra todas as culturas?"})

    with col2:
        if st.button("😊 Comentário Social", use_container_width=True, key="exemplo_social"):
            st.session_state.messages_monitoramento.append({"role": "user", "content": "O campo tá lindo demais!"})

    with col3:
        if st.button("⚠️ Crítica/Problema", use_container_width=True, key="exemplo_critica"):
            st.session_state.messages_monitoramento.append({"role": "user", "content": "Usei e não funcionou."})

    # Exibir histórico de mensagens
    for message in st.session_state.messages_monitoramento:
        with st.chat_message(message["role"]):
            st.markdown(message["content"])

    # Input do usuário
    if prompt := st.chat_input("Digite sua mensagem ou pergunta...", key="chat_monitoramento"):
        # Adicionar mensagem do usuário
        st.session_state.messages_monitoramento.append({"role": "user", "content": prompt})
        with st.chat_message("user"):
            st.markdown(prompt)
        
        # Gerar resposta do agente
        with st.chat_message("assistant"):
            with st.spinner("🌱 Consultando base técnica..."):
                resposta = gerar_resposta_agente(
                    prompt, 
                    st.session_state.messages_monitoramento,
                    agente_monitoramento,
                    modelo_monitoramento
                )
                st.markdown(resposta)
                
                # Adicionar ao histórico
                st.session_state.messages_monitoramento.append({"role": "assistant", "content": resposta})

    # Seção de análise de performance
    st.markdown("---")
    st.header("📊 Análise da Resposta")

    if st.session_state.messages_monitoramento:
        ultima_resposta = st.session_state.messages_monitoramento[-1]["content"] if st.session_state.messages_monitoramento[-1]["role"] == "assistant" else ""
        
        if ultima_resposta:
            col_analise1, col_analise2, col_analise3 = st.columns(3)
            
            with col_analise1:
                # Análise de tom
                if "😊" in ultima_resposta or "😍" in ultima_resposta:
                    st.metric("Tom Identificado", "Social/Engajador", delta="Leve")
                elif "🔬" in ultima_resposta or "📊" in ultima_resposta:
                    st.metric("Tom Identificado", "Técnico", delta="Preciso")
                else:
                    st.metric("Tom Identificado", "Balanceado", delta="Adaptado")
            
            with col_analise2:
                # Comprimento da resposta
                palavras = len(ultima_resposta.split())
                st.metric("Tamanho", f"{palavras} palavras")
            
            with col_analise3:
                # Uso de emojis
                emojis = sum(1 for char in ultima_resposta if char in "😀😃😄😁😆😅😂🤣☺️😊😇🙂🙃😉😌😍🥰😘😗😙😚😋😛😝😜🤪🤨🧐🤓😎🤩🥳😏😒😞😔😟😕🙁☹️😣😖😫😩🥺😢😭😤😠😡🤬🤯😳🥵🥶😱😨😰😥😓🤗🤔🤭🤫🤥😶😐😑😬🙄😯😦😧😮😲🥱😴🤤😪😵🤐🥴🤢🤮🤧😷🤒🤕🤑🤠😈👿👹👺🤡💩👻💀☠️👽👾🤖🎃😺😸😹😻😼😽🙀😿😾")
                st.metric("Emojis", emojis, delta="Moderado" if emojis <= 2 else "Alto")

    # Seção de exemplos de uso
    with st.expander("📋 Exemplos de Respostas do Agente"):
        st.markdown("""
        **🎯 PERGUNTA TÉCNICA:**
        *Usuário:* "Qual a diferença entre os nematoides de galha e de cisto na soja?"
        
        **🤖 AGENTE:** "Boa pergunta! Os nematoides de galha (Meloidogyne) formam aquelas 'inchações' nas raízes, enquanto os de cisto (Heterodera) ficam mais externos. Ambos roubam nutrientes, mas o manejo pode ser diferente. Temos soluções específicas para cada caso! 🌱"
        
        **🎯 COMENTÁRIO SOCIAL:**
        *Usuário:* "Adorei ver as fotos da lavoura no stories!"
        
        **🤖 AGENTE:** "A gente também ama compartilhar esses momentos! Quando a tecnologia encontra o cuidado certo, o campo fica ainda mais bonito 😍 Compartilhe suas fotos também!"
        
        **🎯 CRÍTICA/PROBLEMA:**
        *Usuário:* "A aplicação não deu o resultado esperado"
        
        **🤖 AGENTE:** "Poxa, que pena saber disso! Vamos entender melhor o que aconteceu. Pode me contar sobre as condições de aplicação? Assim conseguimos te orientar melhor da próxima vez. A equipe técnica também está à disposição! 📞"
        """)


# --- Funções auxiliares para busca web ---
def buscar_perplexity(pergunta: str, contexto_agente: str = None) -> str:
    """Realiza busca na web usando API do Perplexity"""
    try:
        headers = {
            "Authorization": f"Bearer {perp_api_key}",
            "Content-Type": "application/json"
        }
        
        # Construir o conteúdo da mensagem
        messages = []
        
        if contexto_agente:
            messages.append({
                "role": "system",
                "content": f"Contexto do agente: {contexto_agente}"
            })
        
        messages.append({
            "role": "user",
            "content": pergunta
        })
        
        data = {
            "model": "sonar-medium-online",
            "messages": messages,
            "max_tokens": 2000,
            "temperature": 0.1
        }
        
        response = requests.post(
            "https://api.perplexity.ai/chat/completions",
            headers=headers,
            json=data,
            timeout=30
        )
        
        if response.status_code == 200:
            result = response.json()
            return result['choices'][0]['message']['content']
        else:
            return f"❌ Erro na busca: {response.status_code} - {response.text}"
            
    except Exception as e:
        return f"❌ Erro ao conectar com Perplexity: {str(e)}"

def analisar_urls_perplexity(urls: List[str], pergunta: str, contexto_agente: str = None) -> str:
    """Analisa URLs específicas usando Perplexity"""
    try:
        headers = {
            "Authorization": f"Bearer {perp_api_key}",
            "Content-Type": "application/json"
        }
        
        # Construir contexto com URLs
        urls_contexto = "\n".join([f"- {url}" for url in urls])
        
        messages = []
        
        if contexto_agente:
            messages.append({
                "role": "system",
                "content": f"Contexto do agente: {contexto_agente}"
            })
        
        messages.append({
            "role": "user",
            "content": f"""Analise as seguintes URLs e responda à pergunta:

URLs para análise:
{urls_contexto}

Pergunta: {pergunta}

Forneça uma análise detalhada baseada no conteúdo dessas URLs."""
        })
        
        data = {
            "model": "sonar-medium-online",
            "messages": messages,
            "max_tokens": 3000,
            "temperature": 0.1
        }
        
        response = requests.post(
            "https://api.perplexity.ai/chat/completions",
            headers=headers,
            json=data,
            timeout=45
        )
        
        if response.status_code == 200:
            result = response.json()
            return result['choices'][0]['message']['content']
        else:
            return f"❌ Erro na análise: {response.status_code} - {response.text}"
            
    except Exception as e:
        return f"❌ Erro ao analisar URLs: {str(e)}"

def transcrever_audio_video(arquivo, tipo):
    """Função placeholder para transcrição de áudio/vídeo"""
    return f"Transcrição do {tipo} {arquivo.name} - Esta funcionalidade requer configuração adicional de APIs de transcrição."

# --- Estilização ---
st.markdown("""
<style>
    .stChatMessage {
        padding: 1rem;
        border-radius: 0.5rem;
        margin-bottom: 1rem;
    }
    [data-testid="stChatMessageContent"] {
        font-size: 1rem;
    }
    .stChatInput {
        bottom: 20px;
        position: fixed;
        width: calc(100% - 5rem);
    }
    div[data-testid="stTabs"] {
        margin-top: -30px;
    }
    div[data-testid="stVerticalBlock"] > div:has(>.stTextArea) {
        border-left: 3px solid #4CAF50;
        padding-left: 1rem;
    }
    .segment-indicator {
        background-color: #f0f2f6;
        padding: 0.5rem;
        border-radius: 0.5rem;
        margin: 0.5rem 0;
        border-left: 4px solid #4CAF50;
    }
    .video-analysis-card {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        color: white;
        padding: 1.5rem;
        border-radius: 10px;
        margin: 1rem 0;
    }
    .inheritance-badge {
        background-color: #e3f2fd;
        color: #1976d2;
        padding: 0.2rem 0.5rem;
        border-radius: 12px;
        font-size: 0.8rem;
        margin-left: 0.5rem;
    }
    .web-search-card {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        color: white;
        padding: 1.5rem;
        border-radius: 10px;
        margin: 1rem 0;
    }
    .seo-analysis-card {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        color: white;
        padding: 1.5rem;
        border-radius: 10px;
        margin: 1rem 0;
    }
    .spelling-review-card {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        color: white;
        padding: 1.5rem;
        border-radius: 10px;
        margin: 1rem 0;
    }
    .validation-unified-card {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        color: white;
        padding: 1.5rem;
        border-radius: 10px;
        margin: 1rem 0;
    }
    .user-indicator {
        background-color: #e8f5e8;
        padding: 0.3rem 0.8rem;
        border-radius: 15px;
        font-size: 0.8rem;
        color: #2e7d32;
        border: 1px solid #c8e6c9;
        margin-left: 0.5rem;
    }
</style>
""", unsafe_allow_html=True)

# --- Informações do sistema na sidebar ---
with st.sidebar:
    st.markdown("---")
    st.subheader("🔐 Sistema de Isolamento")
    
    current_user = get_current_user()
    if current_user == "admin":
        st.success("👑 **Modo Administrador**")
        st.info("Visualizando e gerenciando TODOS os agentes do sistema")
    else:
        st.success(f"👤 **Usuário: {current_user}**")
        st.info("Visualizando e gerenciando apenas SEUS agentes")
    
    # Estatísticas rápidas
    agentes_usuario = listar_agentes()
    if agentes_usuario:
        categorias_count = {}
        for agente in agentes_usuario:
            cat = agente.get('categoria', 'Social')
            categorias_count[cat] = categorias_count.get(cat, 0) + 1
        
        st.markdown("### 📊 Seus Agentes")
        for categoria, count in categorias_count.items():
            st.write(f"- **{categoria}:** {count} agente(s)")
        
        st.write(f"**Total:** {len(agentes_usuario)} agente(s)")

# --- Rodapé ---
st.markdown("---")
st.caption(f"🤖 Agente Social v2.0 | Usuário: {get_current_user()} | {datetime.datetime.now().strftime('%d/%m/%Y %H:%M')}")
